# Copyright Thales 2026
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.

"""Offline tests for the shared PPTX→PDF conversion helper.

These do not invoke the real ``soffice`` binary: they stub ``shutil.which`` and
``subprocess.run`` so the timeout / missing-binary / success behaviours are asserted
deterministically without any external dependency. A separate ``integration``-marked
test exercises the real LibreOffice path when the binary is present.
"""

from __future__ import annotations

import shutil
import subprocess  # nosec B404: only used to construct result/exception stubs in tests
from pathlib import Path

import pytest

from fred_core.conversion import convert_pptx_bytes_to_pdf, convert_pptx_file_to_pdf

_MODULE = "fred_core.conversion.pptx_pdf"


@pytest.mark.asyncio
async def test_convert_bytes_returns_pdf_on_success(monkeypatch) -> None:
    monkeypatch.setattr(f"{_MODULE}.shutil.which", lambda _: "/usr/bin/soffice")

    # Emulate soffice: write a PDF next to the input in the requested --outdir.
    def fake_run(cmd, **kwargs):
        outdir = Path(cmd[cmd.index("--outdir") + 1])
        src = Path(cmd[-1])
        (outdir / src.with_suffix(".pdf").name).write_bytes(b"%PDF-1.5 fake")
        return subprocess.CompletedProcess(cmd, 0, b"", b"")

    monkeypatch.setattr(f"{_MODULE}.subprocess.run", fake_run)

    pdf = await convert_pptx_bytes_to_pdf(b"fake pptx bytes")

    assert pdf is not None
    assert pdf.startswith(b"%PDF")


@pytest.mark.asyncio
async def test_convert_bytes_returns_none_on_timeout(monkeypatch) -> None:
    monkeypatch.setattr(f"{_MODULE}.shutil.which", lambda _: "/usr/bin/soffice")

    def hung_run(cmd, **kwargs):
        # Mirror a soffice call that exceeds the deadline: to_thread surfaces the
        # TimeoutExpired the helper is expected to swallow into a None result.
        raise subprocess.TimeoutExpired(cmd, kwargs.get("timeout", 0))

    monkeypatch.setattr(f"{_MODULE}.subprocess.run", hung_run)

    pdf = await convert_pptx_bytes_to_pdf(b"fake pptx bytes", timeout_seconds=0.01)

    assert pdf is None


@pytest.mark.asyncio
async def test_convert_bytes_returns_none_when_soffice_missing(monkeypatch) -> None:
    monkeypatch.setattr(f"{_MODULE}.shutil.which", lambda _: None)

    def unexpected_run(cmd, **kwargs):  # pragma: no cover - must not be reached
        raise AssertionError("subprocess.run must not run when soffice is absent")

    monkeypatch.setattr(f"{_MODULE}.subprocess.run", unexpected_run)

    pdf = await convert_pptx_bytes_to_pdf(b"fake pptx bytes")

    assert pdf is None


@pytest.mark.asyncio
async def test_convert_bytes_returns_none_when_no_pdf_produced(monkeypatch) -> None:
    monkeypatch.setattr(f"{_MODULE}.shutil.which", lambda _: "/usr/bin/soffice")

    # soffice exits 0 but writes nothing — a real failure mode worth guarding.
    monkeypatch.setattr(
        f"{_MODULE}.subprocess.run",
        lambda cmd, **kwargs: subprocess.CompletedProcess(cmd, 0, b"", b""),
    )

    pdf = await convert_pptx_bytes_to_pdf(b"fake pptx bytes")

    assert pdf is None


def test_convert_file_returns_none_on_called_process_error(
    monkeypatch, tmp_path
) -> None:
    monkeypatch.setattr(f"{_MODULE}.shutil.which", lambda _: "/usr/bin/soffice")

    def failing_run(cmd, **kwargs):
        raise subprocess.CalledProcessError(1, cmd, stderr=b"boom")

    monkeypatch.setattr(f"{_MODULE}.subprocess.run", failing_run)

    src = tmp_path / "deck.pptx"
    src.write_bytes(b"fake")

    assert convert_pptx_file_to_pdf(src) is None


@pytest.mark.integration
@pytest.mark.asyncio
async def test_convert_bytes_with_real_libreoffice() -> None:
    """End-to-end conversion via the real ``soffice`` binary (skipped when absent)."""
    if shutil.which("soffice") is None:
        pytest.skip("soffice not installed")
    pytest.importorskip("pptx")

    import io

    from pptx import Presentation  # pyright: ignore[reportMissingImports]

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(0, 0, prs.slide_width, prs.slide_height)
    box.text_frame.text = "Hello preview"
    buf = io.BytesIO()
    prs.save(buf)

    pdf = await convert_pptx_bytes_to_pdf(buf.getvalue())

    assert pdf is not None
    assert pdf.startswith(b"%PDF")
