# Copyright Thales 2026
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.

"""Resolution / geometry-dependent validation layer for the PPT Filler toolkit.

Story 01's :func:`fred_capability_ppt_filler.parser.parse` stays pure and
offline: it never touches Knowledge Flow and has no slide-geometry opinion. This module
adds the two validations that *do* need an outside lookup or the slide shapes, behind
injected seams so callers (the analyze endpoint and the save processor in Story 05) can
wire in a real, scoped resolver while tests pass a fake one.

Two image-only checks are layered on top of a :class:`ParseResult`:

- ``folder_not_found`` — a non-empty author ``folder`` that does not resolve to a folder
  (DOCUMENT tag) in the current space. Resolving also writes the resolved
  ``folder_tag_id`` back into the image :class:`KeyField` so downstream save/fill can use
  it.
- ``image_key_invalid_location`` — an image key whose ``{{key}}`` sits in a shape that
  cannot hold a picture (a table cell), detected via
  :func:`list_image_anchors_on_slide`.

Both reuse the same ``{slide, key, code, message}`` :class:`TemplateError` contract and
the same ``{schema, errors}`` :class:`ParseResult`, so a caller can hand the augmented
result straight back (or raise on it) exactly like the bare parse result.

Empty / missing folders are NOT handled here: those are already ``image_without_folder``
from Story 01 and must not be double-reported. Text keys are skipped entirely (no folder,
no geometry constraint).
"""

from __future__ import annotations

import io
from pathlib import Path
from typing import Dict, List, Optional, Protocol, Set, Union

from pptx import Presentation

from fred_capability_ppt_filler.parser import (
    KeyField,
    ParseResult,
    SlideSchema,
    TemplateError,
)
from fred_capability_ppt_filler.traversal import (
    list_image_anchors_on_slide,
)

# Resolution / geometry error codes (RFC "Image support"). The metadata codes live in
# ``parser.py``; these two need an outside lookup or the slide shapes, so they belong to
# this validation layer.
CODE_FOLDER_NOT_FOUND = "folder_not_found"
CODE_IMAGE_KEY_INVALID_LOCATION = "image_key_invalid_location"


class FolderResolver(Protocol):
    """Seam for resolving an author folder string to a DOCUMENT tag id.

    The real implementation (Story 04/05) wraps :class:`KfTagClient` with the request's
    scope (personal vs a specific team) already bound, so :meth:`resolve` takes only the
    folder string. Tests pass an in-memory fake. Resolution is async because the real
    lookup hits Knowledge Flow.
    """

    async def resolve(self, folder: str) -> Optional[str]:
        """Return the DOCUMENT tag id for ``folder`` (e.g. ``"images/flags"``) in the
        current space, or ``None`` if it does not exist."""
        ...


def _folder_not_found_message(key: str, slide_number: int, folder: str) -> str:
    """English fallback message for ``folder_not_found`` (frontend i18n's by code)."""
    return (
        f"{{{{{key}}}}} on slide {slide_number} points to the folder "
        f"'{folder}', which does not exist in this space."
    )


def _image_key_invalid_location_message(slide_number: int, keys: List[str]) -> str:
    """English fallback message for ``image_key_invalid_location`` (grouped per slide).

    Keys are grouped by slide to match how the UI groups errors; the heading tells the
    author how to fix it (move the key into a picture-capable shape).
    """
    placeholders = ", ".join(f"{{{{{key}}}}}" for key in keys)
    return (
        f"On slide {slide_number}, {placeholders} cannot be turned into a picture "
        "because the key sits in a table cell. Move it into a text box or a rectangle "
        "so a picture can replace it."
    )


async def resolve_and_validate_images(
    pptx_source: Union[bytes, str, Path],
    parse_result: ParseResult,
    resolver: Optional[FolderResolver],
) -> ParseResult:
    """Augment ``parse_result`` with folder-resolution and image-location validation.

    ``resolver=None`` skips the folder→tag lookup entirely (no ``folder_not_found``
    errors, ``folder_tag_id`` stays ``None``) while STILL running the pure geometry
    check. The stateless analyze route uses this: it has no platform access, so it
    reports every offline error immediately and leaves folder existence to the save
    round-trip, which always resolves with a real resolver (#1903).

    For every IMAGE key in ``parse_result.slides``:

    - If its ``folder`` is non-empty, resolve it via ``resolver`` (each distinct folder
      string is resolved at most once per call — see the ``resolved`` cache). On success
      the resolved ``folder_tag_id`` is written back into the schema's :class:`KeyField`;
      on ``None`` a ``folder_not_found`` ``(slide, key)`` error is appended and
      ``folder_tag_id`` stays ``None``. Empty / missing folders are skipped (they are
      already ``image_without_folder`` from Story 01 and must not be double-reported).
    - If ANY anchor for that key on that slide reports ``invalid_location`` (i.e. the key
      sits in a table cell), an ``image_key_invalid_location`` ``(slide, key)`` error is
      appended.

    Returns a NEW :class:`ParseResult` with the same ``{schema, errors}`` contract: the
    original base errors first, then the appended resolution/location errors, and a schema
    whose image keys carry their resolved ``folder_tag_id``. The input ``parse_result`` is
    not mutated.
    """
    if isinstance(pptx_source, bytes):
        presentation = Presentation(io.BytesIO(pptx_source))
    else:
        presentation = Presentation(str(pptx_source))

    # Per-slide set of keys whose anchors (any occurrence) sit in an invalid location.
    # Geometry is read once per slide from the shared image-anchor traversal.
    invalid_location_keys_by_slide: Dict[int, Set[str]] = {}
    for index, slide in enumerate(presentation.slides):
        slide_number = index + 1
        invalid_keys: Set[str] = set()
        for anchor in list_image_anchors_on_slide(slide):
            if anchor.invalid_location:
                invalid_keys.add(anchor.key)
        if invalid_keys:
            invalid_location_keys_by_slide[slide_number] = invalid_keys

    # Resolve each DISTINCT non-empty folder once (several keys may share a folder).
    resolved: Dict[str, Optional[str]] = {}

    # Rebuild the schema so the resolved folder_tag_id is reflected in the returned
    # result. KeyField is a frozen-by-convention pydantic model; we mutate via
    # ``model_copy(update=...)`` to produce well-formed copies rather than poking fields.
    new_slides: List[SlideSchema] = []
    new_errors: List[TemplateError] = list(parse_result.errors)

    # Track which (slide, key) image keys we have already errored on for location, so a
    # multi-occurrence key is reported once; group keys per slide for the message order.
    location_error_keys_by_slide: Dict[int, List[str]] = {}

    for slide_schema in parse_result.slides:
        slide_number = slide_schema.slide
        invalid_keys = invalid_location_keys_by_slide.get(slide_number, set())
        new_keys: List[KeyField] = []

        for key_field in slide_schema.keys:
            if key_field.type != "image":
                new_keys.append(key_field)
                continue

            updated = key_field

            # --- folder resolution (non-empty folders only) ---
            folder = key_field.folder
            if folder and resolver is not None:
                if folder not in resolved:
                    resolved[folder] = await resolver.resolve(folder)
                tag_id = resolved[folder]
                if tag_id is None:
                    new_errors.append(
                        TemplateError(
                            slide=slide_number,
                            key=key_field.key,
                            code=CODE_FOLDER_NOT_FOUND,
                            message=_folder_not_found_message(
                                key_field.key, slide_number, folder
                            ),
                        )
                    )
                else:
                    updated = key_field.model_copy(update={"folder_tag_id": tag_id})

            new_keys.append(updated)

            # --- invalid location (table cell) ---
            if key_field.key in invalid_keys:
                location_error_keys_by_slide.setdefault(slide_number, [])
                if key_field.key not in location_error_keys_by_slide[slide_number]:
                    location_error_keys_by_slide[slide_number].append(key_field.key)

        new_slides.append(SlideSchema(slide=slide_number, keys=new_keys))

    # Emit one image_key_invalid_location error per (slide, key), in slide/key order,
    # carrying a per-slide grouped message (matching how the UI groups).
    for slide_number, keys in location_error_keys_by_slide.items():
        message = _image_key_invalid_location_message(slide_number, keys)
        for key in keys:
            new_errors.append(
                TemplateError(
                    slide=slide_number,
                    key=key,
                    code=CODE_IMAGE_KEY_INVALID_LOCATION,
                    message=message,
                )
            )

    return ParseResult(schema=new_slides, errors=new_errors)
