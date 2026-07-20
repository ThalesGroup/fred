# Copyright Thales 2026
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.

"""Chat-time fill tools of the ``ppt_filler`` capability (#1903).

Port of Kea's ``integrations/ppt_filler/toolkit.py`` re-seated on the Swift
capability seams: the agent context becomes a typed
:class:`~fred_sdk.contracts.capability.CapabilityContext`, the KF clients
become typed `RuntimeServices` ports, and the output upload goes through the
workspace filesystem port.

Pipeline (all derived from the instance's persisted stored config, never from
hardcoded slide indices):

1. **Dynamic per-slide ``args_schema``** — built at middleware-construction
   time from the persisted ``schema_slides``. Nested by slide: one
   ``slide_<n>`` object per schema slide, each leaf carrying its note
   description so the model gets inline guidance.
2. **Fill** — fetch the template through ``services.agent_assets`` (the
   per-instance config-asset store), open with python-pptx, and fill each
   provided ``slide_<n>`` group via the SHARED traversal
   (:func:`~fred_capability_ppt_filler.traversal.replace_keys_on_slide`).
   Image keys are placed as pictures via the image-anchor traversal, fetching
   bytes through ``services.document_content``.
3. **Deliver** — write the filled deck to the agent workspace
   (``services.workspace_fs``), best-effort convert to PDF
   (:func:`fred_core.convert_pptx_bytes_to_pdf`) and return either a
   :class:`PptPreviewPart` (preview pane + download) or a plain download
   :class:`LinkPart` fallback.

Shape coverage is defined by the shared traversal, not re-implemented here:
text boxes, table cells, and shapes nested inside groups are all filled.
SmartArt and chart text remain out of scope (no clean ``text_frame`` API in
python-pptx).
"""

from __future__ import annotations

import hashlib
import io
import json
import logging
import time
from typing import TYPE_CHECKING, Any, Dict, List, Literal, Optional, Tuple, cast

from fred_core import convert_pptx_bytes_to_pdf
from fred_sdk.contracts.capability import CapabilityContext, EmptyModel
from fred_sdk.contracts.context import (
    LinkKind,
    LinkPart,
    PublishedArtifact,
    ToolInvocationResult,
    UiPart,
)
from fred_sdk.contracts.runtime import RuntimeServices
from langchain_core.tools import BaseTool, StructuredTool
from PIL import Image, UnidentifiedImageError
from pptx import Presentation
from pydantic import BaseModel, Field, create_model

from fred_capability_ppt_filler.parser import apply_kept_notes_to_slide
from fred_capability_ppt_filler.traversal import (
    ImageAnchor,
    list_image_anchors_on_slide,
    replace_keys_on_slide,
)

if TYPE_CHECKING:  # pragma: no cover - typing only
    from fred_capability_ppt_filler.capability import PptFillerConfig
    from fred_capability_ppt_filler.parser import KeyField, SlideSchema

logger = logging.getLogger(__name__)

_FILL_TOOL_NAME = "fill_ppt_template"
_LIST_TOOL_NAME = "list_images_in_folder"
_TOOL_REF = "ppt_filler_fill"
_PPTX_CONTENT_TYPE = (
    "application/vnd.openxmlformats-officedocument.presentationml.presentation"
)
_PDF_CONTENT_TYPE = "application/pdf"
_OUTPUT_FILE_NAME = "filled_presentation.pptx"
_PPTX_SUFFIX = ".pptx"
_PDF_SUFFIX = ".pdf"

# Pillow ``format`` values python-pptx's ``add_picture`` can actually embed.
# Anything KF ingests outside this set (WEBP, ICO, ...) decodes in Pillow but is
# rejected by ``add_picture`` with "unsupported image format"; such images are
# transcoded to PNG in-memory before embedding.
_PPTX_EMBEDDABLE_FORMATS = frozenset({"BMP", "GIF", "JPEG", "PNG", "TIFF", "WMF"})

# Top-level (non-slide) arg the model fills with a human-friendly name for the
# output deck. Optional: when missing/blank we fall back to _OUTPUT_FILE_NAME.
_OUTPUT_NAME_FIELD = "output_file_name"

# The top-level args field for each slide is ``slide_<n>`` (1-based
# ``SlideSchema.slide``). Building and parsing the prefix go through the same
# helpers so they can never drift.
_SLIDE_FIELD_PREFIX = "slide_"

# Appended to every TEXT key's leaf description so the agent knows it may
# emphasize part of a value with inline Markdown (see inline_markdown.py).
_TEXT_FORMATTING_HINT = (
    "You may use inline Markdown to emphasize part of this value: **bold** and "
    "*italic* (or _italic_). Markup is optional — omit it for plain text."
)


class PptPreviewPart(BaseModel):
    """
    Contributed ``ppt_preview`` chat part (#1903, preview-pane extension).

    Why this exists:
    - after a fill, the deck is shown as a read-only PDF preview in a resizable
      side pane, so the user can glance at the result and ask for a correction
      without downloading anything. The card carries both open-preview and
      .pptx download, so a separate download chip is redundant.

    Durability:
    - both URLs are the DURABLE, bearer-protected KF ``/fs/download/{path}``
      hrefs (never a signed short-TTL link) — this part is persisted in the
      conversation, so a frozen signed URL would 403 when the chat is reopened
      later. The frontend fetches them with the session bearer at open time.

    Freshness (the edit→preview loop):
    - ``version`` is stamped per fill (PDF content hash) and used as the
      react-pdf remount key, so a re-fill under the same storage key still
      yields a fresh fetch and an open pane updates live.
    """

    type: Literal["ppt_preview"] = "ppt_preview"
    preview_id: str  # stable id for this deck within the session (storage key)
    title: str
    pdf_download_url: str  # durable KF /fs/download href for the preview PDF
    version: str  # per-fill token; drives cache-busting + react-pdf remount
    pptx_download_url: Optional[str] = None  # durable href for the source .pptx
    file_name: Optional[str] = None  # .pptx download file name


def _slide_field_name(slide_number: int) -> str:
    return f"{_SLIDE_FIELD_PREFIX}{slide_number}"


def _slide_number_from_field(field_name: str) -> Optional[int]:
    if not field_name.startswith(_SLIDE_FIELD_PREFIX):
        return None
    suffix = field_name[len(_SLIDE_FIELD_PREFIX) :]
    try:
        return int(suffix)
    except ValueError:
        return None


def _sanitize_output_file_name(raw: object) -> str:
    """Turn a model-provided output name into a safe ``*.pptx`` file name.

    The name is used both as a storage key segment and as the download file
    name, so it must not contain path separators or be empty. We keep only the
    final path component, strip control/quote characters, ensure a single
    ``.pptx`` suffix, and fall back to :data:`_OUTPUT_FILE_NAME` when nothing
    usable remains.
    """
    if not isinstance(raw, str):
        return _OUTPUT_FILE_NAME
    name = raw.replace("\\", "/").split("/")[-1].strip().strip('"').strip("'").strip()
    name = "".join(c for c in name if c.isprintable() and c not in '<>:"\\|?*')
    if name.lower().endswith(_PPTX_SUFFIX):
        name = name[: -len(_PPTX_SUFFIX)].rstrip()
    if not name:
        return _OUTPUT_FILE_NAME
    return f"{name}{_PPTX_SUFFIX}"


def _pdf_key_for(pptx_key: str) -> str:
    """Storage key for the preview PDF: the ``.pptx`` key with a ``.pdf`` suffix."""
    base = pptx_key
    if base.lower().endswith(_PPTX_SUFFIX):
        base = base[: -len(_PPTX_SUFFIX)]
    return f"{base}{_PDF_SUFFIX}"


def _preview_version(pdf_bytes: bytes) -> str:
    """Per-fill freshness token derived from the PDF content.

    A re-fill that changes the deck yields different bytes → a different token
    → the frontend refetches. Content-hashing is deterministic, which keeps the
    "version changes on re-fill" behaviour testable offline.
    """
    return hashlib.sha256(pdf_bytes).hexdigest()[:16]


def _image_leaf_description(key_field: "KeyField") -> str:
    """Build the leaf description for an IMAGE key.

    The agent does not pass image bytes: it passes the *document id* of a
    document it picked by listing the key's folder with the
    ``list_images_in_folder`` tool. The per-field description stays minimal —
    the generic "how to work with an image field" procedure lives once in the
    fill tool's main description, not repeated on every key.
    """
    note = (key_field.description or "").strip()
    folder = key_field.folder

    parts: list[str] = []
    if note:
        parts.append(note)
    if folder:
        parts.append(f"Image field; its images come from the folder '{folder}'.")
    else:
        parts.append(
            "Image field, but no source folder is configured — do your best "
            "(e.g. leave it unset) and explain to the user that an owner/editor "
            "must fix the template's image folder."
        )
    return " ".join(parts)


def _build_args_schema(schema_slides: "list[SlideSchema]") -> type[BaseModel]:
    """Build the dynamic, nested-by-slide ``args_schema`` from the persisted schema.

    For each slide we synthesize a per-slide model whose fields are that
    slide's keys, each carrying the key's note description as its schema-level
    ``description``. The top-level model then has one ``slide_<n>`` field per
    slide. Everything is derived from ``schema_slides`` — no hardcoded indices.

    Every leaf is OPTIONAL (default ``None``): an omitted text key fills as
    empty string (unchanged), an omitted image key removes its empty
    placeholder slot.
    """
    top_fields: Dict[str, Any] = {}
    for slide_schema in schema_slides:
        leaf_fields: Dict[str, Any] = {}
        for key_field in slide_schema.keys:
            if key_field.type == "image":
                leaf_description = _image_leaf_description(key_field)
            else:
                note = (key_field.description or "").strip()
                leaf_description = (
                    f"{note} {_TEXT_FORMATTING_HINT}" if note else _TEXT_FORMATTING_HINT
                )
            leaf_fields[key_field.key] = (
                Optional[str],
                Field(default=None, description=leaf_description),
            )
        slide_model = create_model(
            f"PptFillSlide{slide_schema.slide}",
            __base__=BaseModel,
            **leaf_fields,
        )
        top_fields[_slide_field_name(slide_schema.slide)] = (
            slide_model,
            Field(
                ...,
                description=(
                    f"Values for the {{key}} fields on slide {slide_schema.slide}."
                ),
            ),
        )
    top_fields[_OUTPUT_NAME_FIELD] = (
        Optional[str],
        Field(
            default=None,
            description=(
                "A short, relevant file name for the generated PowerPoint, based "
                "on its content (e.g. 'Proposition_ACME_2026'). The '.pptx' "
                "extension is added automatically. Optional — omit to use a "
                "generic default."
            ),
        ),
    )
    return create_model(
        "FillPptTemplateArgs",
        __base__=BaseModel,
        **top_fields,
    )


def _raw_values_for_slide(slide_args: object) -> Dict[str, object]:
    """Coerce one ``slide_<n>`` group (a pydantic model) into a raw ``{key: value}`` map."""
    if isinstance(slide_args, BaseModel):
        raw = slide_args.model_dump()
    elif isinstance(slide_args, dict):
        raw = slide_args
    else:
        raw = dict(getattr(slide_args, "__dict__", {}) or {})
    return dict(raw)


def _text_values(raw: Dict[str, object], image_keys: set[str]) -> Dict[str, str]:
    """Project the raw slide values onto the TEXT keys only, as fill strings."""
    return {
        key: "" if value is None else str(value)
        for key, value in raw.items()
        if key not in image_keys
    }


def _fit_inside_box(
    img_w: int, img_h: int, box_left: int, box_top: int, box_w: int, box_h: int
) -> Tuple[int, int, int, int]:
    """Scale ``(img_w, img_h)`` to fit inside ``(box_w, box_h)`` preserving aspect
    ratio, centered within the box. Inputs/outputs are EMU; the image pixel dims
    only set the aspect ratio. Returns ``(left, top, width, height)`` in EMU.
    """
    if img_w <= 0 or img_h <= 0 or box_w <= 0 or box_h <= 0:
        return (box_left, box_top, box_w, box_h)

    ar = img_w / img_h
    if box_w / box_h > ar:
        height = box_h
        width = box_h * ar
        top = box_top
        left = box_left + (box_w - width) / 2
    else:
        width = box_w
        height = box_w / ar
        left = box_left
        top = box_top + (box_h - height) / 2
    return (round(left), round(top), round(width), round(height))


def _remove_anchor_shape(anchor) -> None:
    """Remove a placeholder shape from its slide (drops the empty ``{{key}}`` box)."""
    element = anchor.shape._element
    element.getparent().remove(element)


def _looks_like_path_not_document_id(value: str) -> bool:
    """True when ``value`` is clearly a file PATH rather than a document id.

    An image key expects the opaque document id returned by
    ``list_images_in_folder``. The agent sometimes mistakenly passes the file's
    PATH instead, which then fails the document fetch with an unclear error; we
    detect the mistake up front and hard-fail with an actionable message.
    """
    return "/" in value or "\\" in value


async def _place_images_on_slide(
    *,
    slide,
    image_keys: set[str],
    raw_values: Dict[str, object],
    services: RuntimeServices,
) -> Optional[str]:
    """Place chosen images (or remove empty image slots) on one slide.

    Returns ``None`` on success or an error message string to HARD-FAIL the
    whole fill (the agent's correctable mistake → re-pick), mirroring Kea.
    """
    anchors = list_image_anchors_on_slide(slide)
    anchors_by_key: Dict[str, List[ImageAnchor]] = {}
    for anchor in anchors:
        if anchor.key in image_keys:
            anchors_by_key.setdefault(anchor.key, []).append(anchor)

    for key in image_keys:
        key_anchors = anchors_by_key.get(key, [])
        if not key_anchors:
            continue  # key not present on this slide as an image anchor

        # Guard: an image key in a table cell can never hold a picture. Save
        # rejects this, so a saved agent won't hit it, but refuse loudly.
        if any(anchor.invalid_location for anchor in key_anchors):
            placeholder = f"{{{{{key}}}}}"
            return (
                f"The image field {placeholder} sits in a table cell, which "
                "cannot hold a picture. The template must be fixed by an "
                "owner/editor."
            )

        raw_value = raw_values.get(key)
        document_uid = (
            raw_value.strip()
            if isinstance(raw_value, str) and raw_value.strip()
            else None
        )

        if document_uid is None:
            # Omitted image key: remove every empty placeholder slot.
            for anchor in key_anchors:
                _remove_anchor_shape(anchor)
            continue

        if _looks_like_path_not_document_id(document_uid):
            placeholder = f"{{{{{key}}}}}"
            return (
                f"The value '{document_uid}' you passed for the image field "
                f"{placeholder} looks like a file path, not a document id. "
                "Image fields need the document id returned by the "
                "list_images_in_folder tool (an opaque id with no '/' or '\\'), "
                "not the file's path or name. List the field's folder again "
                "with list_images_in_folder and pass the chosen file's "
                "document id."
            )

        port = services.document_content
        if port is None:
            # No platform port injected (e.g. a bare test harness). Fail LOUD
            # rather than silently skipping the image (RFC §3.9).
            raise RuntimeError(
                "ppt_filler: RuntimeServices.document_content is not available "
                "on this execution path."
            )
        try:
            raw_blob = await port.fetch_raw(document_uid)
            image_bytes = raw_blob.content
        except Exception as exc:
            placeholder = f"{{{{{key}}}}}"
            return (
                f"Could not fetch the document '{document_uid}' you chose for "
                f"the image field {placeholder} [{type(exc).__name__}: {exc}]. "
                "List the field's folder again with list_images_in_folder and "
                "pick a valid image document."
            )

        # Validate the bytes ARE a usable image and read pixel dimensions for
        # the aspect ratio. Formats python-pptx cannot embed are transcoded to
        # PNG in-memory rather than rejected later by ``add_picture``.
        try:
            with Image.open(io.BytesIO(image_bytes)) as image:
                img_w, img_h = image.size
                if image.format not in _PPTX_EMBEDDABLE_FORMATS:
                    buffer = io.BytesIO()
                    image.convert("RGBA").save(buffer, format="PNG")
                    image_bytes = buffer.getvalue()
        except (UnidentifiedImageError, OSError, ValueError) as exc:
            placeholder = f"{{{{{key}}}}}"
            return (
                f"The document '{document_uid}' you chose for the image field "
                f"{placeholder} is not a usable image [{type(exc).__name__}: "
                f"{exc}]. List the field's folder again with "
                "list_images_in_folder and pick a valid image document."
            )

        for anchor in key_anchors:
            left, top, width, height = _fit_inside_box(
                img_w,
                img_h,
                anchor.left,
                anchor.top,
                anchor.width,
                anchor.height,
            )
            try:
                slide.shapes.add_picture(
                    io.BytesIO(image_bytes), left, top, width, height
                )
            except Exception as exc:
                placeholder = f"{{{{{key}}}}}"
                return (
                    f"The document '{document_uid}' you chose for the image "
                    f"field {placeholder} could not be inserted as a picture "
                    f"[{type(exc).__name__}: {exc}]. Pick a valid image document."
                )
            _remove_anchor_shape(anchor)

    return None


def _build_list_images_tool(
    ctx: "CapabilityContext[PptFillerConfig, EmptyModel]",
) -> BaseTool:
    """The folder-listing companion tool for image fields (#1903).

    Kea leaned on the document-access ``list_document_tree`` tool; that KF
    surface does not exist on Swift, so the capability ships its own listing
    backed by the ``document_folders`` port and the folder tag ids resolved at
    save time — the agent never sees a tag id, only the author's folder path.
    """
    config = ctx.config
    services = ctx.services

    # folder path -> resolved tag id, from the persisted schema (save-time).
    tag_by_folder: Dict[str, str] = {}
    for slide_schema in config.schema_slides:
        for key_field in slide_schema.keys:
            if (
                key_field.type == "image"
                and key_field.folder
                and key_field.folder_tag_id
            ):
                tag_by_folder[key_field.folder] = key_field.folder_tag_id

    async def _list_images(folder: str) -> str:
        normalized = folder.strip()
        tag_id = tag_by_folder.get(normalized)
        if tag_id is None:
            known = ", ".join(sorted(tag_by_folder)) or "none"
            return (
                f"Unknown image folder '{normalized}'. This template's image "
                f"folders are: {known}. Pass the folder path exactly as written "
                "in the image field's description."
            )
        port = services.document_folders
        if port is None:
            raise RuntimeError(
                "ppt_filler: RuntimeServices.document_folders is not available "
                "on this execution path."
            )
        entries = await port.list_folder_documents(tag_id)
        if not entries:
            return (
                f"The folder '{normalized}' is empty. Leave the image field "
                "unset and tell the user an owner/editor must add images to it."
            )
        return json.dumps(
            [
                {"document_id": entry.document_uid, "name": entry.document_name}
                for entry in entries
            ]
        )

    return StructuredTool.from_function(
        coroutine=_list_images,
        name=_LIST_TOOL_NAME,
        description=(
            "List the image documents available in one of the PowerPoint "
            "template's image folders. Pass the folder path exactly as written "
            "in an image field's description. Returns a JSON list of "
            "{document_id, name}; pass the chosen file's document_id as the "
            "image field's value when calling fill_ppt_template."
        ),
    )


def build_fill_tools(
    ctx: "CapabilityContext[PptFillerConfig, EmptyModel]",
) -> list[BaseTool]:
    """Build the capability's chat-time tools from one turn's typed context.

    With no persisted schema (template never configured) there is nothing to
    fill, so no tools are exposed. The listing companion is added only when the
    schema actually carries image fields with resolved folders.
    """
    config = ctx.config
    schema_slides = config.schema_slides
    if not schema_slides:
        logger.debug(
            "ppt_filler tools requested but the instance has no persisted "
            "schema; returning no tools (template not configured)."
        )
        return []

    services = ctx.services
    identity = ctx.identity
    args_schema = _build_args_schema(schema_slides)
    template_key = config.template_key

    # Per-slide set of IMAGE keys, derived once from the persisted schema.
    image_keys_by_slide: Dict[int, set[str]] = {
        slide_schema.slide: {
            key_field.key
            for key_field in slide_schema.keys
            if key_field.type == "image"
        }
        for slide_schema in schema_slides
    }
    has_image_fields = any(image_keys_by_slide.values())

    async def _fill(**kwargs: object) -> tuple[str, ToolInvocationResult]:
        # Validate/normalize the model-provided args against the dynamic schema
        # so we get typed per-slide groups regardless of how LangChain passed
        # them.
        validated = args_schema(**kwargs)
        provided = validated.model_dump()

        session_id = identity.session_id
        started = time.monotonic()

        assets = services.agent_assets
        if assets is None:
            raise RuntimeError(
                "ppt_filler: RuntimeServices.agent_assets is not available on "
                "this execution path."
            )
        workspace = services.workspace_fs
        if workspace is None:
            raise RuntimeError(
                "ppt_filler: RuntimeServices.workspace_fs is not available on "
                "this execution path."
            )
        # Checked HERE (not inside the fill loop) so a missing platform port
        # fails LOUD like the two guards above — inside the loop it would be
        # swallowed by the broad fill except and degrade to a soft tool error.
        if has_image_fields and services.document_content is None:
            raise RuntimeError(
                "ppt_filler: RuntimeServices.document_content is not available "
                "on this execution path."
            )

        # 1. Fetch the template from the per-instance config-asset store.
        try:
            template_bytes = await assets.fetch(template_key)
        except Exception as exc:
            elapsed = time.monotonic() - started
            message = (
                f"Could not fetch the PPT template '{template_key}' from agent "
                f"config storage after {elapsed:.0f}s [{type(exc).__name__}: {exc}]."
            )
            logger.exception("[PPTFILL][TOOL] template fetch FAILED -> %s", message)
            return message, ToolInvocationResult(tool_ref=_TOOL_REF, is_error=True)

        # 2. Open the template and fill, per-slide. Text keys go through the
        #    SHARED text traversal; image keys are placed as pictures via the
        #    image-anchor traversal (or their empty slots removed).
        try:
            presentation = Presentation(io.BytesIO(template_bytes))
            slides = list(presentation.slides)
            for field_name, slide_args in provided.items():
                slide_number = _slide_number_from_field(field_name)
                if slide_number is None:
                    continue
                index = slide_number - 1
                if index < 0 or index >= len(slides):
                    # The persisted schema references a slide the template no
                    # longer has; skip rather than crash.
                    logger.warning(
                        "[PPTFILL][TOOL] schema references slide %d but "
                        "template has only %d slides; skipping.",
                        slide_number,
                        len(slides),
                    )
                    continue
                slide = slides[index]
                raw_values = _raw_values_for_slide(slide_args)
                image_keys = image_keys_by_slide.get(slide_number, set())

                # Text keys: omitted -> empty string. IMAGE keys are preserved
                # verbatim ({{key}}) so the image pass below can still locate
                # each occurrence.
                text_values = _text_values(raw_values, image_keys)

                def value_for(
                    key: str,
                    _values: Dict[str, str] = text_values,
                    _image_keys: set[str] = image_keys,
                ) -> str:
                    if key in _image_keys:
                        return f"{{{{{key}}}}}"
                    return _values.get(key, "")

                replace_keys_on_slide(slide, value_for)

                if image_keys:
                    error = await _place_images_on_slide(
                        slide=slide,
                        image_keys=image_keys,
                        raw_values=raw_values,
                        services=services,
                    )
                    if error is not None:
                        logger.warning(
                            "[PPTFILL][TOOL] image fill rejected -> %s", error
                        )
                        return error, ToolInvocationResult(
                            tool_ref=_TOOL_REF, is_error=True
                        )

            # Strip template-authoring notes from EVERY slide so the
            # ``{{key}}:`` descriptions never leak into the deliverable.
            for slide in slides:
                apply_kept_notes_to_slide(slide)

            buffer = io.BytesIO()
            presentation.save(buffer)
            filled_bytes = buffer.getvalue()
        except Exception as exc:
            elapsed = time.monotonic() - started
            message = (
                f"Could not fill the PPT template after {elapsed:.0f}s "
                f"[{type(exc).__name__}: {exc}]."
            )
            logger.exception("[PPTFILL][TOOL] fill FAILED -> %s", message)
            return message, ToolInvocationResult(tool_ref=_TOOL_REF, is_error=True)

        # 3. Write to the agent workspace (session-scoped path so outputs of
        #    different conversations never clobber each other) and build the
        #    durable download href.
        output_file_name = _sanitize_output_file_name(provided.get(_OUTPUT_NAME_FIELD))
        output_path = (
            f"outputs/{session_id}/{output_file_name}"
            if session_id
            else f"outputs/{output_file_name}"
        )
        try:
            pptx_artifact: PublishedArtifact = await workspace.write(
                output_path,
                filled_bytes,
                content_type=_PPTX_CONTENT_TYPE,
                title=output_file_name,
            )
        except Exception as exc:
            elapsed = time.monotonic() - started
            message = (
                f"Filled the template but could not upload the result after "
                f"{elapsed:.0f}s [{type(exc).__name__}: {exc}]."
            )
            logger.exception("[PPTFILL][TOOL] upload FAILED -> %s", message)
            return message, ToolInvocationResult(tool_ref=_TOOL_REF, is_error=True)

        logger.info(
            "[PPTFILL][TOOL] filled deck uploaded session=%s key=%s",
            session_id,
            pptx_artifact.key,
        )

        # 4. Best-effort PDF preview: convert the filled deck and write the PDF
        #    beside the .pptx. Any failure (conversion unavailable/timeout,
        #    upload error) degrades to the plain download chip so a preview
        #    problem never costs the user the deck.
        pdf_download_url: Optional[str] = None
        pdf_bytes = await convert_pptx_bytes_to_pdf(filled_bytes)
        if pdf_bytes is not None:
            try:
                pdf_artifact = await workspace.write(
                    _pdf_key_for(output_path),
                    pdf_bytes,
                    content_type=_PDF_CONTENT_TYPE,
                    title=_pdf_key_for(output_file_name),
                )
                pdf_download_url = pdf_artifact.href
            except Exception as exc:
                logger.warning(
                    "[PPTFILL][TOOL] preview upload failed (deck still returned): %s",
                    exc,
                )

        if pdf_download_url and pdf_bytes is not None:
            preview = PptPreviewPart(
                preview_id=pptx_artifact.key,
                title=output_file_name,
                pdf_download_url=pdf_download_url,
                version=_preview_version(pdf_bytes),
                pptx_download_url=pptx_artifact.href,
                file_name=pptx_artifact.file_name,
            )
            artifact = ToolInvocationResult(
                tool_ref=_TOOL_REF,
                ui_parts=(cast_ui_part(preview),),
            )
            return (
                f"The presentation '{output_file_name}' has been filled. A "
                "preview opens beside the chat and a download button is shown "
                "automatically; do not write a download link yourself.",
                artifact,
            )

        # Preview unavailable → fall back to the download chip and say so.
        link = LinkPart(
            href=pptx_artifact.href,
            title=f"Download {pptx_artifact.file_name}",
            kind=LinkKind.download,
            mime=_PPTX_CONTENT_TYPE,
            document_uid=pptx_artifact.document_uid,
            file_name=pptx_artifact.file_name,
        )
        artifact = ToolInvocationResult(
            tool_ref=_TOOL_REF,
            ui_parts=(cast_ui_part(link),),
        )
        return (
            f"The presentation '{output_file_name}' has been filled (the "
            "preview could not be generated, so only the download is "
            "available). A download button is shown to the user automatically; "
            "do not write a download link yourself.",
            artifact,
        )

    fill_tool = StructuredTool.from_function(
        coroutine=_fill,
        name=_FILL_TOOL_NAME,
        description=(
            "Fill the agent's configured PowerPoint template with the provided "
            "values. Input is one 'slide_<n>' object per slide, each holding "
            "that slide's {{key}} fields; each field's own description says "
            "what to put there. Optionally set 'output_file_name' to a short, "
            "relevant name for the deck (the '.pptx' extension is added "
            "automatically).\n"
            "IMAGE FIELDS: some fields are images (their description says 'its "
            "images come from the folder <path>'). Before calling THIS tool you "
            "MUST, for EACH such folder, call the 'list_images_in_folder' tool "
            "with that exact folder path to see the files it contains. Never "
            "invent or guess image values without having listed the folder "
            "first. Then, for each image field, pick the file whose name best "
            "fits and pass ONLY that file's document_id as the field's value "
            "(not its name, not text). If a folder is missing or empty, leave "
            "that field unset and tell the user an owner/editor must fix the "
            "template's image folder.\n"
            "After the tool runs, the filled PowerPoint is surfaced to the user "
            "automatically (a preview beside the chat when available, otherwise "
            "a download button) — do NOT write, invent, or repeat any download "
            "link or URL in your reply, and do not tell the user to click a "
            "link you wrote. Just briefly confirm the deck is ready."
        ),
        args_schema=args_schema,
        response_format="content_and_artifact",
    )
    tools: List[BaseTool] = [fill_tool]
    if has_image_fields:
        tools.append(_build_list_images_tool(ctx))
    return tools


def cast_ui_part(part: BaseModel) -> UiPart:
    """Widen a concrete part to the `UiPart` union for `ToolInvocationResult`.

    The static `UiPart` alias is the frozen base union; the registry extends
    the RUNTIME union with capability parts at boot (#1977), so this cast is
    the reference pattern for capability-contributed parts (see `demo.py`).
    """
    return cast(UiPart, part)
