"""Offline, fixture-driven tests for the image-anchor shape-walking traversal.

These exercise :func:`list_image_anchors_on_slide` — the geometry seam that gives each
``{{key}}`` its absolute placement box on a slide. Decks are built in-test with
python-pptx (no checked-in binaries), so every fixture is inspectable in the test that
uses it. The build helpers mirror the ``_build_table_deck`` / ``_build_group_deck``
pattern in ``test_parser.py`` but author the group's ``<a:xfrm>`` explicitly so
the expected absolute geometry is deterministic and the group transform is non-trivial.
"""

import io
from typing import List, Optional, Tuple

from fred_capability_ppt_filler.traversal import (
    ImageAnchor,
    list_image_anchors_on_slide,
)
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

# An EMU box: (left, top, width, height). Used both to place fixture shapes and to assert
# the absolute geometry returned by the traversal.
Box = Tuple[int, int, int, int]


def _slide_of(deck: bytes):
    return Presentation(io.BytesIO(deck)).slides[0]


def _build_textbox_deck(body: str, box: Box) -> bytes:
    """One-slide deck with a single top-level text box at ``box`` containing ``body``."""
    left, top, width, height = box
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(left, top, width, height)
    textbox.text_frame.text = body
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _build_two_textbox_deck(body: str, box_a: Box, box_b: Box) -> bytes:
    """One-slide deck with two top-level text boxes that share the same ``body``."""
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    for left, top, width, height in (box_a, box_b):
        textbox = slide.shapes.add_textbox(left, top, width, height)
        textbox.text_frame.text = body
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _build_split_run_textbox_deck(runs: List[str], box: Box) -> bytes:
    """One-slide deck whose text box paragraph is split across several runs.

    ``runs`` are the run texts of the paragraph (e.g. ``["{{lo", "go}}"]`` so a key
    straddles a run boundary), exercising the run-merging idiom in the geometry seam.
    """
    left, top, width, height = box
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(left, top, width, height)
    paragraph = textbox.text_frame.paragraphs[0]
    for run_text in runs:
        run = paragraph.add_run()
        run.text = run_text
        run.font.size = Pt(18)
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _build_table_deck(cell_texts: List[str], box: Box) -> bytes:
    """One-slide deck whose only key-bearing shape is a 1xN TABLE at ``box``."""
    left, top, width, height = box
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    table = slide.shapes.add_table(1, len(cell_texts), left, top, width, height).table
    for col, text in enumerate(cell_texts):
        table.cell(0, col).text = text
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _author_group_xfrm(grp_sp_pr, *, off: Box, ext: Box, ch_off: Box, ch_ext: Box):
    """Author ``<a:xfrm>`` with explicit ``off``/``ext``/``chOff``/``chExt`` on a group.

    ``_build_group_deck`` in the parser tests wraps shapes in a ``<p:grpSp>`` but leaves
    the group transform unset (so ``.left/.width`` are ``None``). The geometry tests need
    a non-trivial, known transform, so this sets all four child-coordinate elements. ``off``
    and ``ch_off`` use only the first two tuple entries (x, y); ``ext`` and ``ch_ext`` only
    the first two (cx, cy).
    """
    xfrm = grp_sp_pr.makeelement(qn("a:xfrm"), {})
    off_el = xfrm.makeelement(qn("a:off"), {"x": str(off[0]), "y": str(off[1])})
    ext_el = xfrm.makeelement(qn("a:ext"), {"cx": str(ext[0]), "cy": str(ext[1])})
    ch_off_el = xfrm.makeelement(
        qn("a:chOff"), {"x": str(ch_off[0]), "y": str(ch_off[1])}
    )
    ch_ext_el = xfrm.makeelement(
        qn("a:chExt"), {"cx": str(ch_ext[0]), "cy": str(ch_ext[1])}
    )
    for child in (off_el, ext_el, ch_off_el, ch_ext_el):
        xfrm.append(child)
    grp_sp_pr.append(xfrm)


def _make_group(spTree, *, off: Box, ext: Box, ch_off: Box, ch_ext: Box):
    """Build an empty ``<p:grpSp>`` (with authored xfrm) and return it, not yet attached.

    Children are appended by the caller (text boxes or nested groups); the group must then
    be appended to its parent element.
    """
    grpSp = spTree.makeelement(qn("p:grpSp"), {})
    nvGrpSpPr = grpSp.makeelement(qn("p:nvGrpSpPr"), {})
    cNvPr = nvGrpSpPr.makeelement(qn("p:cNvPr"), {"id": "999", "name": "TestGroup"})
    nvGrpSpPr.append(cNvPr)
    nvGrpSpPr.append(nvGrpSpPr.makeelement(qn("p:cNvGrpSpPr"), {}))
    nvGrpSpPr.append(nvGrpSpPr.makeelement(qn("p:nvPr"), {}))
    grpSp.append(nvGrpSpPr)
    grpSpPr = grpSp.makeelement(qn("p:grpSpPr"), {})
    _author_group_xfrm(grpSpPr, off=off, ext=ext, ch_off=ch_off, ch_ext=ch_ext)
    grpSp.append(grpSpPr)
    return grpSp


def _build_group_deck(
    body: str,
    child_box: Box,
    *,
    off: Box,
    ext: Box,
    ch_off: Box,
    ch_ext: Box,
) -> bytes:
    """One-slide deck: a text box at ``child_box`` wrapped in a group with a known xfrm.

    The group's transform maps child (group-local) coordinates to absolute coordinates, so
    the expected absolute box of the text box is fully computable from the supplied EMU
    values.
    """
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    spTree = slide.shapes._spTree

    left, top, width, height = child_box
    box = slide.shapes.add_textbox(left, top, width, height)
    box.text_frame.text = body

    grpSp = _make_group(spTree, off=off, ext=ext, ch_off=ch_off, ch_ext=ch_ext)
    grpSp.append(box._element)  # reparent the sp under the group
    spTree.append(grpSp)

    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _build_nested_group_deck(
    body: str,
    child_box: Box,
    *,
    outer: dict,
    inner: dict,
) -> bytes:
    """One-slide deck: a text box inside an inner group inside an outer group.

    ``outer`` / ``inner`` each carry ``off``/``ext``/``ch_off``/``ch_ext`` keyword dicts.
    The absolute box is the outer transform applied to the inner transform applied to the
    text box, exercising multiplicative composition of nested group transforms.
    """
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    spTree = slide.shapes._spTree

    left, top, width, height = child_box
    box = slide.shapes.add_textbox(left, top, width, height)
    box.text_frame.text = body

    inner_grp = _make_group(spTree, **inner)
    inner_grp.append(box._element)

    outer_grp = _make_group(spTree, **outer)
    outer_grp.append(inner_grp)
    spTree.append(outer_grp)

    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _only(anchors: List[ImageAnchor], key: Optional[str] = None) -> ImageAnchor:
    candidates = [a for a in anchors if key is None or a.key == key]
    assert len(candidates) == 1, f"expected exactly one anchor, got {candidates}"
    return candidates[0]


def _box(anchor: ImageAnchor) -> Box:
    return (anchor.left, anchor.top, anchor.width, anchor.height)


# --- Top-level text box: the shape's own geometry, verbatim -----------------------------


def test_top_level_textbox_reports_its_own_geometry():
    box = (Inches(2), Inches(3), Inches(4), Inches(1))  # left/top/width/height EMU
    deck = _build_textbox_deck("Our logo: {{logo}}", box)

    anchors = list_image_anchors_on_slide(_slide_of(deck))

    anchor = _only(anchors, "logo")
    assert _box(anchor) == box
    assert anchor.invalid_location is False
    assert anchor.shape.has_text_frame


# --- Group child: absolute geometry via the composed transform (HIGHEST-RISK) -----------


def test_group_child_absolute_geometry_with_non_identity_scale():
    """Highest-risk geometry (RFC watch-point #2): a text box inside a group whose extent
    differs from its child-extent (a 2x scale on both axes), with the child *not* sitting
    at ``chOff`` so the offset term is also scaled.

    Group:  off=(2in,3in)  ext=(8in,4in)  chOff=(1in,1in)  chExt=(4in,2in)
            => scale_x = 8/4 = 2.0,  scale_y = 4/2 = 2.0
    Child:  left=2in top=2in width=1in height=0.5in
            abs_left = off_x + (left - chOff_x)*scale_x = 2in + (2in-1in)*2 = 4in
            abs_top  = off_y + (top  - chOff_y)*scale_y = 3in + (2in-1in)*2 = 5in
            abs_w    = width  * scale_x = 1in   * 2 = 2in
            abs_h    = height * scale_y = 0.5in * 2 = 1in
    """
    deck = _build_group_deck(
        "Flag here {{flag}}",
        child_box=(Inches(2), Inches(2), Inches(1), Inches(0.5)),
        off=(Inches(2), Inches(3)),
        ext=(Inches(8), Inches(4)),
        ch_off=(Inches(1), Inches(1)),
        ch_ext=(Inches(4), Inches(2)),
    )

    anchor = _only(list_image_anchors_on_slide(_slide_of(deck)), "flag")

    # Expected absolute EMU (computed above): left=4in, top=5in, width=2in, height=1in.
    assert _box(anchor) == (Inches(4), Inches(5), Inches(2), Inches(1))
    # Sanity: the scale really is non-identity, i.e. width grew vs. the authored 1in.
    assert anchor.width == 2 * Inches(1)
    assert anchor.invalid_location is False


def test_nested_groups_compose_multiplicatively():
    """Two nested groups, each a 2x scale, compose to a 4x scale on the inner child.

    Inner group:  off=(1in,1in)  ext=(4in,4in)  chOff=(0,0)   chExt=(2in,2in)  => 2x
    Outer group:  off=(0,0)      ext=(20in,20in) chOff=(0,0)  chExt=(10in,10in) => 2x
    Child:        left=1in top=1in width=1in height=1in

    Inner maps child to inner-local:
        i_left = 1in + (1in-0)*2 = 3in ; i_w = 1in*2 = 2in   (same for the y axis)
    Outer maps inner-local to absolute:
        abs_left = 0 + (3in-0)*2 = 6in ; abs_w = 2in*2 = 4in (same for the y axis)
    """
    deck = _build_nested_group_deck(
        "Nested {{logo}}",
        child_box=(Inches(1), Inches(1), Inches(1), Inches(1)),
        inner=dict(
            off=(Inches(1), Inches(1)),
            ext=(Inches(4), Inches(4)),
            ch_off=(0, 0),
            ch_ext=(Inches(2), Inches(2)),
        ),
        outer=dict(
            off=(0, 0),
            ext=(Inches(20), Inches(20)),
            ch_off=(0, 0),
            ch_ext=(Inches(10), Inches(10)),
        ),
    )

    anchor = _only(list_image_anchors_on_slide(_slide_of(deck)), "logo")

    assert _box(anchor) == (Inches(6), Inches(6), Inches(4), Inches(4))
    assert anchor.invalid_location is False


def test_zero_child_extent_falls_back_to_unit_scale():
    """A zero child-extent would divide by zero; the scale must fall back to 1.0, leaving
    the group a pure translation by its offset (no NaN / no crash)."""
    deck = _build_group_deck(
        "{{logo}}",
        child_box=(Inches(2), Inches(2), Inches(1), Inches(1)),
        off=(Inches(5), Inches(6)),
        ext=(Inches(8), Inches(4)),
        ch_off=(Inches(1), Inches(1)),
        ch_ext=(0, 0),  # degenerate: forces the div-by-zero guard
    )

    anchor = _only(list_image_anchors_on_slide(_slide_of(deck)), "logo")

    # scale == 1.0 on both axes => abs = off + (child - chOff):
    #   left = 5in + (2in-1in) = 6in ; top = 6in + (2in-1in) = 7in ; w/h unchanged (1in).
    assert _box(anchor) == (Inches(6), Inches(7), Inches(1), Inches(1))


# --- Table cell: invalid location, geometry irrelevant ----------------------------------


def test_table_cell_key_is_invalid_location():
    deck = _build_table_deck(
        ["{{flag}}", "plain text"],
        box=(Inches(1), Inches(1), Inches(8), Inches(1)),
    )

    anchor = _only(list_image_anchors_on_slide(_slide_of(deck)), "flag")

    assert anchor.invalid_location is True
    assert anchor.shape.has_table


# --- Same key in two shapes: one anchor per shape (every shape gets filled) --------------


def test_same_key_in_two_shapes_yields_two_anchors():
    box_a = (Inches(1), Inches(1), Inches(3), Inches(1))
    box_b = (Inches(5), Inches(4), Inches(2), Inches(1))
    deck = _build_two_textbox_deck("{{logo}}", box_a, box_b)

    anchors = list_image_anchors_on_slide(_slide_of(deck))

    assert [a.key for a in anchors] == ["logo", "logo"]
    assert {_box(a) for a in anchors} == {box_a, box_b}
    assert all(a.invalid_location is False for a in anchors)


# --- Run-merging: a key split across runs is still anchored ------------------------------


def test_key_split_across_runs_is_still_anchored():
    box = (Inches(2), Inches(2), Inches(3), Inches(1))
    deck = _build_split_run_textbox_deck(["Our {{lo", "go}} here"], box)

    anchor = _only(list_image_anchors_on_slide(_slide_of(deck)), "logo")

    assert _box(anchor) == box
    assert anchor.invalid_location is False


# --- A picture/empty slide contributes no anchors ---------------------------------------


def test_shape_without_keys_contributes_no_anchors():
    deck = _build_textbox_deck(
        "No placeholders here", (Inches(1), Inches(1), Inches(4), Inches(1))
    )

    assert list_image_anchors_on_slide(_slide_of(deck)) == []
