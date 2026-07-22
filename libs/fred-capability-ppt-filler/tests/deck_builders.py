# Copyright Thales 2026
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.

"""Shared in-code ``.pptx`` builders for the capability/fill/analyze tests.

Decks are built with python-pptx (no checked-in binaries), mirroring the helper
style already used by ``tests/test_parser.py`` and ``tests/test_anchors.py``.
Keeping them here avoids duplicating the same builders across ``test_fill.py``,
``test_validate_config.py`` and ``test_analyze_route.py``.
"""

from __future__ import annotations

import io
from typing import List, Tuple

from fred_capability_ppt_filler.parser import SlideSchema, parse
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

# (text-box body, notes text) per slide.
SlideSpec = Tuple[str, str]

# An image-key notes block: type:image + a folder, plus prose so it is well-formed.
IMAGE_NOTES = "{{logo}}:\n- type: image\n- folder: Brand/Logos\nThe company logo"


def build_deck(slides: List[SlideSpec]) -> bytes:
    """Build a ``.pptx`` with one text box and one notes block per slide spec."""
    presentation = Presentation()
    blank_layout = presentation.slide_layouts[6]  # blank
    for body, notes in slides:
        slide = presentation.slides.add_slide(blank_layout)
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
        textbox.text_frame.text = body
        if notes:
            slide.notes_slide.notes_text_frame.text = notes
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def build_image_deck(slides: List[Tuple[str, List[str]]]) -> bytes:
    """Build a deck whose slides each carry one or more textbox specs.

    ``slides`` is a list of ``(notes, [textbox_text, ...])``. Each textbox is
    added at a distinct, fixed box (4in wide, 2in tall) so image anchors have a
    real geometry to fit a picture inside.
    """
    presentation = Presentation()
    blank_layout = presentation.slide_layouts[6]  # blank
    for notes, textboxes in slides:
        slide = presentation.slides.add_slide(blank_layout)
        for offset, text in enumerate(textboxes):
            box = slide.shapes.add_textbox(
                Inches(1),
                Inches(1 + 2 * offset),
                Inches(4),  # wider than tall -> a square image fits the HEIGHT
                Inches(2),
            )
            box.text_frame.text = text
        if notes:
            slide.notes_slide.notes_text_frame.text = notes
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def build_table_deck(notes: str, cell_texts: List[str]) -> bytes:
    """Build a one-slide deck whose key-bearing shape is a 1xN TABLE, plus notes.

    A ``{{key}}`` sitting in a table cell is an invalid picture location, so this
    exercises the image-location refusal path.
    """
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    table = slide.shapes.add_table(
        1, len(cell_texts), Inches(1), Inches(1), Inches(8), Inches(1)
    ).table
    for col, text in enumerate(cell_texts):
        table.cell(0, col).text = text
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def schema_slides(deck: bytes) -> List[SlideSchema]:
    """Parse a deck and return its persisted-style ``schema_slides`` list."""
    return parse(deck).slides


def image_schema(
    deck: bytes, *, slide: int, key: str, tag_id: str = "tag-xyz"
) -> List[SlideSchema]:
    """Parse a deck and stamp a resolved ``folder_tag_id`` on one image key.

    The save-time processor is what resolves the folder to a tag id; here we set
    it directly on the parsed schema so the fill tool sees the same persisted
    shape it would after a real save.
    """
    slides = parse(deck).slides
    for slide_schema in slides:
        if slide_schema.slide != slide:
            continue
        for key_field in slide_schema.keys:
            if key_field.key == key:
                key_field.folder_tag_id = tag_id
    return slides


def png_bytes(width: int = 10, height: int = 10, color: str = "red") -> bytes:
    """A tiny, valid PNG with the given pixel dimensions (no checked-in binaries)."""
    buffer = io.BytesIO()
    Image.new("RGB", (width, height), color).save(buffer, format="PNG")
    return buffer.getvalue()


def webp_bytes(width: int = 10, height: int = 10, color: str = "red") -> bytes:
    """A tiny, valid WEBP. Pillow decodes it but python-pptx's ``add_picture``
    cannot embed WEBP, so the fill tool must transcode it to PNG first."""
    buffer = io.BytesIO()
    Image.new("RGB", (width, height), color).save(buffer, format="WEBP")
    return buffer.getvalue()


def picture_shapes(slide) -> list:
    """Every PICTURE shape on a slide (python-pptx exposes ``shape_type`` PICTURE)."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    return [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
