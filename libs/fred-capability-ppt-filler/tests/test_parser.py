"""Offline, fixture-driven tests for the PPT Filler parser and shared traversal.

Decks are built in-test with python-pptx (no checked-in binaries), so each fixture is
inspectable in the test that uses it. Style mirrors
``tests/test_kf_vector_search_tools.py`` (same inprocess-toolkit family).
"""

import io
from typing import List, Optional, Tuple

import pytest
from fred_capability_ppt_filler.parser import (
    CODE_DESCRIBED_BUT_NOT_IN_SLIDE,
    CODE_DUPLICATED_METADATA,
    CODE_EMPTY_FOLDER,
    CODE_FOLDER_WITHOUT_IMAGE_TYPE,
    CODE_IMAGE_WITHOUT_FOLDER,
    CODE_KEY_WITHOUT_DESCRIPTION,
    CODE_UNKNOWN_METADATA,
    CODE_UNKNOWN_TYPE,
    parse,
)
from fred_capability_ppt_filler.traversal import (
    list_keys_on_slide,
    replace_keys_on_slide,
)
from pptx import Presentation
from pptx.util import Inches, Pt

# (text-box body, notes text) per slide. A blank text-box keeps a slide that only has
# notes (to exercise described_but_not_in_slide).
SlideSpec = Tuple[str, str]


def _build_deck(slides: List[SlideSpec]) -> bytes:
    """Build a ``.pptx`` with one text box and one notes block per slide spec.

    The text box uses a single run; the run-split case is built separately by
    ``_build_split_run_deck``.
    """
    presentation = Presentation()
    blank_layout = presentation.slide_layouts[6]  # blank
    for body, notes in slides:
        slide = presentation.slides.add_slide(blank_layout)
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
        textbox.text_frame.text = body
        if notes:
            slide.notes_slide.notes_text_frame.text = notes
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _build_split_run_deck(
    runs: List[str], notes: str, *, extra_runs: Optional[List[str]] = None
) -> bytes:
    """Build a one-slide deck whose text box paragraph is split into several runs.

    ``runs`` are the run texts of the first paragraph (e.g. ``["{{na", "me}}"]`` so the
    placeholder straddles a run boundary). ``extra_runs`` optionally adds a second
    paragraph.
    """
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
    text_frame = textbox.text_frame
    paragraph = text_frame.paragraphs[0]
    for run_text in runs:
        run = paragraph.add_run()
        run.text = run_text
        run.font.size = Pt(18)
    if extra_runs:
        second = text_frame.add_paragraph()
        for run_text in extra_runs:
            run = second.add_run()
            run.text = run_text
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _build_table_deck(notes: str, cell_texts: List[str]) -> bytes:
    """Build a one-slide deck whose only key-bearing shape is a TABLE.

    ``cell_texts`` fills the cells of a 1xN table row-major; ``{{key}}`` placeholders in
    those cells must be discovered (and fillable) exactly like text-box keys.
    """
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    rows, cols = 1, len(cell_texts)
    table = slide.shapes.add_table(
        rows, cols, Inches(1), Inches(1), Inches(8), Inches(1)
    ).table
    for col, text in enumerate(cell_texts):
        table.cell(0, col).text = text
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _build_group_deck(notes: str, textbox_bodies: List[str]) -> bytes:
    """Build a one-slide deck whose key-bearing text boxes live inside a GROUP shape.

    python-pptx cannot author groups via the public API, so we build separate text boxes
    and then wrap their XML in a ``<p:grpSp>`` element — mirroring how PowerPoint nests
    shapes inside ``Groupe 1`` in real decks.
    """
    from pptx.oxml.ns import qn

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    spTree = slide.shapes._spTree

    boxes = []
    for i, body in enumerate(textbox_bodies):
        box = slide.shapes.add_textbox(Inches(1), Inches(1 + i), Inches(4), Inches(1))
        box.text_frame.text = body
        boxes.append(box)

    # Wrap the just-added text boxes into a single group shape.
    grpSp = spTree.makeelement(qn("p:grpSp"), {})
    nvGrpSpPr = grpSp.makeelement(qn("p:nvGrpSpPr"), {})
    cNvPr = nvGrpSpPr.makeelement(qn("p:cNvPr"), {"id": "999", "name": "TestGroup"})
    nvGrpSpPr.append(cNvPr)
    nvGrpSpPr.append(nvGrpSpPr.makeelement(qn("p:cNvGrpSpPr"), {}))
    nvGrpSpPr.append(nvGrpSpPr.makeelement(qn("p:nvPr"), {}))
    grpSp.append(nvGrpSpPr)
    grpSp.append(grpSp.makeelement(qn("p:grpSpPr"), {}))
    for box in boxes:
        grpSp.append(box._element)  # reparent the sp under the group
    spTree.append(grpSp)

    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _build_bulleted_notes_deck(
    header: str, bullet_lines: List[str], prose: str, body: str
) -> bytes:
    """Build a one-slide deck whose metadata lines are real PowerPoint BULLETS.

    Reproduces what PowerPoint's notes editor does to a line the author types as
    ``- type: image``: it DROPS the literal ``"- "`` from the run text and records the
    dash as a paragraph bullet (``a:pPr`` → ``a:buChar char="-"``) instead. ``header`` and
    ``prose`` are plain (unbulleted) paragraphs; each entry in ``bullet_lines`` becomes a
    bulleted paragraph whose run text has NO leading dash. The parser must still read these
    as ``- key: value`` metadata. ``body`` is the slide's text-box content (the ``{{key}}``
    placeholder).
    """
    from pptx.oxml.ns import qn

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
    textbox.text_frame.text = body

    notes_frame = slide.notes_slide.notes_text_frame
    notes_frame.text = header  # first paragraph: the {{key}}: header (no bullet)
    for line in bullet_lines:
        paragraph = notes_frame.add_paragraph()
        paragraph.text = line  # run text carries NO leading dash
        # Attach a dash bullet at the paragraph level, exactly as PowerPoint does.
        pPr = paragraph._p.get_or_add_pPr()
        pPr.append(pPr.makeelement(qn("a:buFontTx"), {}))
        pPr.append(pPr.makeelement(qn("a:buChar"), {"char": "-"}))
    if prose:
        notes_frame.add_paragraph().text = prose  # trailing prose (no bullet)

    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _slide(result, slide_number: int):
    return next(s for s in result.slides if s.slide == slide_number)


def _keys(result, slide_number: int) -> List[str]:
    return [field.key for field in _slide(result, slide_number).keys]


def _description(result, slide_number: int, key: str) -> str:
    return next(
        field.description
        for field in _slide(result, slide_number).keys
        if field.key == key
    )


def _field(result, slide_number: int, key: str):
    return next(
        field for field in _slide(result, slide_number).keys if field.key == key
    )


def test_per_slide_grouping_of_keys():
    deck = _build_deck(
        [
            (
                "Hello {{name}} and {{role}}",
                "{{name}}:\nThe person's name\n{{role}}:\nTheir role",
            ),
            ("City: {{city}}", "{{city}}:\nThe city"),
        ]
    )
    result = parse(deck)

    assert [s.slide for s in result.slides] == [1, 2]
    assert _keys(result, 1) == ["name", "role"]
    assert _keys(result, 2) == ["city"]
    assert result.errors == []


def test_same_key_on_two_slides_is_two_independent_fields():
    deck = _build_deck(
        [
            ("{{name}}", "{{name}}:\nName on slide 1"),
            ("{{name}}", "{{name}}:\nName on slide 2"),
        ]
    )
    result = parse(deck)

    assert _keys(result, 1) == ["name"]
    assert _keys(result, 2) == ["name"]
    assert _description(result, 1, "name") == "Name on slide 1"
    assert _description(result, 2, "name") == "Name on slide 2"
    assert result.errors == []


def test_same_key_twice_on_one_slide_is_one_field():
    deck = _build_deck([("Header {{name}} ... footer {{name}}", "{{name}}:\nThe name")])
    result = parse(deck)

    assert _keys(result, 1) == ["name"]  # deduped to a single field
    assert result.errors == []


def test_multi_key_header_shares_description():
    deck = _build_deck(
        [("{{first}} {{last}}", "{{first}}, {{last}}:\nA part of the full name")]
    )
    result = parse(deck)

    assert _keys(result, 1) == ["first", "last"]
    assert _description(result, 1, "first") == "A part of the full name"
    assert _description(result, 1, "last") == "A part of the full name"
    assert result.errors == []


def test_inline_braces_in_description_not_mistaken_for_header():
    # The second notes line mentions {{name}} inline; it must be captured as part of the
    # description, NOT treated as a new header (which would otherwise reset the block).
    notes = "{{name}}:\nWrite the name like {{name}} here, e.g. Jane Doe"
    deck = _build_deck([("{{name}}", notes)])
    result = parse(deck)

    assert _keys(result, 1) == ["name"]
    assert (
        _description(result, 1, "name")
        == "Write the name like {{name}} here, e.g. Jane Doe"
    )
    assert result.errors == []


def test_multiline_description_with_internal_blank_lines_captured():
    notes = "{{bio}}:\nFirst line\n\nThird line after a blank\n"
    deck = _build_deck([("{{bio}}", notes)])
    result = parse(deck)

    assert _description(result, 1, "bio") == "First line\n\nThird line after a blank"
    assert result.errors == []


def test_leading_and_trailing_blank_lines_trimmed():
    notes = "{{bio}}:\n\n\nThe biography\n\n"
    deck = _build_deck([("{{bio}}", notes)])
    result = parse(deck)

    assert _description(result, 1, "bio") == "The biography"


def test_key_without_description_error():
    deck = _build_deck([("Hi {{name}} aged {{age}}", "{{name}}:\nThe name")])
    result = parse(deck)

    matching = [e for e in result.errors if e.code == CODE_KEY_WITHOUT_DESCRIPTION]
    assert len(matching) == 1
    error = matching[0]
    assert error.slide == 1
    assert error.key == "age"
    assert "slide 1" in error.message
    assert "age" in error.message


def test_described_but_not_in_slide_error():
    # Slide 1 text box has {{name}} but notes also describe {{ghost}}.
    deck = _build_deck(
        [("{{name}}", "{{name}}:\nThe name\n{{ghost}}:\nA stale description")]
    )
    result = parse(deck)

    matching = [e for e in result.errors if e.code == CODE_DESCRIBED_BUT_NOT_IN_SLIDE]
    assert len(matching) == 1
    error = matching[0]
    assert error.slide == 1
    assert error.key == "ghost"
    assert "slide 1" in error.message
    assert "ghost" in error.message


def test_error_slide_numbers_are_per_slide():
    deck = _build_deck(
        [
            ("{{ok}}", "{{ok}}:\nFine"),
            ("{{name}}", "{{name}}:\nName\n{{ghost}}:\nGhost"),  # ghost on slide 2
        ]
    )
    result = parse(deck)

    ghost_errors = [
        e for e in result.errors if e.code == CODE_DESCRIBED_BUT_NOT_IN_SLIDE
    ]
    assert len(ghost_errors) == 1
    assert ghost_errors[0].slide == 2


def test_key_split_across_runs_is_reconstructed():
    deck = _build_split_run_deck(["Hello {{na", "me}} world"], "{{name}}:\nThe name")
    result = parse(deck)

    assert _keys(result, 1) == ["name"]
    assert result.errors == []


def test_schema_serializes_with_schema_key():
    deck = _build_deck([("{{name}}", "{{name}}:\nThe name")])
    result = parse(deck)

    dumped = result.model_dump(by_alias=True)
    assert "schema" in dumped
    assert "errors" in dumped
    # Serialization is additive and backward compatible: a plain TEXT key keeps the legacy
    # ``{key, description}`` shape (the default image fields are suppressed on the wire).
    assert dumped["schema"] == [
        {"slide": 1, "keys": [{"key": "name", "description": "The name"}]}
    ]


def test_image_key_serializes_its_metadata_fields():
    """An IMAGE key serializes its ``type``/``folder`` (additive), while ``folder_tag_id``
    stays suppressed because it is still unresolved (``None``) in this story."""
    notes = "{{flag}}:\n- type: image\n- folder: images/flags\nPick it."
    deck = _build_deck([("{{flag}}", notes)])
    result = parse(deck)

    dumped = result.model_dump(by_alias=True)
    assert dumped["schema"] == [
        {
            "slide": 1,
            "keys": [
                {
                    "key": "flag",
                    "description": "Pick it.",
                    "type": "image",
                    "folder": "images/flags",
                }
            ],
        }
    ]


# --- Shared traversal direct tests -----------------------------------------------


def test_list_keys_traversal_handles_split_runs_directly():
    deck = _build_split_run_deck(["{{fir", "st}} and {{las", "t}}"], "")
    presentation = Presentation(io.BytesIO(deck))
    slide = presentation.slides[0]

    assert list_keys_on_slide(slide) == ["first", "last"]


def test_replace_keys_traversal_handles_split_runs():
    deck = _build_split_run_deck(["{{na", "me}}!"], "")
    presentation = Presentation(io.BytesIO(deck))
    slide = presentation.slides[0]

    replace_keys_on_slide(slide, lambda key: {"name": "Jane"}[key])

    assert list_keys_on_slide(slide) == []  # no placeholder remains
    text = "".join(
        run.text
        for shape in slide.shapes
        if shape.has_text_frame
        for paragraph in shape.text_frame.paragraphs
        for run in paragraph.runs
    )
    assert text == "Jane!"


# --- Tables and grouped shapes are walked too -------------------------------------


def test_keys_in_table_cells_are_discovered():
    deck = _build_table_deck(
        "{{header}}, {{body}}:\nA cell value",
        ["{{header}}", "{{body}}"],
    )
    result = parse(deck)

    assert _keys(result, 1) == ["header", "body"]
    assert result.errors == []


def test_keys_in_grouped_shapes_are_discovered():
    deck = _build_group_deck(
        "{{a}}, {{b}}:\nValues in a group",
        ["{{a}}", "{{b}}"],
    )
    result = parse(deck)

    assert _keys(result, 1) == ["a", "b"]
    assert result.errors == []


def test_fill_replaces_keys_in_table_cells():
    deck = _build_table_deck("{{x}}:\nThe x", ["before {{x}} after"])
    presentation = Presentation(io.BytesIO(deck))
    slide = presentation.slides[0]

    assert list_keys_on_slide(slide) == ["x"]
    replace_keys_on_slide(slide, lambda key: {"x": "FILLED"}[key])
    assert list_keys_on_slide(slide) == []  # nothing left in the table cell

    # And the literal replacement landed in the table cell text.
    table = next(s.table for s in slide.shapes if s.has_table)
    assert "FILLED" in table.cell(0, 0).text


def test_fill_replaces_keys_in_grouped_shapes():
    deck = _build_group_deck("{{g}}:\nThe g", ["start {{g}} end"])
    presentation = Presentation(io.BytesIO(deck))
    slide = presentation.slides[0]

    assert list_keys_on_slide(slide) == ["g"]
    replace_keys_on_slide(slide, lambda key: {"g": "DONE"}[key])
    assert list_keys_on_slide(slide) == []


# --- The round-trip regression guard ---------------------------------------------


def test_round_trip_parse_fill_reparse_leaves_no_placeholders():
    """parse(deck) -> fill every field via the shared Replace traversal -> re-parse
    leaves NO remaining {{keys}}. This proves the parser and filler share one traversal
    and cannot diverge."""
    deck = _build_deck(
        [
            ("Hello {{name}}, the {{role}}", "{{name}}, {{role}}:\nName then role"),
            (
                "Repeated {{name}} on slide 2 and {{city}}",
                "{{name}}:\nName again\n{{city}}:\nThe city",
            ),
        ]
    )

    result = parse(deck)
    assert result.errors == []  # well-formed deck

    # Build a value per (slide, key) from the parsed schema and fill every slide.
    presentation = Presentation(io.BytesIO(deck))
    for slide_schema in result.slides:
        slide = presentation.slides[slide_schema.slide - 1]
        values = {
            field.key: f"VAL_{slide_schema.slide}_{field.key}"
            for field in slide_schema.keys
        }
        replace_keys_on_slide(slide, lambda key, _values=values: _values[key])

    buffer = io.BytesIO()
    presentation.save(buffer)
    filled_bytes = buffer.getvalue()

    # Re-parse the filled deck: no field should remain.
    reparsed = parse(filled_bytes)
    assert reparsed.slides == []

    # And no raw {{...}} survives anywhere the shared traversal can reach (text frames,
    # table cells, and grouped shapes).
    refilled = Presentation(io.BytesIO(filled_bytes))
    for slide in refilled.slides:
        assert list_keys_on_slide(slide) == []


def test_parse_accepts_path(tmp_path):
    deck = _build_deck([("{{name}}", "{{name}}:\nThe name")])
    path = tmp_path / "deck.pptx"
    path.write_bytes(deck)

    from_bytes = parse(deck)
    from_path = parse(path)

    assert from_path.model_dump(by_alias=True) == from_bytes.model_dump(by_alias=True)


@pytest.mark.parametrize(
    "line,is_header",
    [
        ("{{name}}:", True),
        ("  {{name}}  :  ", True),
        ("{{a}}, {{b}}:", True),
        ("{{a}},{{b}} , {{c}}:", True),
        ("Write {{name}} here", False),  # inline mention, not a header
        ("{{name}}", False),  # no trailing colon
        ("{{name}}: extra text", False),  # text after the colon
        ("Just prose", False),
    ],
)
def test_header_detection(line, is_header):
    from fred_capability_ppt_filler.parser import _HEADER_PATTERN

    assert bool(_HEADER_PATTERN.match(line)) is is_header


# --- Keep-separator: notes after a "---" line are kept verbatim, never parsed ----------


@pytest.mark.parametrize(
    "notes,authoring,kept",
    [
        # No separator -> everything is authoring, nothing kept.
        ("{{name}}:\nThe name", "{{name}}:\nThe name", ""),
        # Separator -> split; the dash line itself is dropped.
        ("{{name}}:\nThe name\n---\nReal note", "{{name}}:\nThe name", "Real note"),
        # A single leading blank line after the separator is trimmed.
        ("{{name}}:\nDesc\n----\n\nKept", "{{name}}:\nDesc", "Kept"),
        # Fewer than 3 dashes is NOT a separator.
        ("{{name}}:\nDesc\n--\nstill desc", "{{name}}:\nDesc\n--\nstill desc", ""),
        # Multi-line kept content is preserved verbatim (including inner blank lines).
        ("{{a}}:\nx\n---\nl1\n\nl2", "{{a}}:\nx", "l1\n\nl2"),
    ],
)
def test_split_authoring_and_kept_notes(notes, authoring, kept):
    from fred_capability_ppt_filler.parser import (
        split_authoring_and_kept_notes,
    )

    assert split_authoring_and_kept_notes(notes) == (authoring, kept)


def test_kept_notes_are_not_parsed_as_headers():
    """A ``{{key}}:`` line AFTER the keep-separator is kept content, not a description —
    so it must not register a description (and would otherwise cause a false
    described_but_not_in_slide for a key absent from the slide)."""
    deck = _build_deck(
        [
            (
                "Hello {{name}}",
                "{{name}}:\nThe name\n---\nReminder: mention {{ghost}}: the sponsor",
            )
        ]
    )
    result = parse(deck)

    # {{name}} is described (above the separator); the kept text mentioning {{ghost}}: is
    # opaque, so there is NO described_but_not_in_slide error for ghost.
    assert result.errors == []
    assert result.slides[0].keys[0].key == "name"
    assert result.slides[0].keys[0].description == "The name"


# --- Image metadata: type / folder parsing ---------------------------------------------


def test_default_type_is_text_with_no_metadata():
    """Backward compat: a plain text key has type=text and no folder fields."""
    deck = _build_deck([("{{name}}", "{{name}}:\nThe name")])
    result = parse(deck)

    field = _field(result, 1, "name")
    assert field.type == "text"
    assert field.folder is None
    assert field.folder_tag_id is None
    assert result.errors == []


def test_image_key_with_folder_parsed():
    notes = '{{flag}}:\n- type: image\n- folder: "images/flags"\nPick the flag.'
    deck = _build_deck([("{{flag}}", notes)])
    result = parse(deck)

    field = _field(result, 1, "flag")
    assert field.type == "image"
    assert field.folder == "images/flags"
    assert field.folder_tag_id is None
    assert field.description == "Pick the flag."
    assert result.errors == []


@pytest.mark.parametrize(
    "folder_line",
    [
        "- folder: 'images/flags'",  # single quotes
        '- folder: "images/flags"',  # double quotes
        "- folder: images/flags",  # bare, no quotes
    ],
)
def test_folder_quote_stripping(folder_line):
    notes = f"{{{{flag}}}}:\n- type: image\n{folder_line}\nPick it."
    deck = _build_deck([("{{flag}}", notes)])
    result = parse(deck)

    assert _field(result, 1, "flag").folder == "images/flags"
    assert result.errors == []


def test_metadata_keys_and_type_are_case_insensitive():
    notes = "{{flag}}:\n- TYPE: Image\n- Folder: X\nGuidance."
    deck = _build_deck([("{{flag}}", notes)])
    result = parse(deck)

    field = _field(result, 1, "flag")
    assert field.type == "image"  # normalized lowercase
    assert field.folder == "X"
    assert result.errors == []


def test_powerpoint_bullet_metadata_is_parsed_like_dashed_metadata():
    """PowerPoint turns ``- type: image`` typed in the notes into a real bullet: the
    literal ``"- "`` is stripped from the text and stored as a paragraph ``buChar``.
    Such metadata must parse identically to a literal-dash line (regression: it used to
    fall through to a TEXT key, so a chosen image id was rendered as text)."""
    deck = _build_bulleted_notes_deck(
        header="{{flag}}:",
        bullet_lines=["type: image", 'folder: "images/flags"'],
        prose="Pick the flag.",
        body="{{flag}}",
    )
    result = parse(deck)

    field = _field(result, 1, "flag")
    assert field.type == "image"
    assert field.folder == "images/flags"
    assert field.description == "Pick the flag."
    assert result.errors == []


def test_leading_dash_prose_line_ends_metadata_and_is_description():
    """A leading-dash line that is not ``key: value`` ends the metadata block and is
    captured as prose, not metadata, not an error."""
    notes = (
        "{{flag}}:\n- type: image\n- folder: images/flags\n"
        "- choose the most recent\nMore guidance."
    )
    deck = _build_deck([("{{flag}}", notes)])
    result = parse(deck)

    field = _field(result, 1, "flag")
    assert field.type == "image"
    assert field.folder == "images/flags"
    # The non-key:value dash line and the line after it are prose.
    assert field.description == "- choose the most recent\nMore guidance."
    assert result.errors == []


def test_metadata_lines_are_not_in_description():
    notes = "{{flag}}:\n- type: image\n- folder: images/flags\nThe guidance prose."
    deck = _build_deck([("{{flag}}", notes)])
    result = parse(deck)

    assert _field(result, 1, "flag").description == "The guidance prose."


def test_multi_key_header_shares_type_folder_description():
    notes = "{{a}}, {{b}}:\n- type: image\n- folder: shared/dir\nShared guidance."
    deck = _build_deck([("{{a}} {{b}}", notes)])
    result = parse(deck)

    for key in ("a", "b"):
        field = _field(result, 1, key)
        assert field.type == "image"
        assert field.folder == "shared/dir"
        assert field.description == "Shared guidance."
    assert result.errors == []


# --- Image metadata: error codes -------------------------------------------------------


def _codes(result):
    return [e.code for e in result.errors]


def test_unknown_metadata_error():
    notes = "{{flag}}:\n- type: image\n- folder: images/flags\n- bogus: 1\nGuidance."
    deck = _build_deck([("{{flag}}", notes)])
    result = parse(deck)

    matching = [e for e in result.errors if e.code == CODE_UNKNOWN_METADATA]
    assert len(matching) == 1
    assert matching[0].slide == 1
    assert matching[0].key == "flag"


def test_unknown_type_error():
    notes = "{{flag}}:\n- type: video\n- folder: images/flags\nGuidance."
    deck = _build_deck([("{{flag}}", notes)])
    result = parse(deck)

    matching = [e for e in result.errors if e.code == CODE_UNKNOWN_TYPE]
    assert len(matching) == 1
    assert matching[0].slide == 1
    assert matching[0].key == "flag"


def test_duplicated_metadata_error():
    notes = "{{flag}}:\n- type: image\n- type: image\n- folder: x\nGuidance."
    deck = _build_deck([("{{flag}}", notes)])
    result = parse(deck)

    matching = [e for e in result.errors if e.code == CODE_DUPLICATED_METADATA]
    assert len(matching) == 1
    assert matching[0].slide == 1
    assert matching[0].key == "flag"


def test_image_without_folder_when_no_folder_line():
    notes = "{{flag}}:\n- type: image\nGuidance with no folder."
    deck = _build_deck([("{{flag}}", notes)])
    result = parse(deck)

    matching = [e for e in result.errors if e.code == CODE_IMAGE_WITHOUT_FOLDER]
    assert len(matching) == 1
    assert matching[0].slide == 1
    assert matching[0].key == "flag"
    # An empty folder line is NOT what triggered it here, but ensure precedence: no
    # empty_folder is raised for an image.
    assert CODE_EMPTY_FOLDER not in _codes(result)


def test_image_with_blank_folder_maps_to_image_without_folder():
    """type: image + a blank folder line -> image_without_folder (not empty_folder)."""
    notes = "{{flag}}:\n- type: image\n- folder:\nGuidance."
    deck = _build_deck([("{{flag}}", notes)])
    result = parse(deck)

    assert CODE_IMAGE_WITHOUT_FOLDER in _codes(result)
    assert CODE_EMPTY_FOLDER not in _codes(result)
    matching = [e for e in result.errors if e.code == CODE_IMAGE_WITHOUT_FOLDER]
    assert matching[0].slide == 1
    assert matching[0].key == "flag"


def test_empty_folder_on_text_key():
    """A blank folder line on a non-image key -> empty_folder (the value problem is the
    most specific issue; folder_without_image_type is NOT also raised)."""
    notes = "{{name}}:\n- folder:\nThe name."
    deck = _build_deck([("{{name}}", notes)])
    result = parse(deck)

    assert CODE_EMPTY_FOLDER in _codes(result)
    assert CODE_FOLDER_WITHOUT_IMAGE_TYPE not in _codes(result)
    matching = [e for e in result.errors if e.code == CODE_EMPTY_FOLDER]
    assert matching[0].slide == 1
    assert matching[0].key == "name"


def test_folder_without_image_type_on_text_key():
    notes = "{{name}}:\n- folder: images/flags\nThe name."
    deck = _build_deck([("{{name}}", notes)])
    result = parse(deck)

    matching = [e for e in result.errors if e.code == CODE_FOLDER_WITHOUT_IMAGE_TYPE]
    assert len(matching) == 1
    assert matching[0].slide == 1
    assert matching[0].key == "name"
    assert CODE_EMPTY_FOLDER not in _codes(result)


def test_image_key_without_description_still_raises():
    # A key on the slide with NO describing header in the notes still raises
    # key_without_description (image keys are not special-cased out of the check). Here
    # only {{other}} is described, so {{flag}} has no header at all.
    notes = "{{other}}:\n- type: image\n- folder: images/other\nOther image."
    deck = _build_deck([("{{flag}} {{other}}", notes)])
    result = parse(deck)

    matching = [e for e in result.errors if e.code == CODE_KEY_WITHOUT_DESCRIPTION]
    assert len(matching) == 1
    assert matching[0].slide == 1
    assert matching[0].key == "flag"


def test_image_key_described_but_not_in_slide_still_raises():
    # The slide text box has {{here}}; the notes describe an absent image key {{ghost}}.
    notes = (
        "{{here}}:\nA real key\n"
        "{{ghost}}:\n- type: image\n- folder: images/flags\nAbsent image."
    )
    deck = _build_deck([("{{here}}", notes)])
    result = parse(deck)

    matching = [e for e in result.errors if e.code == CODE_DESCRIBED_BUT_NOT_IN_SLIDE]
    assert len(matching) == 1
    assert matching[0].slide == 1
    assert matching[0].key == "ghost"


def test_multi_key_header_attributes_metadata_error_to_each_key():
    notes = "{{a}}, {{b}}:\n- type: video\n- folder: x\nGuidance."
    deck = _build_deck([("{{a}} {{b}}", notes)])
    result = parse(deck)

    unknown_type = [e for e in result.errors if e.code == CODE_UNKNOWN_TYPE]
    assert {e.key for e in unknown_type} == {"a", "b"}
    assert all(e.slide == 1 for e in unknown_type)
