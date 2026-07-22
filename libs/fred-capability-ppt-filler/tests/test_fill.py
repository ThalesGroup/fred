# Copyright Thales 2026
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.

"""Offline tests for the chat-time fill tools of the ppt_filler capability.

Port of Kea's ``tests/test_ppt_filler_toolkit.py`` re-seated on the Swift
capability seams: the tools are built from a typed ``CapabilityContext`` whose
``services`` are in-memory ``RuntimeServices`` port fakes (see
``tests/port_fakes.py``), the template lives in the ``agent_assets`` config-asset
store, and the filled deck is written through the ``workspace_fs`` port.

The fill tool declares ``response_format="content_and_artifact"``, so the
``(content, ToolInvocationResult)`` tuple is observed by calling the underlying
coroutine directly (``tool.coroutine(**args)``); one test also drives the tool
through ``ainvoke`` with a ToolCall dict (the runtime idiom, mirroring
``test_capability_document_access_1906.py``) to prove the artifact survives.
"""

from __future__ import annotations

import io
import json

import fred_capability_ppt_filler.fill as fill_mod
import pytest
from deck_builders import (
    IMAGE_NOTES,
    build_deck,
    build_image_deck,
    build_table_deck,
    image_schema,
    picture_shapes,
    png_bytes,
    schema_slides,
    webp_bytes,
)
from fred_capability_ppt_filler.capability import (
    PPT_FILLER_TEMPLATE_KEY,
    PptFillerConfig,
)
from fred_capability_ppt_filler.fill import PptPreviewPart, build_fill_tools
from fred_capability_ppt_filler.traversal import KEY_PATTERN, list_keys_on_slide
from fred_sdk.contracts.capability import (
    CapabilityContext,
    CapabilityIdentity,
    EmptyModel,
)
from fred_sdk.contracts.context import LinkKind, LinkPart, ToolInvocationResult
from fred_sdk.contracts.runtime import RuntimeServices
from port_fakes import FakeAssets, FakeDocs, FakeFolders, FakeWorkspace
from pptx import Presentation
from pptx.util import Inches

_PPTX_CONTENT_TYPE = (
    "application/vnd.openxmlformats-officedocument.presentationml.presentation"
)


# --- Context / tool assembly helpers ----------------------------------------------


def _ctx(
    slides,
    *,
    session_id="sess-9",
    template_key=PPT_FILLER_TEMPLATE_KEY,
    template_bytes=None,
    assets=None,
    docs=None,
    folders=None,
    workspace=None,
):
    """Build a CapabilityContext with the given persisted schema + service fakes.

    When ``template_bytes`` is given and no explicit ``assets`` fake is passed, a
    FakeAssets seeded with the template under ``template_key`` is created so the
    fill tool can fetch it.
    """
    if assets is None:
        blobs = {template_key: template_bytes} if template_bytes is not None else {}
        assets = FakeAssets(blobs)
    return CapabilityContext(
        identity=CapabilityIdentity(user_id="u-1", session_id=session_id),
        config=PptFillerConfig(schema_slides=list(slides), template_key=template_key),
        turn_options=EmptyModel(),
        services=RuntimeServices(
            agent_assets=assets,
            document_content=docs,
            document_folders=folders,
            workspace_fs=workspace if workspace is not None else FakeWorkspace(),
        ),
    )


def _fill_tool(ctx):
    tools = build_fill_tools(ctx)
    assert tools, "expected at least the fill tool"
    return tools[0]


@pytest.fixture
def no_pdf(monkeypatch):
    """Disable the best-effort PDF preview by default (fill returns the download chip)."""

    async def _none(_content, *args, **kwargs):
        return None

    monkeypatch.setattr(fill_mod, "convert_pptx_bytes_to_pdf", _none)


def _enable_preview(monkeypatch, *, pdf_bytes: bytes = b"%PDF-1.5 fake"):
    async def _pdf(_content, *args, **kwargs):
        return pdf_bytes

    monkeypatch.setattr(fill_mod, "convert_pptx_bytes_to_pdf", _pdf)


@pytest.fixture
def preview_union():
    """Fold ``PptPreviewPart`` into the ``UiPart`` union for the duration of a test.

    In production the capability registry calls ``rebuild_ui_part_union`` at boot
    from ``manifest.chat_parts`` (#1977); a standalone capability-package test does
    not boot the registry, so the runtime union is the frozen ``link | geo`` base
    and a ``PptPreviewPart`` cannot validate into ``ToolInvocationResult.ui_parts``
    until it is registered. This simulates that registration and restores the base
    union on teardown.
    """
    from fred_sdk.contracts.ui_part_union import rebuild_ui_part_union

    rebuild_ui_part_union([PptPreviewPart])
    yield
    rebuild_ui_part_union(())


# --- A. Dynamic per-slide args_schema ---------------------------------------------


def test_args_schema_is_nested_by_slide_with_leaf_descriptions():
    deck = build_deck(
        [
            (
                "Hello {{name}} and {{role}}",
                "{{name}}:\nThe person's name\n{{role}}:\nTheir role",
            ),
            ("City: {{city}}", "{{city}}:\nThe city"),
        ]
    )
    tool = _fill_tool(_ctx(schema_slides(deck)))

    schema = tool.args_schema.model_json_schema()

    # Top-level groups keyed by slide_<n>, plus an optional output_file_name.
    assert set(schema["properties"]) == {"slide_1", "slide_2", "output_file_name"}
    assert set(schema["required"]) == {"slide_1", "slide_2"}
    assert "output_file_name" not in schema.get("required", [])

    defs = schema["$defs"]

    def _leaf(slide_field: str) -> dict:
        ref = schema["properties"][slide_field]["$ref"]
        return defs[ref.split("/")[-1]]

    slide1 = _leaf("slide_1")
    assert set(slide1["properties"]) == {"name", "role"}
    assert (
        slide1["properties"]["name"]["description"]
        == f"The person's name {fill_mod._TEXT_FORMATTING_HINT}"
    )
    assert (
        slide1["properties"]["role"]["description"]
        == f"Their role {fill_mod._TEXT_FORMATTING_HINT}"
    )

    slide2 = _leaf("slide_2")
    assert set(slide2["properties"]) == {"city"}
    assert (
        slide2["properties"]["city"]["description"]
        == f"The city {fill_mod._TEXT_FORMATTING_HINT}"
    )


def test_output_file_name_field_is_present_and_optional():
    deck = build_deck([("{{name}}", "{{name}}:\nThe name")])
    tool = _fill_tool(_ctx(schema_slides(deck)))

    schema = tool.args_schema.model_json_schema()
    assert "output_file_name" in schema["properties"]
    assert "output_file_name" not in schema.get("required", [])
    assert schema["properties"]["output_file_name"]["description"]


def test_image_leaf_is_optional_and_names_its_folder():
    """An image key's leaf is OPTIONAL and its description carries the note plus the
    source folder guidance; the how-to-browse procedure lives once in the main tool
    description (which references ``list_images_in_folder`` and insists MUST)."""
    deck = build_image_deck([(IMAGE_NOTES, ["{{logo}}"])])
    tool = _fill_tool(_ctx(image_schema(deck, slide=1, key="logo", tag_id="tag-logos")))

    schema = tool.args_schema.model_json_schema()
    defs = schema["$defs"]
    ref = schema["properties"]["slide_1"]["$ref"]
    leaf = defs[ref.split("/")[-1]]

    assert leaf.get("required", []) == []
    desc = leaf["properties"]["logo"]["description"]
    assert "The company logo" in desc
    assert "come from the folder" in desc
    assert "Brand/Logos" in desc
    # The generic image-handling procedure lives once in the main tool description.
    assert "list_images_in_folder" in tool.description
    assert "MUST" in tool.description


def test_no_tools_when_schema_is_empty():
    """No persisted schema (template never configured) -> nothing to fill -> no tools."""
    assert build_fill_tools(_ctx([])) == []


def test_no_list_tool_for_text_only_schema():
    deck = build_deck([("{{name}}", "{{name}}:\nThe name")])
    tools = build_fill_tools(_ctx(schema_slides(deck)))
    names = {t.name for t in tools}
    assert names == {"fill_ppt_template"}


def test_list_tool_present_when_schema_has_image_fields():
    deck = build_image_deck([(IMAGE_NOTES, ["{{logo}}"])])
    tools = build_fill_tools(
        _ctx(image_schema(deck, slide=1, key="logo", tag_id="tag-logos"))
    )
    names = {t.name for t in tools}
    assert names == {"fill_ppt_template", "list_images_in_folder"}


# --- B/C/D. Fill via shared traversal, upload, return artifact --------------------


@pytest.mark.asyncio
async def test_fill_repeated_keys_on_slide_uses_one_value(no_pdf):
    """A key appearing twice on one slide fills with the same value, and the deck is
    written session-scoped with a download LinkPart returned (preview disabled)."""
    deck = build_deck([("Title {{name}} ... footer {{name}}", "{{name}}:\nThe name")])
    workspace = FakeWorkspace()
    ctx = _ctx(
        schema_slides(deck),
        session_id="sess-9",
        template_key="custom.pptx",
        template_bytes=deck,
        workspace=workspace,
    )
    tool = _fill_tool(ctx)

    content, artifact = await tool.coroutine(slide_1={"name": "Jane Doe"})

    # Written to a SESSION-SCOPED outputs path.
    assert len(workspace.writes) == 1
    write = workspace.writes[0]
    assert write["path"] == "outputs/sess-9/filled_presentation.pptx"
    assert write["content_type"] == _PPTX_CONTENT_TYPE

    # Both occurrences filled with the same value, no leftover placeholders.
    filled = Presentation(io.BytesIO(write["content"]))
    text = "".join(
        run.text
        for shape in filled.slides[0].shapes
        if shape.has_text_frame
        for paragraph in shape.text_frame.paragraphs
        for run in paragraph.runs
    )
    assert text.count("Jane Doe") == 2
    assert not KEY_PATTERN.search(text)

    # The returned artifact carries a download LinkPart with sensible name/mime.
    assert isinstance(artifact, ToolInvocationResult)
    assert artifact.is_error is False
    assert len(artifact.ui_parts) == 1
    link = artifact.ui_parts[0]
    assert isinstance(link, LinkPart)
    assert link.kind == LinkKind.download
    assert link.mime == _PPTX_CONTENT_TYPE
    assert link.file_name == "filled_presentation.pptx"
    assert link.href == "/fs/download/outputs/sess-9/filled_presentation.pptx"
    assert isinstance(content, str) and content.strip()


@pytest.mark.asyncio
async def test_fetches_template_from_configured_key(no_pdf):
    deck = build_deck([("{{name}}", "{{name}}:\nThe name")])
    assets = FakeAssets({"custom.pptx": deck})
    ctx = _ctx(schema_slides(deck), template_key="custom.pptx", assets=assets)
    await _fill_tool(ctx).coroutine(slide_1={"name": "X"})
    # The template was fetched from the configured key (only the seeded blob exists).
    assert "custom.pptx" in assets.blobs


@pytest.mark.asyncio
async def test_default_template_key_is_used(no_pdf):
    deck = build_deck([("{{name}}", "{{name}}:\nThe name")])
    assets = FakeAssets({PPT_FILLER_TEMPLATE_KEY: deck})
    ctx = _ctx(schema_slides(deck), assets=assets)
    # Should not raise: it fetches the default key which is seeded.
    content, artifact = await _fill_tool(ctx).coroutine(slide_1={"name": "X"})
    assert artifact.is_error is False


@pytest.mark.asyncio
async def test_key_split_across_runs_is_filled(no_pdf):
    """A key straddling a run boundary is still filled (shared run-merging traversal)."""
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
    para = textbox.text_frame.paragraphs[0]
    for run_text in ("Hello {{na", "me}} world"):
        para.add_run().text = run_text
    slide.notes_slide.notes_text_frame.text = "{{name}}:\nThe name"
    buffer = io.BytesIO()
    presentation.save(buffer)
    deck = buffer.getvalue()

    workspace = FakeWorkspace()
    ctx = _ctx(schema_slides(deck), template_bytes=deck, workspace=workspace)
    await _fill_tool(ctx).coroutine(slide_1={"name": "Jane"})

    filled = Presentation(io.BytesIO(workspace.writes[0]["content"]))
    assert list_keys_on_slide(filled.slides[0]) == []
    text = "".join(
        run.text
        for shape in filled.slides[0].shapes
        if shape.has_text_frame
        for paragraph in shape.text_frame.paragraphs
        for run in paragraph.runs
    )
    assert "Jane" in text


@pytest.mark.asyncio
async def test_omitted_text_key_fills_empty_string(no_pdf):
    """A provided-but-None text key fills as empty string (placeholder removed)."""
    deck = build_deck(
        [("Hello {{name}} and {{role}}", "{{name}}:\nName\n{{role}}:\nRole")]
    )
    workspace = FakeWorkspace()
    ctx = _ctx(schema_slides(deck), template_bytes=deck, workspace=workspace)

    await _fill_tool(ctx).coroutine(slide_1={"name": "Ada", "role": None})

    filled = Presentation(io.BytesIO(workspace.writes[0]["content"]))
    text = "".join(
        run.text
        for shape in filled.slides[0].shapes
        if shape.has_text_frame
        for paragraph in shape.text_frame.paragraphs
        for run in paragraph.runs
    )
    assert "Ada" in text
    assert not KEY_PATTERN.search(text)  # {{role}} gone, replaced by ""


@pytest.mark.asyncio
async def test_fill_all_fields_across_slides_leaves_no_placeholders(no_pdf):
    deck = build_deck(
        [
            ("Hello {{name}}, the {{role}}", "{{name}}, {{role}}:\nName then role"),
            (
                "Repeated {{name}} and {{city}}",
                "{{name}}:\nName again\n{{city}}:\nThe city",
            ),
        ]
    )
    workspace = FakeWorkspace()
    ctx = _ctx(schema_slides(deck), template_bytes=deck, workspace=workspace)

    content, artifact = await _fill_tool(ctx).coroutine(
        slide_1={"name": "Alice", "role": "Engineer"},
        slide_2={"name": "Alice", "city": "Paris"},
    )

    assert artifact.is_error is False
    refilled = Presentation(io.BytesIO(workspace.writes[0]["content"]))
    for slide in refilled.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                assert not KEY_PATTERN.search(shape.text_frame.text)


@pytest.mark.asyncio
async def test_fill_without_session_uses_unscoped_path(no_pdf):
    deck = build_deck([("{{name}}", "{{name}}:\nThe name")])
    workspace = FakeWorkspace()
    ctx = _ctx(
        schema_slides(deck),
        session_id=None,
        template_bytes=deck,
        workspace=workspace,
    )
    await _fill_tool(ctx).coroutine(slide_1={"name": "Bob"})
    assert workspace.writes[0]["path"] == "outputs/filled_presentation.pptx"


@pytest.mark.asyncio
async def test_template_fetch_failure_returns_error_result(no_pdf):
    """A template fetch failure becomes an is_error tool result, not a raised exception."""
    deck = build_deck([("{{name}}", "{{name}}:\nThe name")])
    # assets does NOT contain the template key -> fetch raises inside the tool.
    ctx = _ctx(schema_slides(deck), assets=FakeAssets({}))

    content, artifact = await _fill_tool(ctx).coroutine(slide_1={"name": "X"})

    assert artifact.is_error is True
    assert artifact.ui_parts == ()
    assert "template" in content.lower()


# --- E. content_and_artifact through the ToolCall-dict ainvoke idiom --------------


@pytest.mark.asyncio
async def test_ainvoke_with_tool_call_dict_surfaces_artifact(no_pdf):
    """Driving the fill tool through ``ainvoke`` with a ToolCall dict (the runtime
    idiom) yields a message whose ``.artifact`` carries the download LinkPart."""
    deck = build_deck([("{{name}}", "{{name}}:\nThe name")])
    ctx = _ctx(schema_slides(deck), template_bytes=deck)
    tool = _fill_tool(ctx)

    message = await tool.ainvoke(
        {
            "type": "tool_call",
            "name": "fill_ppt_template",
            "args": {"slide_1": {"name": "Ada"}},
            "id": "call-1",
        }
    )

    artifact = message.artifact
    assert isinstance(artifact, ToolInvocationResult)
    assert artifact.is_error is False
    assert len(artifact.ui_parts) == 1
    assert isinstance(artifact.ui_parts[0], LinkPart)


# --- F. Model-chosen output file name (sanitization / traversal-safety) ------------


@pytest.mark.asyncio
async def test_output_file_name_is_used_for_path_and_download(no_pdf):
    deck = build_deck([("{{name}}", "{{name}}:\nThe name")])
    workspace = FakeWorkspace()
    ctx = _ctx(
        schema_slides(deck),
        session_id="sess-7",
        template_bytes=deck,
        workspace=workspace,
    )

    _content, artifact = await _fill_tool(ctx).coroutine(
        slide_1={"name": "Ada"}, output_file_name="Proposition ACME 2026"
    )

    assert workspace.writes[0]["path"] == "outputs/sess-7/Proposition ACME 2026.pptx"
    assert artifact.ui_parts[0].file_name == "Proposition ACME 2026.pptx"


@pytest.mark.asyncio
async def test_output_file_name_existing_extension_not_doubled(no_pdf):
    deck = build_deck([("{{name}}", "{{name}}:\nThe name")])
    workspace = FakeWorkspace()
    ctx = _ctx(
        schema_slides(deck), session_id="s", template_bytes=deck, workspace=workspace
    )
    await _fill_tool(ctx).coroutine(slide_1={"name": "X"}, output_file_name="Deck.PPTX")
    assert workspace.writes[0]["path"] == "outputs/s/Deck.pptx"


@pytest.mark.asyncio
async def test_output_file_name_strips_path_components(no_pdf):
    """Path separators are stripped (no traversal): only the final component is kept."""
    deck = build_deck([("{{name}}", "{{name}}:\nThe name")])
    workspace = FakeWorkspace()
    ctx = _ctx(
        schema_slides(deck), session_id="s", template_bytes=deck, workspace=workspace
    )
    await _fill_tool(ctx).coroutine(
        slide_1={"name": "X"}, output_file_name="../../etc/evil.pptx"
    )
    assert workspace.writes[0]["path"] == "outputs/s/evil.pptx"


@pytest.mark.asyncio
async def test_blank_output_file_name_falls_back_to_default(no_pdf):
    deck = build_deck([("{{name}}", "{{name}}:\nThe name")])
    workspace = FakeWorkspace()
    ctx = _ctx(
        schema_slides(deck), session_id="s", template_bytes=deck, workspace=workspace
    )
    await _fill_tool(ctx).coroutine(slide_1={"name": "X"}, output_file_name="   ")
    assert workspace.writes[0]["path"] == "outputs/s/filled_presentation.pptx"


# --- G. PDF preview (best-effort) --------------------------------------------------


@pytest.mark.asyncio
async def test_successful_fill_emits_ppt_preview_and_writes_pdf(
    monkeypatch, preview_union
):
    """A successful conversion writes the PDF beside the .pptx and emits a
    ppt_preview part with durable hrefs + a version, not a standalone download."""
    _enable_preview(monkeypatch)
    deck = build_deck([("{{name}}", "{{name}}:\nThe name")])
    workspace = FakeWorkspace()
    ctx = _ctx(
        schema_slides(deck),
        session_id="sess-9",
        template_bytes=deck,
        workspace=workspace,
    )

    _content, artifact = await _fill_tool(ctx).coroutine(slide_1={"name": "Jane"})

    assert artifact.is_error is False
    # Two writes: the .pptx then the preview .pdf, session-scoped sibling paths.
    assert len(workspace.writes) == 2
    pptx_write, pdf_write = workspace.writes
    assert pptx_write["path"] == "outputs/sess-9/filled_presentation.pptx"
    assert pdf_write["path"] == "outputs/sess-9/filled_presentation.pdf"
    assert pdf_write["content_type"] == "application/pdf"

    assert len(artifact.ui_parts) == 1
    part = artifact.ui_parts[0]
    assert isinstance(part, PptPreviewPart)
    assert not any(isinstance(p, LinkPart) for p in artifact.ui_parts)
    assert (
        part.pdf_download_url == "/fs/download/outputs/sess-9/filled_presentation.pdf"
    )
    assert part.pptx_download_url == (
        "/fs/download/outputs/sess-9/filled_presentation.pptx"
    )
    assert part.version  # a freshness token is stamped
    assert part.file_name == "filled_presentation.pptx"
    assert part.preview_id == "outputs/sess-9/filled_presentation.pptx"


@pytest.mark.asyncio
async def test_failed_conversion_returns_link_fallback_with_note(no_pdf):
    """When conversion returns None, the fill still returns the .pptx via a download
    LinkPart and the message states the preview could not be generated."""
    deck = build_deck([("{{name}}", "{{name}}:\nThe name")])
    workspace = FakeWorkspace()
    ctx = _ctx(schema_slides(deck), template_bytes=deck, workspace=workspace)

    content, artifact = await _fill_tool(ctx).coroutine(slide_1={"name": "Jane"})

    assert artifact.is_error is False
    assert len(workspace.writes) == 1  # only the .pptx, no PDF
    assert len(artifact.ui_parts) == 1
    assert isinstance(artifact.ui_parts[0], LinkPart)
    assert artifact.ui_parts[0].kind == LinkKind.download
    assert "preview" in content.lower()


@pytest.mark.asyncio
async def test_preview_version_changes_on_refill(monkeypatch, preview_union):
    deck = build_deck([("{{name}}", "{{name}}:\nThe name")])
    ctx = _ctx(schema_slides(deck), template_bytes=deck)
    tool = _fill_tool(ctx)

    _enable_preview(monkeypatch, pdf_bytes=b"%PDF-1.5 first")
    _c1, a1 = await tool.coroutine(slide_1={"name": "Jane"})
    v1 = a1.ui_parts[0]
    assert isinstance(v1, PptPreviewPart)

    _enable_preview(monkeypatch, pdf_bytes=b"%PDF-1.5 second-different")
    _c2, a2 = await tool.coroutine(slide_1={"name": "Bob"})
    v2 = a2.ui_parts[0]
    assert isinstance(v2, PptPreviewPart)

    assert v1.version != v2.version


# --- H. Notes: authoring guidance stripped; kept content preserved ----------------


def _slide_notes(presentation, index: int) -> str:
    slide = list(presentation.slides)[index]
    if not slide.has_notes_slide:
        return ""
    return slide.notes_slide.notes_text_frame.text or ""


@pytest.mark.asyncio
async def test_filled_deck_strips_authoring_notes(no_pdf):
    deck = build_deck([("Hello {{name}}", "{{name}}:\nThe person's name")])
    workspace = FakeWorkspace()
    ctx = _ctx(schema_slides(deck), template_bytes=deck, workspace=workspace)

    await _fill_tool(ctx).coroutine(slide_1={"name": "Ada"})

    filled = Presentation(io.BytesIO(workspace.writes[0]["content"]))
    notes = _slide_notes(filled, 0)
    assert notes == ""


@pytest.mark.asyncio
async def test_filled_deck_keeps_notes_after_separator(no_pdf):
    deck = build_deck(
        [
            (
                "Hello {{name}}",
                "{{name}}:\nThe person's name\n---\nSpeaker note: pause here.",
            )
        ]
    )
    workspace = FakeWorkspace()
    ctx = _ctx(schema_slides(deck), template_bytes=deck, workspace=workspace)

    await _fill_tool(ctx).coroutine(slide_1={"name": "Ada"})

    filled = Presentation(io.BytesIO(workspace.writes[0]["content"]))
    assert _slide_notes(filled, 0) == "Speaker note: pause here."


# --- I. Image keys: placement / removal / hard-fail refusals -----------------------


@pytest.mark.asyncio
async def test_image_key_with_doc_id_places_picture_and_removes_placeholder(no_pdf):
    deck = build_image_deck([(IMAGE_NOTES, ["{{logo}}"])])
    docs = FakeDocs({"doc-logo": (png_bytes(10, 10, "red"), "image/png")})
    workspace = FakeWorkspace()
    ctx = _ctx(
        image_schema(deck, slide=1, key="logo"),
        template_bytes=deck,
        docs=docs,
        workspace=workspace,
    )

    _content, artifact = await _fill_tool(ctx).coroutine(slide_1={"logo": "doc-logo"})

    assert artifact.is_error is False
    assert docs.fetch_uids == ["doc-logo"]
    filled = Presentation(io.BytesIO(workspace.writes[0]["content"]))
    slide = list(filled.slides)[0]
    assert len(picture_shapes(slide)) == 1
    assert list_keys_on_slide(slide) == []


@pytest.mark.asyncio
async def test_omitted_image_key_removes_placeholder_without_picture(no_pdf):
    deck = build_image_deck([(IMAGE_NOTES, ["{{logo}}"])])
    docs = FakeDocs({})  # no docs; an omitted key must never fetch
    workspace = FakeWorkspace()
    ctx = _ctx(
        image_schema(deck, slide=1, key="logo"),
        template_bytes=deck,
        docs=docs,
        workspace=workspace,
    )

    _content, artifact = await _fill_tool(ctx).coroutine(slide_1={"logo": None})

    assert artifact.is_error is False
    assert docs.fetch_uids == []
    filled = Presentation(io.BytesIO(workspace.writes[0]["content"]))
    slide = list(filled.slides)[0]
    assert picture_shapes(slide) == []
    assert list_keys_on_slide(slide) == []


@pytest.mark.asyncio
async def test_image_key_with_path_instead_of_id_hard_fails(no_pdf):
    deck = build_image_deck([(IMAGE_NOTES, ["{{logo}}"])])
    docs = FakeDocs({})
    workspace = FakeWorkspace()
    ctx = _ctx(
        image_schema(deck, slide=1, key="logo"),
        template_bytes=deck,
        docs=docs,
        workspace=workspace,
    )

    content, artifact = await _fill_tool(ctx).coroutine(
        slide_1={"logo": "Brand/Logos/logo.png"}
    )

    assert artifact.is_error is True
    assert artifact.ui_parts == ()
    lowered = content.lower()
    assert "path" in lowered
    assert "document id" in lowered
    assert "list_images_in_folder" in content
    assert "Brand/Logos/logo.png" in content
    # Detected up front: no fetch, no upload.
    assert docs.fetch_uids == []
    assert workspace.writes == []


@pytest.mark.asyncio
async def test_image_fetch_failure_hard_fails(no_pdf):
    deck = build_image_deck([(IMAGE_NOTES, ["{{logo}}"])])
    docs = FakeDocs({})  # requested uid absent -> fetch raises
    workspace = FakeWorkspace()
    ctx = _ctx(
        image_schema(deck, slide=1, key="logo"),
        template_bytes=deck,
        docs=docs,
        workspace=workspace,
    )

    content, artifact = await _fill_tool(ctx).coroutine(slide_1={"logo": "missing-doc"})

    assert artifact.is_error is True
    assert "missing-doc" in content
    assert workspace.writes == []


@pytest.mark.asyncio
async def test_bad_image_bytes_hard_fails(no_pdf):
    deck = build_image_deck([(IMAGE_NOTES, ["{{logo}}"])])
    docs = FakeDocs({"doc-logo": (b"not an image", "image/png")})
    workspace = FakeWorkspace()
    ctx = _ctx(
        image_schema(deck, slide=1, key="logo"),
        template_bytes=deck,
        docs=docs,
        workspace=workspace,
    )

    content, artifact = await _fill_tool(ctx).coroutine(slide_1={"logo": "doc-logo"})

    assert artifact.is_error is True
    assert artifact.ui_parts == ()
    assert "image" in content.lower()
    assert workspace.writes == []


@pytest.mark.asyncio
async def test_webp_image_is_transcoded_and_embedded(no_pdf):
    deck = build_image_deck([(IMAGE_NOTES, ["{{logo}}"])])
    docs = FakeDocs({"doc-logo": (webp_bytes(10, 10, "red"), "image/webp")})
    workspace = FakeWorkspace()
    ctx = _ctx(
        image_schema(deck, slide=1, key="logo"),
        template_bytes=deck,
        docs=docs,
        workspace=workspace,
    )

    _content, artifact = await _fill_tool(ctx).coroutine(slide_1={"logo": "doc-logo"})

    assert artifact.is_error is False
    filled = Presentation(io.BytesIO(workspace.writes[0]["content"]))
    slide = list(filled.slides)[0]
    assert len(picture_shapes(slide)) == 1
    assert list_keys_on_slide(slide) == []


@pytest.mark.asyncio
async def test_image_key_in_table_cell_is_refused(no_pdf):
    """An image key in a table cell (invalid picture location) hard-fails the fill."""
    deck = build_table_deck(IMAGE_NOTES, ["{{logo}}", "plain"])
    docs = FakeDocs({"doc-logo": (png_bytes(10, 10, "red"), "image/png")})
    workspace = FakeWorkspace()
    ctx = _ctx(
        image_schema(deck, slide=1, key="logo"),
        template_bytes=deck,
        docs=docs,
        workspace=workspace,
    )

    content, artifact = await _fill_tool(ctx).coroutine(slide_1={"logo": "doc-logo"})

    assert artifact.is_error is True
    assert "table cell" in content.lower()
    assert workspace.writes == []


@pytest.mark.asyncio
async def test_mixed_text_and_image_keys_fill_independently(no_pdf):
    deck = build_image_deck(
        [("{{title}}:\nThe title\n" + IMAGE_NOTES, ["Title: {{title}}", "{{logo}}"])]
    )
    docs = FakeDocs({"doc-logo": (png_bytes(20, 10, "red"), "image/png")})
    workspace = FakeWorkspace()
    ctx = _ctx(
        image_schema(deck, slide=1, key="logo"),
        template_bytes=deck,
        docs=docs,
        workspace=workspace,
    )

    _content, artifact = await _fill_tool(ctx).coroutine(
        slide_1={"title": "Quarterly", "logo": "doc-logo"}
    )

    assert artifact.is_error is False
    filled = Presentation(io.BytesIO(workspace.writes[0]["content"]))
    slide = list(filled.slides)[0]
    assert len(picture_shapes(slide)) == 1
    assert list_keys_on_slide(slide) == []
    body = "".join(
        run.text
        for shape in slide.shapes
        if shape.has_text_frame
        for paragraph in shape.text_frame.paragraphs
        for run in paragraph.runs
    )
    assert "Quarterly" in body


# --- J. Fail-loud when required ports are missing ----------------------------------


@pytest.mark.asyncio
async def test_missing_agent_assets_raises_runtime_error(no_pdf):
    deck = build_deck([("{{name}}", "{{name}}:\nThe name")])
    ctx = CapabilityContext(
        identity=CapabilityIdentity(user_id="u-1", session_id="s"),
        config=PptFillerConfig(schema_slides=schema_slides(deck)),
        turn_options=EmptyModel(),
        services=RuntimeServices(agent_assets=None, workspace_fs=FakeWorkspace()),
    )
    with pytest.raises(RuntimeError):
        await _fill_tool(ctx).coroutine(slide_1={"name": "X"})


@pytest.mark.asyncio
async def test_missing_workspace_fs_raises_runtime_error(no_pdf):
    deck = build_deck([("{{name}}", "{{name}}:\nThe name")])
    ctx = CapabilityContext(
        identity=CapabilityIdentity(user_id="u-1", session_id="s"),
        config=PptFillerConfig(schema_slides=schema_slides(deck)),
        turn_options=EmptyModel(),
        services=RuntimeServices(
            agent_assets=FakeAssets({PPT_FILLER_TEMPLATE_KEY: deck}), workspace_fs=None
        ),
    )
    with pytest.raises(RuntimeError):
        await _fill_tool(ctx).coroutine(slide_1={"name": "X"})


@pytest.mark.asyncio
async def test_missing_document_content_for_image_fails_the_fill(no_pdf):
    """A missing ``document_content`` port fails LOUD like the other port guards.

    The guard runs BEFORE the fill try-block (next to the ``agent_assets`` /
    ``workspace_fs`` guards) whenever the schema carries image fields, so a
    misconfigured execution path propagates a ``RuntimeError`` instead of
    degrading to a soft tool error.
    """
    deck = build_image_deck([(IMAGE_NOTES, ["{{logo}}"])])
    workspace = FakeWorkspace()
    ctx = _ctx(
        image_schema(deck, slide=1, key="logo"),
        template_bytes=deck,
        docs=None,  # image schema but no content port
        workspace=workspace,
    )

    with pytest.raises(RuntimeError, match="document_content"):
        await _fill_tool(ctx).coroutine(slide_1={"logo": "doc-logo"})

    assert workspace.writes == []


# --- K. list_images_in_folder companion tool --------------------------------------


def _list_tool(ctx):
    tools = build_fill_tools(ctx)
    listing = [t for t in tools if t.name == "list_images_in_folder"]
    assert listing, "expected the list_images_in_folder tool"
    return listing[0]


@pytest.mark.asyncio
async def test_list_images_known_folder_returns_json():
    deck = build_image_deck([(IMAGE_NOTES, ["{{logo}}"])])
    folders = FakeFolders(
        tag_to_docs={"tag-logos": [("doc-1", "Logo A"), ("doc-2", "Logo B")]}
    )
    ctx = _ctx(
        image_schema(deck, slide=1, key="logo", tag_id="tag-logos"),
        template_bytes=deck,
        folders=folders,
    )

    result = await _list_tool(ctx).coroutine(folder="Brand/Logos")

    assert json.loads(result) == [
        {"document_id": "doc-1", "name": "Logo A"},
        {"document_id": "doc-2", "name": "Logo B"},
    ]
    assert folders.list_calls == ["tag-logos"]


@pytest.mark.asyncio
async def test_list_images_unknown_folder_returns_message():
    deck = build_image_deck([(IMAGE_NOTES, ["{{logo}}"])])
    folders = FakeFolders(tag_to_docs={"tag-logos": [("doc-1", "Logo A")]})
    ctx = _ctx(
        image_schema(deck, slide=1, key="logo", tag_id="tag-logos"),
        template_bytes=deck,
        folders=folders,
    )

    result = await _list_tool(ctx).coroutine(folder="Nope/Missing")

    assert "Unknown image folder" in result
    assert "Brand/Logos" in result  # names the known folders
    assert folders.list_calls == []  # never listed an unknown folder


@pytest.mark.asyncio
async def test_list_images_empty_folder_returns_message():
    deck = build_image_deck([(IMAGE_NOTES, ["{{logo}}"])])
    folders = FakeFolders(tag_to_docs={"tag-logos": []})  # resolves but empty
    ctx = _ctx(
        image_schema(deck, slide=1, key="logo", tag_id="tag-logos"),
        template_bytes=deck,
        folders=folders,
    )

    result = await _list_tool(ctx).coroutine(folder="Brand/Logos")

    assert "empty" in result.lower()
    assert folders.list_calls == ["tag-logos"]


# --- Instructions overlay (system-prompt fragment) ---------------------------------


class _FakeModelRequest:
    """Duck-typed stand-in for langchain's ModelRequest: the middleware only
    reads ``system_prompt`` and calls ``override(system_message=...)``."""

    def __init__(self, system_prompt):
        self.system_prompt = system_prompt
        self.system_message = None

    def override(self, *, system_message):
        self.system_message = system_message
        return self


async def _passthrough(request):
    return request


@pytest.mark.asyncio
async def test_middleware_overlays_fill_instructions_when_template_configured():
    from fred_capability_ppt_filler.capability import PptFillerCapability

    deck = build_deck([("{{name}}", "{{name}}:\nThe name")])
    (mw,) = PptFillerCapability().middleware(_ctx(schema_slides(deck)))
    request = _FakeModelRequest("BASE PROMPT")

    await mw.awrap_model_call(request, _passthrough)

    assert request.system_message is not None
    merged = request.system_message.content
    assert merged.startswith("BASE PROMPT")
    assert "fill_ppt_template" in merged
    assert "NEVER ask the user to upload" in merged


@pytest.mark.asyncio
async def test_middleware_stays_silent_without_a_configured_template():
    from fred_capability_ppt_filler.capability import PptFillerCapability

    (mw,) = PptFillerCapability().middleware(_ctx([]))
    request = _FakeModelRequest("BASE PROMPT")

    await mw.awrap_model_call(request, _passthrough)

    assert request.system_message is None
