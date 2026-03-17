# Copyright Thales 2025
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.

from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

from knowledge_flow_backend.core.processors.input.pptx_markdown_processor.pptx_markdown_processor import (
    PptxMarkdownProcessor,
)
from knowledge_flow_backend.core.processors.input.pptx_markdown_processor.pptx_markdown_processor_v2 import (
    PptxMarkdownProcessorV2,
)


def _blank_layout(presentation: Presentation):
    # `slide_layouts[6]` is usually blank in default templates.
    # Fallback keeps tests robust if template indices differ.
    return presentation.slide_layouts[6] if len(presentation.slide_layouts) > 6 else presentation.slide_layouts[-1]


def _build_pptx(path: Path, slides: list[list[str]]) -> Path:
    presentation = Presentation()
    for slide_lines in slides:
        slide = presentation.slides.add_slide(_blank_layout(presentation))
        for idx, text in enumerate(slide_lines):
            top = Inches(1.0 + (idx * 0.9))
            textbox = slide.shapes.add_textbox(Inches(1.0), top, Inches(9.0), Inches(0.7))
            textbox.text_frame.text = text
    presentation.save(str(path))
    return path


@pytest.fixture
def v1_processor() -> PptxMarkdownProcessor:
    return PptxMarkdownProcessor()


@pytest.fixture
def v2_processor() -> PptxMarkdownProcessorV2:
    return PptxMarkdownProcessorV2()


def test_v1_validity_and_metadata(v1_processor: PptxMarkdownProcessor, tmp_path: Path):
    pptx_path = _build_pptx(tmp_path / "sample_v1.pptx", [["Title", "Body"]])

    assert v1_processor.check_file_validity(pptx_path) is True
    metadata = v1_processor.extract_file_metadata(pptx_path)
    assert metadata["document_name"] == "sample_v1.pptx"
    assert metadata["num_slides"] == 1


def test_v1_convert_generates_markdown(v1_processor: PptxMarkdownProcessor, tmp_path: Path):
    pptx_path = _build_pptx(
        tmp_path / "sample_v1_convert.pptx",
        [
            ["Slide 1 title", "Slide 1 body"],
            ["Slide 2 title", "Slide 2 body"],
        ],
    )
    output_dir = tmp_path / "out_v1"

    result = v1_processor.convert_file_to_markdown(pptx_path, output_dir, "doc-v1")
    md_path = Path(result["md_file"])
    md = md_path.read_text(encoding="utf-8")

    assert md_path.exists()
    assert md.count("### Slide") == 2
    assert "Slide 1 title" in md
    assert "Slide 2 body" in md
    assert "---" in md


def test_v1_convert_handles_no_text(v1_processor: PptxMarkdownProcessor, tmp_path: Path):
    presentation = Presentation()
    _ = presentation.slides.add_slide(_blank_layout(presentation))
    pptx_path = tmp_path / "sample_v1_empty.pptx"
    presentation.save(str(pptx_path))

    result = v1_processor.convert_file_to_markdown(pptx_path, tmp_path / "out_v1_empty", "doc-v1-empty")
    md = Path(result["md_file"]).read_text(encoding="utf-8")

    assert md.strip() == "*No extractable text*"


def test_v1_invalid_file(v1_processor: PptxMarkdownProcessor, tmp_path: Path):
    invalid = tmp_path / "broken.pptx"
    invalid.write_text("not a real pptx", encoding="utf-8")
    assert v1_processor.check_file_validity(invalid) is False


def test_v2_validity_and_metadata(v2_processor: PptxMarkdownProcessorV2, tmp_path: Path):
    pptx_path = _build_pptx(tmp_path / "sample_v2.pptx", [["MAIN TITLE", "Body paragraph."]])

    assert v2_processor.check_file_validity(pptx_path) is True
    metadata = v2_processor.extract_file_metadata(pptx_path)
    assert metadata["document_name"] == "sample_v2.pptx"
    assert metadata["num_slides"] == 1


def test_v2_convert_generates_structured_markdown(v2_processor: PptxMarkdownProcessorV2, tmp_path: Path):
    pptx_path = _build_pptx(
        tmp_path / "sample_v2_convert.pptx",
        [
            ["MAIN TITLE", "Body paragraph."],
            ["SECOND TITLE", "Another body line."],
        ],
    )
    output_dir = tmp_path / "out_v2"

    result = v2_processor.convert_file_to_markdown(pptx_path, output_dir, "doc-v2")
    md_path = Path(result["md_file"])
    md = md_path.read_text(encoding="utf-8")

    assert md_path.exists()
    assert "## Slide 1" in md
    assert "## Slide 2" in md
    assert "MAIN TITLE" in md
    assert "Another body line." in md
    assert "---" in md


def test_v2_convert_handles_no_text(v2_processor: PptxMarkdownProcessorV2, tmp_path: Path):
    presentation = Presentation()
    _ = presentation.slides.add_slide(_blank_layout(presentation))
    pptx_path = tmp_path / "sample_v2_empty.pptx"
    presentation.save(str(pptx_path))

    result = v2_processor.convert_file_to_markdown(pptx_path, tmp_path / "out_v2_empty", "doc-v2-empty")
    md = Path(result["md_file"]).read_text(encoding="utf-8")

    assert "## Slide 1" in md
    assert "*No extractable text*" in md


def test_v2_invalid_file(v2_processor: PptxMarkdownProcessorV2, tmp_path: Path):
    invalid = tmp_path / "broken_v2.pptx"
    invalid.write_text("not a real pptx", encoding="utf-8")
    assert v2_processor.check_file_validity(invalid) is False
