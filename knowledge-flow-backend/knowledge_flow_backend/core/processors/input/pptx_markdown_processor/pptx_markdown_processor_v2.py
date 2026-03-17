# Copyright Thales 2025
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.

import logging
from dataclasses import dataclass
from pathlib import Path
from statistics import median
from typing import Any

from pptx import Presentation

from knowledge_flow_backend.application_context import get_configuration
from knowledge_flow_backend.common.processing_profile_context import get_current_processing_profile
from knowledge_flow_backend.common.structures import IngestionProcessingProfile, ProcessingConfig
from knowledge_flow_backend.core.processors.input.common.base_input_processor import BaseMarkdownProcessor, InputConversionError

logger = logging.getLogger(__name__)


@dataclass(slots=True)
class _SlideTextBlock:
    # We keep one normalized structure per text shape so all later heuristics
    # read the same data and do not repeatedly hit python-pptx objects.
    text: str
    lines: list[str]
    top: int
    left: int
    width: int
    height: int
    max_font_pt: float | None
    uppercase_ratio: float

    @property
    def right(self) -> int:
        return self.left + self.width

    @property
    def bottom(self) -> int:
        return self.top + self.height

    @property
    def center_x(self) -> float:
        return self.left + (self.width / 2)


def _normalize_space(value: str) -> str:
    # Normalize line breaks/tabs early so matching/grouping is stable
    # across exporters and template variants.
    return " ".join(value.replace("\r", " ").replace("\n", " ").replace("\t", " ").split()).strip()


def _uppercase_ratio(value: str) -> float:
    # Uppercase ratio is a language-agnostic signal for "banner/header"-like text.
    # It avoids hardcoding deck-specific words.
    letters = [ch for ch in value if ch.isalpha()]
    if not letters:
        return 0.0
    return sum(1 for ch in letters if ch.isupper()) / len(letters)


def _extract_shape_lines(shape) -> list[str]:
    # Some templates store text in plain .text, others in text_frame paragraphs.
    # We normalize both paths here so downstream logic stays generic.
    if not bool(getattr(shape, "has_text_frame", False)):
        raw_text = _normalize_space(str(getattr(shape, "text", "") or ""))
        return [raw_text] if raw_text else []

    text_frame = getattr(shape, "text_frame", None)
    if text_frame is None:
        return []

    # Detect whether this text frame represents a list structure. In python-pptx,
    # para.level == 0 is used both for non-bulleted paragraphs and top-level
    # bullet paragraphs. If any paragraph has level > 0, we treat the whole
    # text frame as a list so that top-level items are preserved as list items
    # instead of being flattened into plain text.
    para_levels = [
        int(getattr(p, "level", 0) or 0) for p in getattr(text_frame, "paragraphs", [])
    ]
    has_any_level_gt0 = any(lvl > 0 for lvl in para_levels)

    lines: list[str] = []
    for para in text_frame.paragraphs:
        text = _normalize_space(str(getattr(para, "text", "") or ""))
        if not text:
            continue
        level = int(getattr(para, "level", 0) or 0)
        # Preserve explicit paragraph nesting when present.
        # This prevents flattening real bullet structures.
        if level > 0:
            lines.append(f"{'  ' * (level - 1)}- {text}")
        elif has_any_level_gt0:
            # Treat level-0 paragraphs as top-level list items when the text
            # frame contains nested levels, so list structure is retained.
            lines.append(f"- {text}")
        else:
            lines.append(text)
    return lines


def _extract_shape_max_font_pt(shape) -> float | None:
    # Font size is one of the few robust signals available in PPTX text-only parsing.
    # We use the maximum size in the block to detect likely headers.
    if not bool(getattr(shape, "has_text_frame", False)):
        return None
    text_frame = getattr(shape, "text_frame", None)
    if text_frame is None:
        return None

    font_sizes: list[float] = []
    for para in text_frame.paragraphs:
        for run in para.runs:
            run_font = getattr(run, "font", None)
            run_size = getattr(run_font, "size", None)
            if run_size is not None:
                font_sizes.append(float(run_size.pt))
        para_font = getattr(para, "font", None)
        para_size = getattr(para_font, "size", None)
        if para_size is not None:
            font_sizes.append(float(para_size.pt))

    if not font_sizes:
        return None
    return max(font_sizes)


def _collect_slide_blocks(slide) -> list[_SlideTextBlock]:
    # Convert raw slide shapes into normalized text blocks with geometry.
    # All higher-level logic operates on this generic intermediate form.
    blocks: list[_SlideTextBlock] = []
    for shape in slide.shapes:
        lines = _extract_shape_lines(shape)
        if not lines:
            continue
        text = _normalize_space(" ".join(lines))
        if not text:
            continue

        block = _SlideTextBlock(
            text=text,
            lines=lines,
            top=int(getattr(shape, "top", 0) or 0),
            left=int(getattr(shape, "left", 0) or 0),
            width=int(getattr(shape, "width", 0) or 0),
            height=int(getattr(shape, "height", 0) or 0),
            max_font_pt=_extract_shape_max_font_pt(shape),
            uppercase_ratio=_uppercase_ratio(text),
        )
        blocks.append(block)

    blocks.sort(key=lambda b: (b.top, b.left))
    # Default reading order baseline (top-to-bottom, left-to-right).
    return blocks


def _pick_slide_title(
    blocks: list[_SlideTextBlock],
    *,
    slide_width: int,
    slide_height: int,
    options: ProcessingConfig.PptxPipelineConfig,
) -> _SlideTextBlock | None:
    # Title selection is score-based (position + width + font + concise bonus),
    # not keyword-based, so it works across different branding/templates.
    if not blocks:
        return None

    top_candidates = [b for b in blocks if b.top <= int(slide_height * options.title_top_ratio)]
    candidates = top_candidates or blocks

    def score(block: _SlideTextBlock) -> float:
        # Top blocks are preferred because most decks place titles there.
        # Width/font act as secondary signals when top area is crowded.
        top_score = 1.0 - (block.top / max(1, slide_height))
        width_score = block.width / max(1, slide_width)
        font_score = (block.max_font_pt or 0.0) / 36.0
        concise_bonus = 0.40 if len(block.text) <= options.title_concise_max_chars else 0.0
        return (top_score * 2.0) + width_score + font_score + concise_bonus

    return max(candidates, key=score)


def _looks_like_section_header(
    block: _SlideTextBlock,
    *,
    median_font_pt: float | None,
    slide_height: int,
    options: ProcessingConfig.PptxPipelineConfig,
) -> bool:
    # Header detection deliberately combines several weak signals:
    # short text, position, capitalization, font size.
    # This is more generic than depending on exact wording.
    text_len = len(block.text)
    if text_len == 0 or text_len > options.header_max_text_length:
        return False
    if len(block.text.split()) < options.header_min_word_count:
        return False
    if any(block.text.endswith(punct) for punct in (".", "!", "?", ";")):
        return False

    is_short_block = len(block.lines) <= options.header_max_lines
    looks_upper = block.uppercase_ratio >= options.header_uppercase_ratio_threshold
    looks_large = median_font_pt is not None and block.max_font_pt is not None and block.max_font_pt >= (median_font_pt + options.header_min_font_delta_pt)
    near_top = block.top <= int(slide_height * options.header_top_ratio)
    return near_top and is_short_block and (looks_upper or looks_large)


def _looks_like_short_label(block: _SlideTextBlock, *, options: ProcessingConfig.PptxPipelineConfig) -> bool:
    # Short labels are used for row-like process steps (e.g., 5 verbs across top row).
    # We keep constraints simple and language-neutral.
    token_count = len(block.text.split())
    if token_count == 0 or token_count > options.short_label_max_words:
        return False
    if len(block.text) > options.short_label_max_chars:
        return False
    if any(ch in block.text for ch in ".!?;:"):
        return False
    return True


def _horizontal_overlap_ratio(a: _SlideTextBlock, b: _SlideTextBlock) -> float:
    # Horizontal overlap helps attach content to the correct section in multi-column slides.
    # It is geometric, so it stays template-agnostic.
    overlap = max(0, min(a.right, b.right) - max(a.left, b.left))
    smallest_width = max(1, min(a.width, b.width))
    return overlap / smallest_width


def _sort_blocks_by_visual_rows(
    blocks: list[_SlideTextBlock],
    *,
    slide_height: int,
    options: ProcessingConfig.PptxPipelineConfig,
) -> list[_SlideTextBlock]:
    # Optional row clustering exists for decks with "same-row section headers".
    # When disabled, we keep pure geometric order to avoid reordering regressions.
    if len(blocks) <= 1:
        return blocks

    # Default-safe mode: no row clustering, just geometric reading order.
    if options.section_row_grouping_ratio <= 0:
        return sorted(blocks, key=lambda b: (b.top, b.left))

    row_threshold = max(1, int(slide_height * options.section_row_grouping_ratio))
    rows: list[list[_SlideTextBlock]] = []
    for block in sorted(blocks, key=lambda b: (b.top, b.left)):
        matched = False
        for row in rows:
            row_top = int(sum(item.top for item in row) / len(row))
            if abs(block.top - row_top) <= row_threshold:
                row.append(block)
                matched = True
                break
        if not matched:
            rows.append([block])

    ordered: list[_SlideTextBlock] = []
    rows.sort(key=lambda row: int(sum(item.top for item in row) / len(row)))
    for row in rows:
        ordered.extend(sorted(row, key=lambda b: b.left))
    return ordered


def _assign_to_section(
    block: _SlideTextBlock,
    headers: list[_SlideTextBlock],
    slide_width: int,
    slide_height: int,
    options: ProcessingConfig.PptxPipelineConfig,
) -> int | None:
    # Assign a content block to the nearest plausible section header.
    # Distance + overlap is safer than simple "previous header wins".
    best_index: int | None = None
    best_score: float | None = None
    best_overlap = 0.0
    best_center_gap = 1.0

    for idx, header in enumerate(headers):
        vertical_gap = block.top - header.bottom
        # Reject impossible matches first (too far above or below).
        if vertical_gap < -int(slide_height * options.section_backward_tolerance_ratio):
            continue
        if vertical_gap > int(slide_height * options.section_forward_tolerance_ratio):
            continue

        overlap_ratio = _horizontal_overlap_ratio(block, header)
        center_gap = abs(block.center_x - header.center_x) / max(1, slide_width)
        vertical_penalty = vertical_gap / max(1, slide_height) if vertical_gap >= 0 else (abs(vertical_gap) / max(1, slide_height)) + 0.35
        overlap_penalty = 0.0 if overlap_ratio > 0 else 0.75
        score = (2.2 * vertical_penalty) + center_gap + overlap_penalty

        if best_score is None or score < best_score:
            best_score = score
            best_index = idx
            best_overlap = overlap_ratio
            best_center_gap = center_gap

    if best_score is None:
        return None
    # Extra guard: if there is no horizontal overlap, require center proximity.
    # This prevents cross-column mis-attachments.
    if best_overlap <= 0 and best_center_gap > options.section_nonoverlap_max_center_gap_ratio:
        return None
    return best_index if best_score <= options.section_score_threshold else None


def _extract_horizontal_label_row(
    blocks: list[_SlideTextBlock],
    *,
    slide_width: int,
    slide_height: int,
    options: ProcessingConfig.PptxPipelineConfig,
) -> list[_SlideTextBlock]:
    # Detect a row of short labels used as phase/step banners.
    # We emit them as one "Key Steps" line to reduce markdown noise.
    if not options.detect_horizontal_label_row:
        return []
    if len(blocks) < options.horizontal_label_min_labels:
        return []

    top_band_limit = int(slide_height * options.horizontal_label_top_band_ratio)
    candidates = [b for b in blocks if b.top <= top_band_limit and _looks_like_short_label(b, options=options)]
    if len(candidates) < options.horizontal_label_min_labels:
        return []

    y_threshold = max(1, int(slide_height * options.horizontal_label_row_y_tolerance_ratio))
    rows: list[list[_SlideTextBlock]] = []
    for block in sorted(candidates, key=lambda b: (b.top, b.left)):
        matched = False
        for row in rows:
            row_top = int(sum(item.top for item in row) / len(row))
            if abs(block.top - row_top) <= y_threshold:
                row.append(block)
                matched = True
                break
        if not matched:
            rows.append([block])

    if not rows:
        return []

    best_row: list[_SlideTextBlock] = []
    best_score = -1.0
    for row in rows:
        if len(row) < options.horizontal_label_min_labels:
            continue
        row_sorted = sorted(row, key=lambda b: b.left)
        span = row_sorted[-1].right - row_sorted[0].left
        span_ratio = span / max(1, slide_width)
        row_score = len(row_sorted) + span_ratio
        if span_ratio >= options.horizontal_label_min_span_ratio and row_score > best_score:
            best_score = row_score
            best_row = row_sorted

    return best_row


def _render_section_items(blocks: list[_SlideTextBlock]) -> list[str]:
    # Convert raw blocks into readable markdown bullets/paragraphs.
    # Heuristic: long/sentence-like blocks stay as paragraphs; short ones become bullets.
    lines: list[str] = []
    for block in blocks:
        # If the shape already contains explicit markdown bullets, preserve them line by line.
        if any(part.lstrip().startswith("- ") for part in block.lines):
            for part in block.lines:
                if part.lstrip().startswith("- "):
                    lines.append(part.strip())
                else:
                    lines.append(f"- {_normalize_space(part)}")
            continue

        # Normalize each original line separately so we can preserve line/paragraph boundaries.
        normalized_lines = []
        for part in block.lines:
            normalized = _normalize_space(part)
            if normalized:
                normalized_lines.append(normalized)

        if not normalized_lines:
            continue

        # Use a space-joined version for sentence detection, but keep newlines in the output text.
        content = " ".join(normalized_lines)
        paragraph_text = "\n".join(normalized_lines)

        looks_like_sentence = len(content) > 120 or any(
            token in content for token in (". ", "? ", "! ", "; ", ": ")
        )
        if looks_like_sentence:
            lines.append(paragraph_text)
        else:
            lines.append(f"- {paragraph_text}")
    return lines


def _render_slide_markdown(
    slide_index: int,
    slide,
    slide_width: int,
    slide_height: int,
    options: ProcessingConfig.PptxPipelineConfig,
) -> str:
    # End-to-end slide renderer:
    # 1) normalize blocks
    # 2) pick title
    # 3) detect section headers
    # 4) assign remaining blocks
    # 5) render markdown
    #
    # Kept deterministic so identical inputs/configs produce stable output.
    blocks = _collect_slide_blocks(slide)
    if not blocks:
        return f"## Slide {slide_index}\n\n*No extractable text*"

    title_block = _pick_slide_title(
        blocks,
        slide_width=slide_width,
        slide_height=slide_height,
        options=options,
    )
    title_text = title_block.text if title_block else ""
    remaining = [block for block in blocks if block is not title_block]

    font_sizes = [block.max_font_pt for block in remaining if block.max_font_pt is not None]
    median_font_pt = float(median(font_sizes)) if font_sizes else None

    section_headers = [
        block
        for block in remaining
        if _looks_like_section_header(
            block,
            median_font_pt=median_font_pt,
            slide_height=slide_height,
            options=options,
        )
    ]
    section_headers = _sort_blocks_by_visual_rows(section_headers, slide_height=slide_height, options=options)

    section_items: list[list[_SlideTextBlock]] = [[] for _ in section_headers]
    intro_blocks: list[_SlideTextBlock] = []
    orphan_blocks: list[_SlideTextBlock] = []
    header_ids = {id(header) for header in section_headers}
    horizontal_label_row = _extract_horizontal_label_row(
        [block for block in remaining if id(block) not in header_ids],
        slide_width=slide_width,
        slide_height=slide_height,
        options=options,
    )
    horizontal_label_ids = {id(block) for block in horizontal_label_row}

    for block in remaining:
        if id(block) in header_ids:
            continue
        if id(block) in horizontal_label_ids:
            continue
        if not section_headers:
            intro_blocks.append(block)
            continue
        section_index = _assign_to_section(
            block,
            section_headers,
            slide_width=slide_width,
            slide_height=slide_height,
            options=options,
        )
        if section_index is None:
            # Unassigned top-area blocks are kept as intro context.
            # Lower unassigned blocks are treated as orphan content.
            if block.top <= int(slide_height * options.intro_top_ratio):
                intro_blocks.append(block)
            else:
                orphan_blocks.append(block)
        else:
            section_items[section_index].append(block)

    lines: list[str] = []
    heading = f"## Slide {slide_index}"
    if title_text:
        heading += f": {title_text}"
    lines.append(heading)
    lines.append("")

    if horizontal_label_row:
        # Compact rendering avoids repeating many tiny labels as separate bullets.
        labels = ", ".join(block.text for block in sorted(horizontal_label_row, key=lambda b: b.left))
        lines.append(f"**Key Steps:** {labels}")
        lines.append("")

    if intro_blocks:
        lines.extend(_render_section_items(sorted(intro_blocks, key=lambda b: (b.top, b.left))))
        lines.append("")

    for idx, header in enumerate(section_headers):
        lines.append(f"### {_normalize_space(' '.join(header.lines))}")
        rendered_items = _render_section_items(sorted(section_items[idx], key=lambda b: (b.top, b.left)))
        lines.extend(rendered_items or ["*No details extracted*"])
        lines.append("")

    if orphan_blocks:
        orphan_sorted = sorted(orphan_blocks, key=lambda b: (b.top, b.left))
        # "Highlights" catches short header-like leftovers (common for callouts/badges).
        # Everything else is preserved under Additional Content, so we do not silently drop text.
        callout_blocks = [
            block
            for block in orphan_sorted
            if _looks_like_section_header(
                block,
                median_font_pt=median_font_pt,
                slide_height=slide_height,
                options=options,
            )
            and len(block.text.split()) <= options.callout_max_words
        ]
        callout_ids = {id(block) for block in callout_blocks}
        additional_blocks = [block for block in orphan_sorted if id(block) not in callout_ids]

        if callout_blocks:
            lines.append("### Highlights")
            lines.extend(f"- {block.text}" for block in callout_blocks)
            lines.append("")

        if additional_blocks:
            lines.append("### Additional Content")
            lines.extend(_render_section_items(additional_blocks))
            lines.append("")

    while lines and lines[-1] == "":
        lines.pop()
    return "\n".join(lines)


class PptxMarkdownProcessorV2(BaseMarkdownProcessor):
    description = "Converts PPTX slide decks into Markdown sections, slide by slide."

    def _resolve_effective_options(
        self,
    ) -> tuple[IngestionProcessingProfile, ProcessingConfig.PptxPipelineConfig]:
        # Resolve profile-aware options once per document.
        # Medium/rich can inherit fast defaults through ProcessingConfig.
        processing = get_configuration().processing
        current_profile = get_current_processing_profile()
        active_profile = processing.normalize_profile(current_profile)
        return active_profile, processing.get_effective_pptx_config(active_profile)

    def check_file_validity(self, file_path: Path) -> bool:
        """Checks if the PPTX file is valid and can be opened."""
        try:
            Presentation(str(file_path))
            return True
        except Exception as e:
            logger.error(f"Invalid or corrupted PPTX file: {file_path} - {e}")
            return False

    def extract_file_metadata(self, file_path: Path) -> dict[str, Any]:
        """Extracts basic metadata from the PPTX file."""
        metadata: dict[str, Any] = {"document_name": file_path.name}
        try:
            presentation = Presentation(str(file_path))
            metadata["num_slides"] = len(presentation.slides)
        except Exception as e:
            logger.error(f"Error reading PPTX file: {e}")
            metadata["error"] = str(e)
        return metadata

    def convert_file_to_markdown(self, file_path: Path, output_dir: Path, document_uid: str | None) -> dict:
        """Converts each slide's text content into structured Markdown."""
        output_dir.mkdir(parents=True, exist_ok=True)
        md_path = output_dir / "output.md"

        try:
            active_profile, pptx_options = self._resolve_effective_options()
            presentation = Presentation(str(file_path))
            slide_width = int(getattr(presentation, "slide_width", 0) or 0)
            slide_height = int(getattr(presentation, "slide_height", 0) or 0)

            # Render each slide independently so one problematic slide layout
            # does not break global ordering assumptions.
            slide_texts = [
                _render_slide_markdown(
                    slide_index=slide_index,
                    slide=slide,
                    slide_width=slide_width,
                    slide_height=slide_height,
                    options=pptx_options,
                )
                for slide_index, slide in enumerate(presentation.slides, start=1)
            ]

            content = "\n\n---\n\n".join(slide_texts) if slide_texts else "*No extractable text*"
            md_path.write_text(content, encoding="utf-8")
            logger.info("[PROCESSOR][PPTX] Using profile=%s with effective PPTX options.", active_profile.value)
            return {
                "doc_dir": str(output_dir),
                "md_file": str(md_path),
                "message": "PPTX slides converted to structured Markdown.",
            }

        except Exception as exc:
            logger.exception("Failed to convert PPTX to Markdown: %s", file_path)
            raise InputConversionError(f"PptxMarkdownProcessor failed for '{file_path.name}': {exc}") from exc
