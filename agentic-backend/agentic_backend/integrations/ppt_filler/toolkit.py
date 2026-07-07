# Copyright Thales 2025
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.

"""Inprocess toolkit factory for the ``ppt_filler`` provider (PPTFILL-05 / #1834).

The chat-time fill tool that makes a configured ``ppt_filler`` agent produce a filled,
downloadable deck.

Pipeline (all derived from the agent's persisted :class:`PptFillerParams`, never from
hardcoded slide indices):

1. **Dynamic per-slide ``args_schema``** — built AT TOOL-BUILD TIME from the persisted
   ``schema_slides``. It is nested by slide: a top-level object with one ``slide_<n>``
   property per schema slide (1-based ``SlideSchema.slide``), each of which is an object
   whose leaf properties are that slide's ``{{keys}}``. Each leaf carries its note
   description as the schema-level description so the model gets inline guidance.
2. **Fill** — fetch the template bytes from the agent-config blob store under the fixed
   ``PPT_FILLER_TEMPLATE_KEY``, open with python-pptx, and for each provided ``slide_<n>``
   group fill the n-th (1-based) slide via the SHARED traversal
   (:func:`replace_keys_on_slide`). Reusing the traversal guarantees the filler can never
   diverge from the parser, and that every occurrence of a key on a slide is filled
   consistently.
3. **Download** — render to bytes, upload to SESSION-SCOPED user storage, and return a
   :class:`LinkPart` (``kind=download``) so the UI renders a download button.

Shape coverage is defined by the shared traversal, not re-implemented here: text boxes,
table cells, and shapes nested inside groups are all filled. SmartArt and chart text
remain out of scope (no clean ``text_frame`` API in python-pptx).

NOTE on imports: ``build_ppt_filler_tools`` is imported directly from this submodule by
``core/tools/inprocess_toolkit_registry.py`` (NOT re-exported via the package
``__init__``) to avoid an import cycle through ``agent_spec``. Keep it that way.
"""

from __future__ import annotations

import hashlib
import io
import logging
import time
from typing import TYPE_CHECKING, Any, Dict, Optional, Tuple

from fred_core import convert_pptx_bytes_to_pdf
from langchain_core.tools import BaseTool, StructuredTool
from PIL import Image, UnidentifiedImageError
from pptx import Presentation
from pydantic import BaseModel, Field, create_model

from agentic_backend.common.kf_base_client import KnowledgeFlowAgentContext
from agentic_backend.common.kf_document_client import KfDocumentClient
from agentic_backend.common.kf_workspace_client import KfWorkspaceClient
from agentic_backend.core.agents.v2.contracts.context import ToolInvocationResult
from agentic_backend.core.chatbot.chat_schema import LinkKind, LinkPart, PptPreviewPart
from agentic_backend.integrations.ppt_filler.parser import apply_kept_notes_to_slide
from agentic_backend.integrations.ppt_filler.ppt_filler_params import (
    PPT_FILLER_PROVIDER,
    PptFillerParams,
)
from agentic_backend.integrations.ppt_filler.traversal import (
    list_image_anchors_on_slide,
    replace_keys_on_slide,
)

if TYPE_CHECKING:  # pragma: no cover - typing only
    from agentic_backend.integrations.ppt_filler.parser import KeyField, SlideSchema

logger = logging.getLogger(__name__)

_TOOL_NAME = "fill_ppt_template"
_TOOL_REF = "ppt_filler_fill"
_PPTX_CONTENT_TYPE = (
    "application/vnd.openxmlformats-officedocument.presentationml.presentation"
)
_PDF_CONTENT_TYPE = "application/pdf"
_OUTPUT_FILE_NAME = "filled_presentation.pptx"
_PPTX_SUFFIX = ".pptx"
_PDF_SUFFIX = ".pdf"

# Pillow ``format`` values python-pptx's ``add_picture`` can actually embed. Anything KF
# ingests outside this set (WEBP, ICO, ...) decodes in Pillow but is rejected by
# ``add_picture`` with "unsupported image format". Rather than hard-fail and force the
# agent into a blind re-pick, we transcode such images to PNG in-memory before embedding.
_PPTX_EMBEDDABLE_FORMATS = frozenset({"BMP", "GIF", "JPEG", "PNG", "TIFF", "WMF"})

# Top-level (non-slide) arg the model fills with a human-friendly name for the output
# deck. Optional: when missing/blank we fall back to ``_OUTPUT_FILE_NAME``.
_OUTPUT_NAME_FIELD = "output_file_name"

# The top-level args field for each slide is ``slide_<n>`` where ``n`` is the 1-based
# ``SlideSchema.slide``. Building (here) and parsing (in the tool body) the prefix go
# through the same helper so they can never drift.
_SLIDE_FIELD_PREFIX = "slide_"

# Appended to every TEXT key's leaf description so the agent knows it may emphasize part
# of a value with inline Markdown. The supported subset matches the docx export exactly
# (see ``core/markdown/inline.py``); emphasis overlays bold/italic on the placeholder's
# own font, leaving size/color/font untouched.
_TEXT_FORMATTING_HINT = (
    "You may use inline Markdown to emphasize part of this value: **bold** and "
    "*italic* (or _italic_). Markup is optional — omit it for plain text."
)


def _slide_field_name(slide_number: int) -> str:
    return f"{_SLIDE_FIELD_PREFIX}{slide_number}"


def _sanitize_output_file_name(raw: object) -> str:
    """Turn a model-provided output name into a safe ``*.pptx`` file name.

    The name is used both as a storage key segment and as the download file name, so it
    must not contain path separators or be empty. We keep only the final path component,
    strip control/quote characters, ensure a single ``.pptx`` suffix, and fall back to
    :data:`_OUTPUT_FILE_NAME` when nothing usable remains.
    """
    if not isinstance(raw, str):
        return _OUTPUT_FILE_NAME
    # Keep only the final path component to defeat traversal (e.g. "../../x", "a/b.pptx").
    name = raw.replace("\\", "/").split("/")[-1].strip().strip('"').strip("'").strip()
    # Drop characters that are awkward in file names / storage keys.
    name = "".join(c for c in name if c.isprintable() and c not in '<>:"\\|?*')
    # Normalize the extension: case-insensitive, exactly one trailing ".pptx".
    if name.lower().endswith(_PPTX_SUFFIX):
        name = name[: -len(_PPTX_SUFFIX)].rstrip()
    if not name:
        return _OUTPUT_FILE_NAME
    return f"{name}{_PPTX_SUFFIX}"


def _pdf_key_for(pptx_upload_key: str) -> str:
    """Storage key for the preview PDF: the ``.pptx`` key with a ``.pdf`` suffix.

    Sharing the session prefix means the PDF is cleaned up on session delete exactly like
    the ``.pptx`` (``/storage/user`` keys are ``{session_id}``-scoped).
    """
    base = pptx_upload_key
    if base.lower().endswith(_PPTX_SUFFIX):
        base = base[: -len(_PPTX_SUFFIX)]
    return f"{base}{_PDF_SUFFIX}"


def _pdf_presign_href_from_download(
    pptx_download_url: str, pdf_key: str
) -> Optional[str]:
    """Durable KF href the frontend calls to mint a FRESH presigned PDF URL at open time.

    We deliberately store this durable endpoint instead of a presigned URL: a presigned URL
    expires (~1h) but this part is persisted in the conversation, so a frozen URL 403s when
    the chat is reopened later. The `.pptx` download URL is origin-relative
    (``/…/storage/user/{pptx_key}``); the presign endpoint lives at the sibling path
    ``/…/storage/user/presigned/{pdf_key}``. We derive one from the other so the API prefix
    stays whatever KF actually served, with no host/prefix config duplicated here.
    """
    marker = "/storage/user/"
    idx = pptx_download_url.find(marker)
    if idx == -1:
        return None
    prefix = pptx_download_url[: idx + len(marker)]  # ".../storage/user/"
    return f"{prefix}presigned/{pdf_key.lstrip('/')}"


def _preview_version(pdf_bytes: bytes) -> str:
    """Per-fill freshness token derived from the PDF content.

    A re-fill that changes the deck yields different bytes → a different token → the
    frontend refetches (new ``?v=`` URL + react-pdf remount key). Identical re-fills keep
    the same token, so an unchanged deck is not needlessly refetched. Content-hashing is
    deterministic, which keeps the "version changes on re-fill" behaviour testable offline.
    """
    return hashlib.sha256(pdf_bytes).hexdigest()[:16]


def _slide_number_from_field(field_name: str) -> Optional[int]:
    if not field_name.startswith(_SLIDE_FIELD_PREFIX):
        return None
    suffix = field_name[len(_SLIDE_FIELD_PREFIX) :]
    try:
        return int(suffix)
    except ValueError:
        return None


def _get_ppt_filler_params(agent: KnowledgeFlowAgentContext) -> PptFillerParams:
    """Extract this agent's own :class:`PptFillerParams` from its tuning refs.

    Mirrors ``_get_kf_vector_search_params`` in ``kf_document_client``: scan the agent's
    ``tuning.mcp_servers`` for the ref whose params carry ``provider == ppt_filler`` and
    return them; fall back to a default (empty schema) so callers never need a None check.
    """
    agent_settings = getattr(agent, "agent_settings", None)
    tuning = getattr(agent_settings, "tuning", None)
    if tuning is not None:
        for ref in tuning.mcp_servers:
            if (
                ref.params is not None
                and getattr(ref.params, "provider", None) == PPT_FILLER_PROVIDER
            ):
                return ref.params  # type: ignore[return-value]
    return PptFillerParams()


def _image_leaf_description(key_field: "KeyField") -> str:
    """Build the leaf description for an IMAGE key.

    The agent does not pass image bytes: it passes the *document id* of a document it
    picked by browsing the key's folder. The per-field description stays minimal — it just
    names the field's note and points at the source directory by its ``working_directory``
    PATH (the author's ``folder`` string, e.g. ``images/flags``), which is exactly what the
    ``list_document_tree`` tool takes as its ``working_directory`` argument. The
    ``folder_tag_id`` is deliberately NOT surfaced here: it is an opaque tag id used as a
    search filter, not a browsable path. The generic "how to work with an image field"
    procedure (browse that directory with ``list_document_tree``, pick by name, pass the
    document id) lives once in the main tool description, not repeated on every key.
    """
    note = (key_field.description or "").strip()
    folder = key_field.folder

    parts: list[str] = []
    if note:
        parts.append(note)
    if folder:
        parts.append(f"Image field; its images come from working_directory '{folder}'.")
    else:
        parts.append(
            "Image field, but no source directory is configured — do your best (e.g. "
            "leave it unset) and explain to the user that an owner/editor must fix the "
            "template's image directory."
        )
    return " ".join(parts)


def _build_args_schema(schema_slides: "list[SlideSchema]") -> type[BaseModel]:
    """Build the dynamic, nested-by-slide ``args_schema`` from the persisted schema.

    For each slide we synthesize a per-slide model whose fields are that slide's keys,
    each carrying the key's note description as its schema-level ``description``. The
    top-level model then has one ``slide_<n>`` field per slide pointing at the per-slide
    model. Everything is derived from ``schema_slides`` — no hardcoded indices.

    Every leaf is OPTIONAL (default ``None``): an omitted text key fills as empty string
    (unchanged), an omitted image key removes its empty placeholder slot. Image keys carry
    a description that tells the agent to browse the key's folder and pass a document id.
    """
    # ``create_model`` takes each field as a ``(type, FieldInfo)`` tuple; typing the dict
    # values as ``Any`` keeps the ``**`` spread assignable to its keyword parameters.
    top_fields: Dict[str, Any] = {}
    for slide_schema in schema_slides:
        leaf_fields: Dict[str, Any] = {}
        for key_field in slide_schema.keys:
            # Each leaf is an OPTIONAL string whose description is the note description (for
            # image keys, augmented with browse-the-folder guidance), so the model gets
            # inline per-field guidance. Duplicate keys on one slide are already
            # de-duplicated by the parser, so last-wins here is harmless.
            if key_field.type == "image":
                leaf_description = _image_leaf_description(key_field)
            else:
                # Text key: append the inline-formatting hint so the agent knows it may
                # bold/italic part of the value. The note (if any) leads; the hint follows.
                note = (key_field.description or "").strip()
                leaf_description = (
                    f"{note} {_TEXT_FORMATTING_HINT}" if note else _TEXT_FORMATTING_HINT
                )
            leaf_fields[key_field.key] = (
                Optional[str],
                Field(default=None, description=leaf_description),
            )
        slide_model = create_model(
            f"PptFillSlide{slide_schema.slide}",
            __base__=BaseModel,
            **leaf_fields,
        )
        top_fields[_slide_field_name(slide_schema.slide)] = (
            slide_model,
            Field(
                ...,
                description=(
                    f"Values for the {{key}} fields on slide {slide_schema.slide}."
                ),
            ),
        )
    # Optional, model-chosen name for the produced deck. Lets the agent give the file a
    # relevant name (e.g. derived from the content) instead of a generic default. The
    # ".pptx" extension is added automatically if omitted; leave empty for the default.
    top_fields[_OUTPUT_NAME_FIELD] = (
        Optional[str],
        Field(
            default=None,
            description=(
                "A short, relevant file name for the generated PowerPoint, based on "
                "its content (e.g. 'Proposition_ACME_2026'). The '.pptx' extension is "
                "added automatically. Optional — omit to use a generic default."
            ),
        ),
    )
    return create_model(
        "FillPptTemplateArgs",
        __base__=BaseModel,
        **top_fields,
    )


def _raw_values_for_slide(slide_args: object) -> Dict[str, object]:
    """Coerce one ``slide_<n>`` group (a pydantic model) into a raw ``{key: value}`` map.

    Values are returned as-is (``None`` preserved) so the caller can tell an omitted key
    (``None``) apart from a provided empty string; text-key normalization to ``str``
    happens later in :func:`_text_values`.
    """
    if isinstance(slide_args, BaseModel):
        raw = slide_args.model_dump()
    elif isinstance(slide_args, dict):
        raw = slide_args
    else:
        raw = dict(getattr(slide_args, "__dict__", {}) or {})
    return dict(raw)


def _text_values(raw: Dict[str, object], image_keys: set[str]) -> Dict[str, str]:
    """Project the raw slide values onto the TEXT keys only, as fill strings.

    Image keys are excluded (they are placed as pictures, not text). An omitted text key
    (``None``) becomes the empty string, exactly as before.
    """
    return {
        key: "" if value is None else str(value)
        for key, value in raw.items()
        if key not in image_keys
    }


def _fit_inside_box(
    img_w: int, img_h: int, box_left: int, box_top: int, box_w: int, box_h: int
) -> Tuple[int, int, int, int]:
    """Scale ``(img_w, img_h)`` to fit inside ``(box_w, box_h)`` preserving aspect ratio,
    centered within the box. Inputs/outputs are EMU; the image pixel dims only set the
    aspect ratio. Returns ``(left, top, width, height)`` in EMU.

    Given image aspect ``ar = img_w / img_h`` and box ``(W, H)``:
      - if ``W / H > ar``: the box is relatively wider, so fit the height →
        ``h = H, w = H * ar``, centered horizontally;
      - else: fit the width → ``w = W, h = W / ar``, centered vertically.
    """
    # Degenerate boxes/images: fall back to the full box (no scaling possible).
    if img_w <= 0 or img_h <= 0 or box_w <= 0 or box_h <= 0:
        return (box_left, box_top, box_w, box_h)

    ar = img_w / img_h
    if box_w / box_h > ar:
        height = box_h
        width = box_h * ar
        top = box_top
        left = box_left + (box_w - width) / 2
    else:
        width = box_w
        height = box_w / ar
        left = box_left
        top = box_top + (box_h - height) / 2
    return (round(left), round(top), round(width), round(height))


def _remove_anchor_shape(anchor) -> None:
    """Remove a placeholder shape from its slide (drops the empty ``{{key}}`` box)."""
    element = anchor.shape._element
    element.getparent().remove(element)


def _looks_like_path_not_document_id(value: str) -> bool:
    """True when ``value`` is clearly a file PATH rather than a document id.

    An image key expects the opaque document id returned by ``list_document_tree``
    (a UUID-style identifier that never contains a path separator). The agent sometimes
    mistakenly passes the file's PATH instead (e.g. ``Brand/Logos/logo.png``), which then
    fails the document fetch with an unclear "document not found" error. We detect that
    mistake up front and hard-fail with an actionable message so the agent self-corrects.

    Heuristic (deliberately CONSERVATIVE to keep false positives rare — we only flag values
    that cannot be a valid document id): the value contains a ``/`` or ``\\`` path
    separator. Valid document ids never contain one; any path with directory structure does.
    """
    return "/" in value or "\\" in value


async def _place_images_on_slide(
    *,
    slide,
    image_keys: set[str],
    raw_values: Dict[str, object],
    agent: KnowledgeFlowAgentContext,
) -> Optional[str]:
    """Place chosen images (or remove empty image slots) on one slide.

    For each IMAGE key on the slide we look at every anchor (one per ``{{key}}``
    occurrence) found by :func:`list_image_anchors_on_slide`:

    - The key was OMITTED (value ``None``): remove every placeholder shape so no empty box
      remains. No picture is inserted.
    - The key was PROVIDED a document id: fetch the document's original bytes and, for each
      anchor, insert a fit-inside, aspect-preserved, centered picture, then remove the
      placeholder shape.

    Returns ``None`` on success or an error message string to HARD-FAIL the whole fill:
    - a provided doc id that cannot be fetched or whose bytes are not a usable image
      (the agent's correctable mistake -> re-pick);
    - an image key sitting in a table cell (``invalid_location``) -> a cell cannot hold a
      picture (should not happen for a saved agent, but we refuse rather than crash).
    """
    anchors = list_image_anchors_on_slide(slide)
    # Group anchors per key so each occurrence of the same key is handled together.
    anchors_by_key: Dict[str, list] = {}
    for anchor in anchors:
        if anchor.key in image_keys:
            anchors_by_key.setdefault(anchor.key, []).append(anchor)

    # The document client is built lazily on the FIRST fetch: a slide whose image keys are
    # all omitted (only placeholder removals) must not require an initialized app context.
    document_client: Optional[KfDocumentClient] = None

    for key in image_keys:
        key_anchors = anchors_by_key.get(key, [])
        if not key_anchors:
            continue  # key not present on this slide as an image anchor; nothing to do

        # Guard: an image key in a table cell can never hold a picture. Save rejects this,
        # so a saved agent won't hit it, but refuse loudly instead of crashing.
        if any(anchor.invalid_location for anchor in key_anchors):
            placeholder = f"{{{{{key}}}}}"
            return (
                f"The image field {placeholder} sits in a table cell, which cannot hold "
                "a picture. The template must be fixed by an owner/editor."
            )

        raw_value = raw_values.get(key)
        document_uid = (
            raw_value.strip()
            if isinstance(raw_value, str) and raw_value.strip()
            else None
        )

        if document_uid is None:
            # Omitted image key: remove every empty placeholder slot, insert nothing.
            for anchor in key_anchors:
                _remove_anchor_shape(anchor)
            continue

        # Guard BEFORE fetching: the agent sometimes passes the file's PATH instead of its
        # document id, which the fetch would reject with an unclear "document not found".
        # Detect the wrong-format value up front and hard-fail with an actionable message
        # so the agent re-picks the id (from list_document_tree) in the same turn.
        if _looks_like_path_not_document_id(document_uid):
            placeholder = f"{{{{{key}}}}}"
            return (
                f"The value '{document_uid}' you passed for the image field {placeholder} "
                "looks like a file path, not a document id. Image fields need the document "
                "id returned by the list_document_tree tool (an opaque id with no '/' or "
                "'\\'), not the file's path or name. Browse the field's folder again with "
                "list_document_tree and pass the chosen file's document id."
            )

        # Provided a document id: fetch the original bytes. A fetch failure is a HARD
        # fail (the agent's correctable mistake) -> re-pick.
        try:
            if document_client is None:
                document_client = KfDocumentClient(agent=agent)  # type: ignore[arg-type]
            raw_blob = await document_client.fetch_raw_content(
                document_uid=document_uid
            )
            image_bytes = raw_blob.bytes
        except Exception as exc:
            placeholder = f"{{{{{key}}}}}"
            return (
                f"Could not fetch the document '{document_uid}' you chose for the image "
                f"field {placeholder} [{type(exc).__name__}: {exc}]. Browse the field's "
                "folder again with list_document_tree and pick a valid image document."
            )

        # Validate the bytes ARE a usable image and read its pixel dimensions for the
        # aspect ratio. A non-image (or corrupt image) is a HARD fail -> re-pick.
        # Pillow decodes more formats than python-pptx can embed (e.g. WEBP, ICO), so a
        # format outside ``_PPTX_EMBEDDABLE_FORMATS`` is transcoded to PNG in-memory here
        # rather than rejected later by ``add_picture`` — this avoids a blind re-pick.
        try:
            with Image.open(io.BytesIO(image_bytes)) as image:
                img_w, img_h = image.size
                if image.format not in _PPTX_EMBEDDABLE_FORMATS:
                    buffer = io.BytesIO()
                    # PNG keeps any alpha; RGBA is a safe superset for the formats KF
                    # ingests (WEBP/ICO/GIF can carry transparency).
                    image.convert("RGBA").save(buffer, format="PNG")
                    image_bytes = buffer.getvalue()
        except (UnidentifiedImageError, OSError, ValueError) as exc:
            placeholder = f"{{{{{key}}}}}"
            return (
                f"The document '{document_uid}' you chose for the image field "
                f"{placeholder} is not a usable image [{type(exc).__name__}: {exc}]. "
                "Browse the field's folder again with list_document_tree and pick a "
                "valid image document."
            )

        for anchor in key_anchors:
            left, top, width, height = _fit_inside_box(
                img_w,
                img_h,
                anchor.left,
                anchor.top,
                anchor.width,
                anchor.height,
            )
            # ``add_picture`` may itself reject bytes python-pptx can't decode; treat any
            # failure here as a HARD fail too (re-pick) rather than crash the whole tool.
            try:
                slide.shapes.add_picture(
                    io.BytesIO(image_bytes), left, top, width, height
                )
            except Exception as exc:
                placeholder = f"{{{{{key}}}}}"
                return (
                    f"The document '{document_uid}' you chose for the image field "
                    f"{placeholder} could not be inserted as a picture "
                    f"[{type(exc).__name__}: {exc}]. Pick a valid image document."
                )
            _remove_anchor_shape(anchor)

    return None


def build_ppt_filler_tools(agent: KnowledgeFlowAgentContext) -> list[BaseTool]:
    """Return the in-process LangChain tools for the PPT Filler toolkit.

    The single tool's dynamic ``args_schema`` is derived per-slide from the agent's
    persisted ``schema_slides``. With no schema (template never configured) there is
    nothing to fill, so no tool is exposed.
    """
    params = _get_ppt_filler_params(agent)
    schema_slides = params.schema_slides
    if not schema_slides:
        logger.debug(
            "ppt_filler toolkit requested but the agent has no persisted schema; "
            "returning no tools (template not configured)."
        )
        return []

    args_schema = _build_args_schema(schema_slides)
    template_key = params.template_key

    # Per-slide set of IMAGE keys, derived once from the persisted schema. At fill time we
    # branch on this: text keys go through the text traversal, image keys are placed as
    # pictures (or their empty slots removed). Built from ``type`` so parse and fill agree.
    image_keys_by_slide: Dict[int, set[str]] = {
        slide_schema.slide: {
            key_field.key
            for key_field in slide_schema.keys
            if key_field.type == "image"
        }
        for slide_schema in schema_slides
    }

    async def _fill(**kwargs: object) -> tuple[str, ToolInvocationResult]:
        # Validate/normalize the model-provided args against the dynamic schema so we get
        # the typed per-slide groups regardless of whether LangChain passed raw dicts.
        validated = args_schema(**kwargs)
        provided = validated.model_dump()

        # Resolve the per-run identifiers from the agent context (same accessors the
        # AgentFlow helpers use): agent_id for the agent-config template fetch, and
        # session_id to scope the user-storage output blob.
        runtime_context = getattr(agent, "runtime_context", None)
        agent_settings = getattr(agent, "agent_settings", None)
        agent_id = getattr(agent_settings, "id", None) or type(agent).__name__
        session_id = getattr(runtime_context, "session_id", None)
        access_token = getattr(runtime_context, "access_token", None)

        started = time.monotonic()
        client = KfWorkspaceClient(agent=agent)  # type: ignore[arg-type]

        # 1. Fetch the template from the agent-config blob store (fixed per-agent key).
        try:
            blob = await client.fetch_agent_config_blob(
                template_key,
                access_token,
                agent_id,
            )
        except Exception as exc:
            elapsed = time.monotonic() - started
            message = (
                f"Could not fetch the PPT template '{template_key}' from agent config "
                f"storage after {elapsed:.0f}s [{type(exc).__name__}: {exc}]."
            )
            logger.exception("[PPTFILL][TOOL] template fetch FAILED -> %s", message)
            return message, ToolInvocationResult(tool_ref=_TOOL_REF, is_error=True)

        # 2. Open the template and fill, per-slide. Text keys go through the SHARED text
        #    traversal (text boxes, table cells, grouped shapes). Image keys are placed as
        #    pictures via the image-anchor traversal (or their empty slots removed). A bad
        #    image pick is the agent's correctable mistake -> HARD fail so it re-picks; a
        #    missing folder is handled softly by the agent (it simply passes no doc id).
        try:
            presentation = Presentation(io.BytesIO(blob.bytes))
            slides = list(presentation.slides)
            for field_name, slide_args in provided.items():
                slide_number = _slide_number_from_field(field_name)
                if slide_number is None:
                    continue
                index = slide_number - 1
                if index < 0 or index >= len(slides):
                    # The persisted schema references a slide the template no longer has;
                    # skip rather than crash (template may have been edited).
                    logger.warning(
                        "[PPTFILL][TOOL] schema references slide %d but template has "
                        "only %d slides; skipping.",
                        slide_number,
                        len(slides),
                    )
                    continue
                slide = slides[index]
                raw_values = _raw_values_for_slide(slide_args)
                image_keys = image_keys_by_slide.get(slide_number, set())

                # Text keys: omitted -> empty string (unchanged behavior). IMAGE keys are
                # left untouched here (their ``{{key}}`` text is preserved) so the image
                # anchor traversal below can still find each occurrence; they are then
                # replaced by a picture (or their whole shape removed).
                text_values = _text_values(raw_values, image_keys)

                def value_for(
                    key: str,
                    _values: Dict[str, str] = text_values,
                    _image_keys: set[str] = image_keys,
                ) -> str:
                    # Every occurrence of a key on this slide is filled with the same
                    # value; the shared traversal calls this once per placeholder. An
                    # image key is preserved verbatim ({{key}}) so the image pass can
                    # locate and replace its shape.
                    if key in _image_keys:
                        return f"{{{{{key}}}}}"
                    return _values.get(key, "")

                replace_keys_on_slide(slide, value_for)

                # Image keys: place a picture per anchor (provided doc id) or remove the
                # empty placeholder slot (omitted). A bad pick hard-fails the whole fill.
                if image_keys:
                    error = await _place_images_on_slide(
                        slide=slide,
                        image_keys=image_keys,
                        raw_values=raw_values,
                        agent=agent,
                    )
                    if error is not None:
                        logger.warning(
                            "[PPTFILL][TOOL] image fill rejected -> %s", error
                        )
                        return error, ToolInvocationResult(
                            tool_ref=_TOOL_REF, is_error=True
                        )

            # Strip template-authoring notes from EVERY slide so the ``{{key}}:``
            # descriptions never leak into the deliverable. Any content the author placed
            # after a ``---`` keep-separator is preserved as the slide's real notes.
            for slide in slides:
                apply_kept_notes_to_slide(slide)

            buffer = io.BytesIO()
            presentation.save(buffer)
            filled_bytes = buffer.getvalue()
        except Exception as exc:
            elapsed = time.monotonic() - started
            message = (
                f"Could not fill the PPT template after {elapsed:.0f}s "
                f"[{type(exc).__name__}: {exc}]."
            )
            logger.exception("[PPTFILL][TOOL] fill FAILED -> %s", message)
            return message, ToolInvocationResult(tool_ref=_TOOL_REF, is_error=True)

        # 3. Upload to SESSION-SCOPED user storage and return a download LinkPart. The
        #    session prefix mirrors AgentFlow.upload_user_blob so outputs are cleaned up
        #    on session delete; KfWorkspaceClient.upload_user_blob does not prefix itself.
        #    The model-chosen name (sanitized) becomes both the storage key segment and
        #    the download file name; it falls back to a generic default when not given.
        output_file_name = _sanitize_output_file_name(provided.get(_OUTPUT_NAME_FIELD))
        upload_key = (
            f"{session_id}/{output_file_name}" if session_id else output_file_name
        )
        try:
            upload_result = await client.upload_user_blob(
                key=upload_key,
                file_content=filled_bytes,
                filename=output_file_name,
                content_type=_PPTX_CONTENT_TYPE,
            )
        except Exception as exc:
            elapsed = time.monotonic() - started
            message = (
                f"Filled the template but could not upload the result after "
                f"{elapsed:.0f}s [{type(exc).__name__}: {exc}]."
            )
            logger.exception("[PPTFILL][TOOL] upload FAILED -> %s", message)
            return message, ToolInvocationResult(tool_ref=_TOOL_REF, is_error=True)

        logger.info(
            "[PPTFILL][TOOL] filled deck uploaded session=%s key=%s url=%s",
            session_id,
            upload_result.key,
            upload_result.download_url,
        )

        # 4. Best-effort PDF preview: convert the filled deck to PDF and upload it beside the
        #    .pptx under the same session-scoped key. We store a DURABLE presign href (a KF
        #    endpoint the frontend calls at open time), NOT a presigned URL — a presigned URL
        #    expires (~1h) while this part is persisted, so a frozen URL would 403 when the
        #    chat is reopened later. Any failure (conversion unavailable/timeout, upload
        #    error, or no presign href — e.g. local dev storage) degrades to the plain
        #    download chip so a preview problem never costs the user the deck.
        pdf_bytes = await convert_pptx_bytes_to_pdf(filled_bytes)
        presign_href: Optional[str] = None
        if pdf_bytes is not None:
            pdf_key = _pdf_key_for(upload_key)
            try:
                await client.upload_user_blob(
                    key=pdf_key,
                    file_content=pdf_bytes,
                    filename=_pdf_key_for(output_file_name),
                    content_type=_PDF_CONTENT_TYPE,
                )
                if upload_result.download_url:
                    presign_href = _pdf_presign_href_from_download(
                        upload_result.download_url, pdf_key
                    )
            except Exception as exc:  # best-effort: keep the .pptx on any preview error
                logger.warning(
                    "[PPTFILL][TOOL] preview upload failed (deck still returned): %s",
                    exc,
                )

        if presign_href:
            version = _preview_version(pdf_bytes) if pdf_bytes is not None else ""
            preview = PptPreviewPart(
                preview_id=upload_key,
                title=output_file_name,
                pdf_presign_url=presign_href,
                version=version,
                pptx_download_url=upload_result.download_url,
                file_name=upload_result.file_name,
            )
            artifact = ToolInvocationResult(tool_ref=_TOOL_REF, ui_parts=(preview,))
            return (
                f"The presentation '{output_file_name}' has been filled. A preview opens "
                "beside the chat and a download button is shown automatically; do not write "
                "a download link yourself.",
                artifact,
            )

        # Preview unavailable → fall back to the download chip and say so.
        link = LinkPart(
            href=upload_result.download_url,
            title=f"Download {upload_result.file_name}",
            kind=LinkKind.download,
            mime=_PPTX_CONTENT_TYPE,
            document_uid=upload_result.document_uid,
            file_name=upload_result.file_name,
        )
        artifact = ToolInvocationResult(tool_ref=_TOOL_REF, ui_parts=(link,))
        return (
            f"The presentation '{output_file_name}' has been filled (the preview could not "
            "be generated, so only the download is available). A download button is shown to "
            "the user automatically; do not write a download link yourself.",
            artifact,
        )

    # NOTE: do NOT set response_format="content_and_artifact" here. The v2 ReAct
    # resolver invokes inprocess provider tools via `tool.ainvoke(<plain args dict>)`
    # (see react_tool_resolution._resolve_runtime_provider_tool). With
    # content_and_artifact, LangChain returns ONLY the content string on a plain-args
    # invoke and silently drops the artifact, so the download LinkPart never reaches
    # the UI. Returning a bare ``(content, ToolInvocationResult)`` tuple from the
    # coroutine instead lets the resolver normalize the artifact (it explicitly
    # handles the 2-tuple shape), which is how the ui_parts (the download link) and
    # is_error survive.
    fill_tool = StructuredTool.from_function(
        coroutine=_fill,
        name=_TOOL_NAME,
        description=(
            "Fill the agent's configured PowerPoint template with the provided values. "
            "Input is one 'slide_<n>' object per slide, each holding that slide's "
            "{{key}} fields; each field's own description says what to put there. "
            "Optionally set 'output_file_name' to a short, relevant name for the deck "
            "(the '.pptx' extension is added automatically).\n"
            "IMAGE FIELDS: some fields are images (their description says 'its images "
            "come from working_directory <path>', e.g. 'my directory/other directory'). Before calling "
            "THIS tool you MUST, for EACH such directory, call the 'list_document_tree' "
            "tool with working_directory set to that exact PATH (it is a folder path, "
            "NOT an id) to see the files it contains. Never invent or guess image values "
            "without having listed the directory first. Then, for each image field, pick "
            "the file whose name best fits and pass ONLY that file's document id as the "
            "field's value (not its name, not text). If a directory is missing or empty, "
            "leave that field unset and tell the user an owner/editor must fix the "
            "template's image directory.\n"
            "After the tool runs, the filled PowerPoint is surfaced to the user automatically "
            "(a preview beside the chat when available, otherwise a download button) — do NOT "
            "write, invent, or repeat any download link or URL in your reply, and do not tell "
            "the user to click a link you wrote. Just briefly confirm the deck is ready."
        ),
        args_schema=args_schema,
    )
    return [fill_tool]
