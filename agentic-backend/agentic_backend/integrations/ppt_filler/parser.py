"""Parser / validator for the PPT Filler toolkit (PPTFILL-01).

Turns a ``.pptx`` into a per-slide template schema plus a list of structured,
slide-numbered errors. Pure and offline: it reads the deck, lists ``{{keys}}`` via the
shared traversal, parses descriptions from each slide's notes, and validates that the
text-box keys and the described keys match per slide.

The schema/error Pydantic models defined here (:class:`KeyField`, :class:`SlideSchema`,
:class:`TemplateError`, :class:`ParseResult`) are the data contract that downstream
issues (params model, analyze endpoint, fill tool) import.
"""

from __future__ import annotations

import io
import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Dict, List, Literal, Optional, Tuple, Union

from pptx import Presentation
from pydantic import BaseModel, ConfigDict, Field, model_serializer

from agentic_backend.integrations.ppt_filler.traversal import (
    KEY_PATTERN,
    list_keys_on_slide,
)

# Stable, machine-readable error codes. The frontend maps these to i18n messages; the
# ``message`` field on each error is an English fallback.
CODE_KEY_WITHOUT_DESCRIPTION = "key_without_description"
CODE_DESCRIBED_BUT_NOT_IN_SLIDE = "described_but_not_in_slide"

# Image-metadata error codes (RFC "Image support"). The folder *resolution* codes
# (``folder_not_found``) and the geometry code (``image_key_invalid_location``) are
# defined in their own stories, not here.
CODE_UNKNOWN_METADATA = "unknown_metadata"
CODE_UNKNOWN_TYPE = "unknown_type"
CODE_DUPLICATED_METADATA = "duplicated_metadata"
CODE_IMAGE_WITHOUT_FOLDER = "image_without_folder"
CODE_EMPTY_FOLDER = "empty_folder"
CODE_FOLDER_WITHOUT_IMAGE_TYPE = "folder_without_image_type"

# A metadata line is a leading ``- <key>: <value>`` entry at the very top of a key's
# description block. The value (group 2) may be empty (``- folder:``). Recognized keys
# are ``type`` and ``folder``; anything else is ``unknown_metadata``.
_METADATA_LINE_PATTERN = re.compile(r"^\s*-\s*(\w+)\s*:\s*(.*)$")

# Recognized metadata keys (normalized lowercase). Recognized ``type`` values are
# ``text`` / ``image`` (checked inline so the type narrows to the schema Literal).
_RECOGNIZED_METADATA_KEYS = frozenset({"type", "folder"})

# A notes line is a header ONLY if it is one or more comma-separated ``{{key}}`` tokens
# ending in a colon, e.g. ``{{name}}:`` or ``{{a}}, {{b}}:``. Anything else (including a
# line that merely mentions ``{{...}}`` inline) is description text.
_HEADER_PATTERN = re.compile(r"^\s*\{\{[^}]+\}\}(\s*,\s*\{\{[^}]+\}\})*\s*:\s*$")

# A "keep separator" line is a line of only dashes (>= 3). Everything in a slide's notes
# AFTER the first such line is content the author wants kept verbatim in the FILLED deck
# (e.g. real speaker notes); everything before it is template-authoring content (the
# ``{{key}}:`` headers + descriptions) that is stripped from the output. The same split
# is used by the parser (so keep-notes are never mistaken for headers) and by the filler
# (so the output keeps only the keep-notes).
_KEEP_SEPARATOR_PATTERN = re.compile(r"^\s*-{3,}\s*$")


def split_authoring_and_kept_notes(notes_text: str) -> "tuple[str, str]":
    """Split a slide's notes at the first keep-separator line (``---``).

    Returns ``(authoring_text, kept_text)`` where:

    - ``authoring_text`` is everything BEFORE the first all-dashes line (the
      ``{{key}}:`` headers and their descriptions). It is what the parser reads.
    - ``kept_text`` is everything AFTER that line, with one leading blank line trimmed,
      preserved verbatim as the notes of the filled deck. Empty when there is no
      separator.

    The separator line itself is dropped from both parts.
    """
    lines = notes_text.splitlines()
    for i, line in enumerate(lines):
        if _KEEP_SEPARATOR_PATTERN.match(line):
            authoring = "\n".join(lines[:i])
            kept_lines = lines[i + 1 :]
            # Trim a single leading blank line so "desc\n---\n\nnotes" keeps "notes".
            if kept_lines and kept_lines[0].strip() == "":
                kept_lines = kept_lines[1:]
            return authoring, "\n".join(kept_lines)
    return notes_text, ""


class KeyField(BaseModel):
    """A single template field: a ``{{key}}`` and its note description (scoped to one
    slide).

    ``type`` defaults to ``"text"`` so the JSON contract stays backward compatible; image
    keys carry the author's ``folder`` string and (after resolution in a later story) the
    resolved ``folder_tag_id``. Text keys leave both folder fields ``None``.

    Serialization is **additive**: the three image fields (``type`` / ``folder`` /
    ``folder_tag_id``) are emitted only when they are non-default (i.e. for image keys or
    once a folder/tag is set). A plain text key therefore serializes to the exact
    ``{"key", "description"}`` shape it always has, so the existing wire contract is
    unchanged while image keys gain their metadata.
    """

    key: str
    description: str = ""
    type: Literal["text", "image"] = "text"
    folder: Optional[str] = None  # author's folder string; only meaningful for images
    folder_tag_id: Optional[str] = None  # resolved tag id, filled later; None here

    @model_serializer(mode="wrap")
    def _serialize_additive(self, handler):
        data = handler(self)
        # Drop the image fields when they hold their defaults so text keys keep the
        # legacy ``{key, description}`` wire shape (backward compatible). Image keys (or
        # any key with a resolved folder/tag) keep their populated fields.
        if self.type == "text":
            data.pop("type", None)
        if self.folder is None:
            data.pop("folder", None)
        if self.folder_tag_id is None:
            data.pop("folder_tag_id", None)
        return data


class SlideSchema(BaseModel):
    """The fields of one slide, grouped under its 1-based slide number."""

    slide: int
    keys: List[KeyField] = Field(default_factory=list)


class TemplateError(BaseModel):
    """A per-slide validation error with a stable machine-readable ``code``."""

    slide: int
    key: str
    code: str
    message: str


class ParseResult(BaseModel):
    """Result of :func:`parse`: the per-slide schema and the list of errors.

    The per-slide schema field is named ``slides`` in Python (a bare ``schema``
    attribute would shadow ``BaseModel``), but serializes as ``schema`` so the JSON
    contract stays ``{ "schema": [...], "errors": [...] }`` as specified by the RFC.
    Populate it positionally or by either name.
    """

    model_config = ConfigDict(populate_by_name=True)

    slides: List[SlideSchema] = Field(
        default_factory=list, alias="schema", serialization_alias="schema"
    )
    errors: List[TemplateError] = Field(default_factory=list)


@dataclass
class _ParsedKeyMeta:
    """The parsed result for one described key: its prose description, normalized image
    metadata, and any metadata errors discovered in that header's block.

    ``errors`` holds ``(code, key)`` pairs raised for THIS header; the slide number and
    final ``TemplateError`` message are attached by :func:`parse` (which knows the slide
    and re-attributes a multi-key header's errors to each of its keys).
    """

    description: str = ""
    type: Literal["text", "image"] = "text"  # normalized lowercase
    folder: Optional[str] = None
    errors: List[Tuple[str, str]] = field(default_factory=list)


def _strip_matching_quotes(value: str) -> str:
    """Strip a single matching pair of surrounding single/double quotes, if present."""
    if len(value) >= 2 and value[0] == value[-1] and value[0] in ("'", '"'):
        return value[1:-1]
    return value


def _parse_metadata_block(
    block: List[str],
) -> "Tuple[Literal['text', 'image'], Optional[str], List[Tuple[str, str]], int]":
    """Parse the leading metadata block of a description ``block``.

    The metadata block is the run of leading lines matching ``- <key>: <value>``. It ends
    at the first line that is not such a line (that line and everything after is prose). A
    leading-dash line that is not a ``key: value`` shape ends the block and is prose, not
    metadata, not an error.

    Returns ``(type, folder, errors, prose_start)`` where ``type`` is the normalized type
    (defaults to ``"text"``), ``folder`` is the normalized folder value (``None`` when no
    valid ``folder`` line, else the quote-stripped string, possibly empty), ``errors`` is
    the list of ``(code, "")`` metadata errors (key attached later), and ``prose_start``
    is the index in ``block`` where the prose description begins.
    """
    type_value: Optional[Literal["text", "image"]] = None
    folder_value: Optional[str] = None
    folder_present = False
    errors: List[Tuple[str, str]] = []
    seen_keys: set[str] = set()
    duplicated_keys: set[str] = set()

    i = 0
    n = len(block)
    while i < n:
        match = _METADATA_LINE_PATTERN.match(block[i])
        if match is None:
            break  # first non-metadata line ends the block; the rest is prose
        raw_key, raw_value = match.group(1), match.group(2)
        meta_key = raw_key.strip().lower()
        value = raw_value.strip()

        if meta_key in seen_keys:
            # Same metadata key declared twice in one block. Report once per key.
            if meta_key not in duplicated_keys:
                duplicated_keys.add(meta_key)
                errors.append((CODE_DUPLICATED_METADATA, ""))
            i += 1
            continue
        seen_keys.add(meta_key)

        if meta_key not in _RECOGNIZED_METADATA_KEYS:
            errors.append((CODE_UNKNOWN_METADATA, ""))
        elif meta_key == "type":
            normalized = value.lower()
            if normalized == "text" or normalized == "image":
                type_value = normalized
            else:
                errors.append((CODE_UNKNOWN_TYPE, ""))
        else:  # meta_key == "folder"
            folder_present = True
            folder_value = _strip_matching_quotes(value)
        i += 1

    resolved_type = type_value if type_value is not None else "text"

    # Folder/type cross-validation. Precedence is documented per RFC:
    #   - type: image, no folder line        -> image_without_folder
    #   - type: image, folder present but ""  -> image_without_folder
    #     (the empty value "also maps to image_without_folder" so the message reads
    #      grammatically: "you didn't give a folder")
    #   - type != image, folder present but "" -> empty_folder
    #     (the value problem is the most specific issue; we DO NOT also raise
    #      folder_without_image_type for the same blank line)
    #   - type != image, folder present non-empty -> folder_without_image_type
    if resolved_type == "image":
        if not folder_present or (folder_value is not None and folder_value == ""):
            errors.append((CODE_IMAGE_WITHOUT_FOLDER, ""))
    else:  # text key
        if folder_present:
            if folder_value == "":
                errors.append((CODE_EMPTY_FOLDER, ""))
            else:
                errors.append((CODE_FOLDER_WITHOUT_IMAGE_TYPE, ""))

    return resolved_type, folder_value, errors, i


def _parse_notes_descriptions(notes_text: str) -> Dict[str, _ParsedKeyMeta]:
    """Parse a slide's notes into a ``{key: _ParsedKeyMeta}`` mapping.

    A header line (``{{a}}, {{b}}:``) introduces a (possibly multiline) block. The block
    MAY begin with a contiguous metadata block (``- type: image`` / ``- folder: ...``);
    the prose description is everything after it. Blank lines inside the prose are kept;
    leading/trailing blank lines are trimmed. A multi-key header applies the same
    metadata + description to each listed key (the metadata block is parsed once and
    shared). Lines that merely mention ``{{...}}`` inline are description text, not
    headers.

    Only the authoring portion (before the first ``---`` keep-separator) is parsed; any
    kept-notes content after the separator is opaque and never read as headers.
    """
    descriptions: Dict[str, _ParsedKeyMeta] = {}
    authoring_text, _kept = split_authoring_and_kept_notes(notes_text)
    lines = authoring_text.splitlines()

    i = 0
    n = len(lines)
    while i < n:
        line = lines[i]
        if not _HEADER_PATTERN.match(line):
            # Stray description text with no preceding header: ignore it (only
            # described keys matter).
            i += 1
            continue

        # The keys named on this header line (one or more).
        header_keys = [m.strip() for m in KEY_PATTERN.findall(line)]

        # Collect the block: every line until the next header or EOF.
        block: List[str] = []
        i += 1
        while i < n and not _HEADER_PATTERN.match(lines[i]):
            block.append(lines[i])
            i += 1

        # Parse the leading metadata block (if any); prose starts after it.
        meta_type, meta_folder, meta_errors, prose_start = _parse_metadata_block(block)
        prose = block[prose_start:]

        # Trim leading/trailing blank lines but keep internal blank lines.
        start = 0
        end = len(prose)
        while start < end and prose[start].strip() == "":
            start += 1
        while end > start and prose[end - 1].strip() == "":
            end -= 1
        description = "\n".join(prose[start:end])

        for key in header_keys:
            # Last header for a key wins (consistent with later issues storing one
            # description per key per slide). Errors are attributed to this key.
            descriptions[key] = _ParsedKeyMeta(
                description=description,
                type=meta_type,
                folder=meta_folder,
                errors=[(code, key) for code, _ in meta_errors],
            )

    return descriptions


def _slide_notes_text(slide) -> str:
    """Return the plain text of a slide's notes, or ``""`` if it has none."""
    if not slide.has_notes_slide:
        return ""
    notes_slide = slide.notes_slide
    if notes_slide.notes_text_frame is None:
        return ""
    return notes_slide.notes_text_frame.text or ""


def apply_kept_notes_to_slide(slide) -> None:
    """Rewrite a slide's notes for the FILLED deck, dropping template-authoring content.

    The authoring notes (the ``{{key}}:`` headers and their descriptions) are internal
    guidance and must not leak into the deliverable. After filling, each slide's notes
    are replaced with only the *kept* portion (everything after the first ``---``
    keep-separator); if there is no separator, the notes are cleared entirely.

    Slides with no notes slide are left untouched (so we never create empty notes
    slides). The shared :func:`split_authoring_and_kept_notes` keeps this in lock-step
    with what the parser treats as authoring vs kept content.
    """
    if not slide.has_notes_slide:
        return
    notes_frame = slide.notes_slide.notes_text_frame
    if notes_frame is None:
        return
    _authoring, kept = split_authoring_and_kept_notes(notes_frame.text or "")
    # Setting ``.text`` collapses the frame to a single run with ``kept`` (or empties it).
    notes_frame.text = kept


def _dedupe_preserving_order(keys: List[str]) -> List[str]:
    seen: set[str] = set()
    ordered: List[str] = []
    for key in keys:
        if key not in seen:
            seen.add(key)
            ordered.append(key)
    return ordered


def _metadata_error_message(code: str, key: str, slide_number: int) -> str:
    """English fallback message for a metadata error ``code`` (frontend i18n's by code)."""
    placeholder = f"{{{{{key}}}}}"
    where = f"the notes of slide {slide_number}"
    messages = {
        CODE_UNKNOWN_METADATA: (
            f"{placeholder} in {where} declares an unknown metadata key (only "
            "'type' and 'folder' are recognized)."
        ),
        CODE_UNKNOWN_TYPE: (
            f"{placeholder} in {where} declares an unknown type (only 'text' and "
            "'image' are recognized)."
        ),
        CODE_DUPLICATED_METADATA: (
            f"{placeholder} in {where} declares the same metadata key more than once."
        ),
        CODE_IMAGE_WITHOUT_FOLDER: (
            f"{placeholder} in {where} is an image but does not give a folder."
        ),
        CODE_EMPTY_FOLDER: (
            f"{placeholder} in {where} gives a folder but leaves it blank."
        ),
        CODE_FOLDER_WITHOUT_IMAGE_TYPE: (
            f"{placeholder} in {where} gives a folder but is not an image "
            "(set 'type: image' to use a folder)."
        ),
    }
    return messages.get(code, f"{placeholder} in {where} has a metadata error.")


def parse(pptx_source: Union[bytes, str, Path]) -> ParseResult:
    """Parse a ``.pptx`` into a per-slide schema and a list of validation errors.

    ``pptx_source`` may be raw ``.pptx`` bytes, a path string, or a :class:`Path`.

    Validation is per-slide:

    - ``key_without_description`` â€” a key appears in a text box on slide N but has no
      description in slide N's notes.
    - ``described_but_not_in_slide`` â€” slide N's notes describe a key that never appears
      in a text box on slide N.
    - image-metadata codes (``unknown_metadata``, ``unknown_type``,
      ``duplicated_metadata``, ``image_without_folder``, ``empty_folder``,
      ``folder_without_image_type``) â€” raised while parsing a key's leading metadata
      block, attributed to slide N and the offending key.

    The same key string on different slides is two independent fields (independent
    descriptions); the same key multiple times on one slide is one (deduped) field.
    """
    if isinstance(pptx_source, bytes):
        presentation = Presentation(io.BytesIO(pptx_source))
    else:
        presentation = Presentation(str(pptx_source))

    schema: List[SlideSchema] = []
    errors: List[TemplateError] = []

    for index, slide in enumerate(presentation.slides):
        slide_number = index + 1  # 1-based, as the author sees it

        text_keys = _dedupe_preserving_order(list_keys_on_slide(slide))
        described = _parse_notes_descriptions(_slide_notes_text(slide))

        # Build the schema for this slide from the text-box keys (deduped, in order),
        # carrying parsed image metadata (folder_tag_id stays None â€” resolved later).
        slide_keys: List[KeyField] = []
        for key in text_keys:
            meta = described.get(key)
            if meta is None:
                slide_keys.append(KeyField(key=key))
            else:
                slide_keys.append(
                    KeyField(
                        key=key,
                        description=meta.description,
                        type=meta.type,
                        folder=meta.folder,
                    )
                )
        if slide_keys:
            schema.append(SlideSchema(slide=slide_number, keys=slide_keys))

        text_key_set = set(text_keys)

        # Metadata errors discovered while parsing each described key's leading block.
        # These fire regardless of whether the key appears in a text box on this slide
        # (a described-but-absent image key still reports its bad metadata).
        for key, meta in described.items():
            for code, error_key in meta.errors:
                errors.append(
                    TemplateError(
                        slide=slide_number,
                        key=error_key,
                        code=code,
                        message=_metadata_error_message(code, error_key, slide_number),
                    )
                )

        # key_without_description: in a text box but not described in this slide's notes.
        for key in text_keys:
            if key not in described:
                errors.append(
                    TemplateError(
                        slide=slide_number,
                        key=key,
                        code=CODE_KEY_WITHOUT_DESCRIPTION,
                        message=(
                            f"{{{{{key}}}}} appears on slide {slide_number} but has no "
                            "description in the slide notes."
                        ),
                    )
                )

        # described_but_not_in_slide: described in notes but absent from this slide's
        # text boxes.
        for key in described:
            if key not in text_key_set:
                errors.append(
                    TemplateError(
                        slide=slide_number,
                        key=key,
                        code=CODE_DESCRIBED_BUT_NOT_IN_SLIDE,
                        message=(
                            f"{{{{{key}}}}} is described in the notes of slide "
                            f"{slide_number} but never appears in a text box on that "
                            "slide."
                        ),
                    )
                )

    # Construct via the ``schema`` alias (the field is named ``slides`` in Python to
    # avoid shadowing ``BaseModel``; ``populate_by_name=True`` also allows ``slides=``).
    return ParseResult(schema=schema, errors=errors)
