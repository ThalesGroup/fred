"""Shared text-frame traversal for the PPT Filler toolkit.

This is the single seam reused by BOTH directions of the feature:

- **List** (used by the parser, PPTFILL-01): find every ``{{key}}`` occurrence on a
  slide.
- **Replace** (used by the filler, PPTFILL-05): replace every ``{{key}}`` occurrence on
  a slide with a provided value.

Both directions share the same run-merging logic so they can never diverge: any key the
parser surfaces is guaranteed fillable, and vice versa. The round-trip test
(``parse → fill → re-parse``) is the regression guard for that invariant.

PowerPoint frequently splits a single ``{{key}}`` placeholder across several runs inside
one paragraph (e.g. because of autocorrect or spell-check spans). Both directions
therefore merge the run texts of a paragraph, operate on the merged string, and map the
result back onto the runs.

Shape coverage. A placeholder is fillable wherever python-pptx exposes a text frame, so
the traversal walks, recursively:

- plain text boxes, titles, and other placeholders (``has_text_frame`` shapes);
- **table** cells (each cell is a text frame);
- **grouped** shapes — recursing into the group, so a text box / table nested at any
  depth is reached.

Still out of scope (no clean ``text_frame`` API in python-pptx; text lives in low-level
DrawingML XML): **SmartArt** (``DIAGRAM`` / ``IGX_GRAPHIC``) and **chart** text. Keys
placed there are not seen by the parser and therefore not filled; this is documented in
the RFC rather than silently mis-handled.
"""

from __future__ import annotations

import copy
import re
from dataclasses import dataclass
from typing import TYPE_CHECKING, Callable, List, Tuple

from pptx.oxml.ns import qn

from agentic_backend.core.markdown.inline import Span, parse_inline_markdown

if TYPE_CHECKING:  # pragma: no cover - typing only
    from pptx.shapes.base import BaseShape
    from pptx.slide import Slide
    from pptx.text.text import _Paragraph

# A placeholder is ``{{key}}``. The key is everything between the braces that is not a
# closing brace. This matches the RFC/issue contract exactly.
KEY_PATTERN = re.compile(r"\{\{([^}]+)\}\}")


def _iter_shape_paragraphs(shape: "BaseShape") -> List["_Paragraph"]:
    """Yield every text paragraph reachable from ``shape``, recursing as needed.

    Three text-bearing shape kinds are handled, in priority order:

    - a **group** (``shape_type`` GROUP, exposing ``.shapes``): recurse into each child
      so nested text boxes / tables at any depth are reached;
    - a **table** (``has_table``): every cell is a text frame;
    - a plain **text frame** (``has_text_frame``): its own paragraphs.

    Any other shape (pictures, media, OLE, ink, lines, and — for now — SmartArt and
    charts) contributes no paragraphs.
    """
    paragraphs: List["_Paragraph"] = []

    # Group: recurse. Checked first because a group is itself neither has_text_frame nor
    # has_table, but its children may be either.
    if getattr(shape, "shapes", None) is not None and not shape.has_text_frame:
        for child in shape.shapes:  # type: ignore[attr-defined]
            paragraphs.extend(_iter_shape_paragraphs(child))
        return paragraphs

    # Table: each cell carries its own text frame.
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:  # type: ignore[attr-defined]
            for cell in row.cells:
                paragraphs.extend(cell.text_frame.paragraphs)
        return paragraphs

    # Plain text frame (text box, title, other placeholder, auto-shape, ...).
    if shape.has_text_frame:
        paragraphs.extend(shape.text_frame.paragraphs)  # type: ignore[attr-defined]

    return paragraphs


def _iter_text_paragraphs(slide: "Slide") -> List["_Paragraph"]:
    """Yield every fillable paragraph on ``slide``, descending into tables and groups."""
    paragraphs: List["_Paragraph"] = []
    for shape in slide.shapes:
        paragraphs.extend(_iter_shape_paragraphs(shape))
    return paragraphs


def list_keys_on_slide(slide: "Slide") -> List[str]:
    """Return every ``{{key}}`` key found in the text-frame shapes of ``slide``.

    Keys are returned in document order with duplicates preserved; de-duplication (when
    needed) is the caller's responsibility. Keys are reconstructed even when PowerPoint
    splits a placeholder across multiple runs within a paragraph, because the run texts
    of each paragraph are merged before matching.
    """
    keys: List[str] = []
    for paragraph in _iter_text_paragraphs(slide):
        merged = "".join(run.text for run in paragraph.runs)
        for match in KEY_PATTERN.finditer(merged):
            keys.append(match.group(1).strip())
    return keys


def _styled_spans_for_paragraph(
    merged: str, value_for: Callable[[str], str]
) -> List[Span]:
    """Build the styled spans for a paragraph after substitution.

    Inline Markdown is parsed from the SUBSTITUTED VALUES ONLY, never from the surrounding
    static template text — so an author's literal ``*`` / ``_`` in slide text is never
    reinterpreted as emphasis. The merged paragraph string is walked match-by-match: the
    static text between placeholders becomes one plain span each (verbatim, no parsing),
    and each placeholder's substituted value is run through
    :func:`parse_inline_markdown`, so only there can ``**bold**`` / ``*italic*`` take
    effect.
    """
    spans: List[Span] = []
    pos = 0
    for match in KEY_PATTERN.finditer(merged):
        if match.start() > pos:
            # Static template text: kept verbatim, NOT parsed for markup.
            spans.append(Span(merged[pos : match.start()]))
        value = value_for(match.group(1).strip())
        # The substituted value is the only place inline markup is honored.
        spans.extend(parse_inline_markdown(value))
        pos = match.end()
    if pos < len(merged):
        spans.append(Span(merged[pos:]))
    return spans


def _styled_run_element(base_r, text: str, *, bold: bool, italic: bool):
    """Clone ``base_r`` into a new ``<a:r>`` carrying ``text`` and overlaid emphasis.

    The clone inherits every property the author set on the base ``{{key}}`` run (size,
    color, font name, …) — python-pptx has no public font-copy, so cloning the element is
    how the whole base style is carried over without enumerating individual properties.
    Only ``bold`` / ``italic`` are overlaid on top.
    """
    new_r = copy.deepcopy(base_r)
    text_el = new_r.find(qn("a:t"))
    if text_el is None:
        text_el = new_r.makeelement(qn("a:t"), {})
        new_r.append(text_el)
    text_el.text = text
    if bold or italic:
        r_pr = new_r.find(qn("a:rPr"))
        if r_pr is None:
            r_pr = new_r.makeelement(qn("a:rPr"), {})
            new_r.insert(0, r_pr)
        if bold:
            r_pr.set("b", "1")
        if italic:
            r_pr.set("i", "1")
    return new_r


def _line_break_element(base_r):
    """Build an ``<a:br/>`` line break carrying the base run's properties.

    A ``\\n`` inside a run's ``<a:t>`` is NOT a valid DrawingML line break: PowerPoint /
    LibreOffice render it unreliably and bleed the *preceding* run's emphasis across the
    wrap (a bold title leaks bold onto the body line that follows it). Newlines in a value
    must therefore become explicit ``<a:br/>`` elements between runs. The break carries a
    clone of the base ``rPr`` so the line spacing/font of the break matches the text.
    """
    br = base_r.makeelement(qn("a:br"), {})
    base_r_pr = base_r.find(qn("a:rPr"))
    if base_r_pr is not None:
        br.append(copy.deepcopy(base_r_pr))
    return br


def _span_elements(base_r, span: Span):
    """Yield the ``<a:r>`` / ``<a:br/>`` elements for one span, splitting on newlines.

    Each newline in the span text becomes an ``<a:br/>`` between two runs (see
    :func:`_line_break_element`); a code span carries no pptx emphasis (RFC: treated as
    plain text), so only ``bold`` / ``italic`` are overlaid.
    """
    pieces = span.text.split("\n")
    elements = []
    for index, piece in enumerate(pieces):
        if index > 0:
            elements.append(_line_break_element(base_r))
        elements.append(
            _styled_run_element(base_r, piece, bold=span.bold, italic=span.italic)
        )
    return elements


def _write_spans_onto_paragraph(paragraph: "_Paragraph", spans: List[Span]) -> None:
    """Rewrite a paragraph's runs from ``spans``, inheriting the first run's style.

    The first existing run is the BASE STYLE SOURCE: every new run clones its ``<a:r>``
    XML and only overlays ``bold`` / ``italic`` on top (see :func:`_styled_run_element`).
    Newlines inside a value become explicit ``<a:br/>`` breaks, NOT literal ``\\n`` in run
    text, so emphasis never bleeds across a wrapped line (see :func:`_line_break_element`).

    A markup-free, single-line value yields a single plain run with exactly the base run's
    style — today's behavior, unchanged.
    """
    runs = list(paragraph.runs)
    if not runs:
        return
    base_r = runs[0]._r

    new_elements: list = []
    for span in spans:
        new_elements.extend(_span_elements(base_r, span))

    # Swap the new elements in for the old runs in place: insert them in order right after
    # the base run (each after the previous so document order is preserved), then drop
    # every original run. Doing it in this order keeps ``base_r`` valid while we clone, and
    # leaves the paragraph holding exactly the new runs/breaks, in order.
    anchor = base_r
    for element in new_elements:
        anchor.addnext(element)
        anchor = element
    for run in runs:
        run._r.getparent().remove(run._r)


def replace_keys_on_slide(slide: "Slide", value_for: Callable[[str], str]) -> None:
    """Replace every ``{{key}}`` occurrence in the text-frame shapes of ``slide``.

    ``value_for`` maps a (stripped) key to its replacement string. It is called once per
    placeholder occurrence; every occurrence of a key on the slide is filled
    consistently as long as ``value_for`` is deterministic.

    The same run-merging logic as :func:`list_keys_on_slide` is used, so a key split
    across runs is correctly replaced. After substitution, the paragraph's runs are
    rewritten from the styled spans: inline Markdown (``**bold**`` / ``*italic*``) is
    parsed from the SUBSTITUTED VALUES ONLY and overlaid on top of the first run's
    inherited formatting; static template text (including any literal ``*`` / ``_``) is
    kept verbatim. A value with no markup yields a single run with the first run's style,
    exactly as before.
    """
    for paragraph in _iter_text_paragraphs(slide):
        runs = list(paragraph.runs)
        if not runs:
            continue
        merged = "".join(run.text for run in runs)
        if not KEY_PATTERN.search(merged):
            continue

        replaced = KEY_PATTERN.sub(
            lambda match: value_for(match.group(1).strip()), merged
        )
        if replaced == merged:
            continue

        spans = _styled_spans_for_paragraph(merged, value_for)

        # Fast path — no value carried emphasis: collapse onto the first run exactly as
        # before (one run, first run's formatting). This keeps a markup-free fill byte-for-
        # byte identical to the prior behavior; only an actually-emphasized value takes the
        # multi-run path below.
        if not any(span.bold or span.italic for span in spans):
            runs[0].text = replaced
            for extra_run in runs[1:]:
                extra_run.text = ""
            continue

        _write_spans_onto_paragraph(paragraph, spans)


# ---------------------------------------------------------------------------
# Image-anchor shape-walking traversal (geometry seam).
#
# ``list_keys_on_slide`` / ``replace_keys_on_slide`` flatten shapes to paragraphs and
# discard the containing shape, so they cannot give an image its placement box. This
# second traversal walks the *shapes* (recursing into groups) and yields, per ``{{key}}``
# found in a shape's text, the containing shape plus its ABSOLUTE EMU geometry. It drives
# BOTH the analyze-time ``image_key_invalid_location`` check and the fill-time picture
# insertion, so parse and fill can never diverge on geometry.
#
# It deliberately reuses ``KEY_PATTERN`` and the run-merging idiom (a key may be split
# across runs), but stays a separate function from the text seam: same merging, different
# output (shape + geometry vs. text round-trip).
# ---------------------------------------------------------------------------


@dataclass(frozen=True)
class ImageAnchor:
    """An ``{{key}}`` occurrence resolved to a placement box on a slide.

    All coordinates are ABSOLUTE EMU on the slide (group transforms already composed).
    ``shape`` is the python-pptx shape that holds the key, kept for fill-time removal /
    picture insertion. ``invalid_location`` is ``True`` when the key sits in a table cell
    (a cell cannot hold a picture); the geometry is then meaningless and only the flag is
    consumed by the caller.
    """

    key: str
    left: int
    top: int
    width: int
    height: int
    shape: object
    invalid_location: bool


# A transform mapping child (group-local) coordinates to absolute slide coordinates:
# ``abs = offset + (child - child_offset) * scale`` per axis. The identity transform
# (top level, no enclosing group) is ``offset=child_offset=0`` and ``scale=1``.
_Transform = Tuple[float, float, float, float, float, float]
# (off_x, off_y, ch_off_x, ch_off_y, scale_x, scale_y)
_IDENTITY: "_Transform" = (0.0, 0.0, 0.0, 0.0, 1.0, 1.0)


def _merged_keys_of_shape(shape: "BaseShape") -> List[str]:
    """Return the (stripped) ``{{key}}`` keys in a single text-frame shape's own text.

    Run-merging is per paragraph, identical to the text seam, so a placeholder split
    across runs is still recognised. The shape must expose a text frame; callers gate on
    that. Keys are returned in document order with duplicates preserved.
    """
    keys: List[str] = []
    for paragraph in shape.text_frame.paragraphs:  # type: ignore[attr-defined]
        merged = "".join(run.text for run in paragraph.runs)
        for match in KEY_PATTERN.finditer(merged):
            keys.append(match.group(1).strip())
    return keys


def _shape_has_key(shape: "BaseShape") -> bool:
    return bool(_merged_keys_of_shape(shape))


def _group_child_transform(group: "BaseShape") -> "_Transform":
    """Build the child→absolute transform contributed by ``group``'s own ``xfrm``.

    PowerPoint stores a group's placement as ``off``/``ext`` (the box the group occupies
    in its parent's coordinate space) and ``chOff``/``chExt`` (the coordinate system its
    children are authored in). The mapping from a child coordinate to the group's own
    coordinate space is, per axis::

        coord_in_group = off + (child_coord - chOff) * (ext / chExt)

    ``off``/``ext`` are read via the public ``.left/.top/.width/.height`` (which resolve
    the same ``a:off``/``a:ext``), and ``chOff``/``chExt`` from the group's XML
    (``grpSpPr/xfrm`` → ``a:chOff``/``a:chExt``). A zero child-extent would divide by
    zero, so the scale falls back to ``1.0`` on that axis (and the group is then treated
    as a plain translation by ``off``).
    """
    off_x = group.left or 0
    off_y = group.top or 0
    ext_cx = group.width or 0
    ext_cy = group.height or 0

    ch_off_x = ch_off_y = 0
    ch_ext_cx = ch_ext_cy = 0

    grp_sp_pr = group.element.find(qn("p:grpSpPr"))
    xfrm = grp_sp_pr.find(qn("a:xfrm")) if grp_sp_pr is not None else None
    if xfrm is not None:
        ch_off = xfrm.find(qn("a:chOff"))
        ch_ext = xfrm.find(qn("a:chExt"))
        if ch_off is not None:
            ch_off_x = int(ch_off.get("x", "0"))
            ch_off_y = int(ch_off.get("y", "0"))
        if ch_ext is not None:
            ch_ext_cx = int(ch_ext.get("cx", "0"))
            ch_ext_cy = int(ch_ext.get("cy", "0"))

    # Guard div-by-zero: a zero child-extent makes the scale undefined, so treat it as 1.
    scale_x = (ext_cx / ch_ext_cx) if ch_ext_cx else 1.0
    scale_y = (ext_cy / ch_ext_cy) if ch_ext_cy else 1.0

    return (off_x, off_y, ch_off_x, ch_off_y, scale_x, scale_y)


def _apply_transform(
    transform: "_Transform", left: int, top: int, width: int, height: int
) -> Tuple[int, int, int, int]:
    """Map a child's group-local box through ``transform`` to absolute EMU.

    Position is offset-then-scaled; extent is scaled only. Results are rounded to whole
    EMU (PowerPoint geometry is integral)::

        abs_left = off_x + (left - ch_off_x) * scale_x
        abs_w    = width * scale_x
    """
    off_x, off_y, ch_off_x, ch_off_y, scale_x, scale_y = transform
    abs_left = off_x + (left - ch_off_x) * scale_x
    abs_top = off_y + (top - ch_off_y) * scale_y
    abs_w = width * scale_x
    abs_h = height * scale_y
    return (round(abs_left), round(abs_top), round(abs_w), round(abs_h))


def _compose(outer: "_Transform", inner: "_Transform") -> "_Transform":
    """Compose two transforms so ``apply(outer∘inner, p) == apply(outer, apply(inner, p))``.

    Nested groups multiply: a child of an inner group is first mapped into the inner
    group's coordinate space, then that space is mapped into the outer group's, etc. The
    composed transform stays in the same ``(off, ch_off, scale)`` form so the recursion
    can keep carrying a single accumulated transform down each level.
    """
    o_off_x, o_off_y, o_ch_off_x, o_ch_off_y, o_sx, o_sy = outer
    i_off_x, i_off_y, i_ch_off_x, i_ch_off_y, i_sx, i_sy = inner
    # apply(outer, apply(inner, p)) expanded per axis:
    #   inner: a = i_off + (p - i_ch_off) * i_s
    #   outer: b = o_off + (a - o_ch_off) * o_s
    # => b = [o_off + (i_off - o_ch_off) * o_s] + (p - i_ch_off) * (i_s * o_s)
    new_off_x = o_off_x + (i_off_x - o_ch_off_x) * o_sx
    new_off_y = o_off_y + (i_off_y - o_ch_off_y) * o_sy
    return (
        new_off_x,
        new_off_y,
        float(i_ch_off_x),
        float(i_ch_off_y),
        i_sx * o_sx,
        i_sy * o_sy,
    )


def _is_group(shape: "BaseShape") -> bool:
    """A group exposes ``.shapes`` but is itself not a text frame nor a table.

    Matches the duck-typing already used by the text seam (a group is neither
    ``has_text_frame`` nor ``has_table`` but carries children).
    """
    return (
        getattr(shape, "shapes", None) is not None
        and not shape.has_text_frame
        and not getattr(shape, "has_table", False)
    )


def _collect_image_anchors(
    shape: "BaseShape", transform: "_Transform", out: List[ImageAnchor]
) -> None:
    """Append the image anchors contributed by ``shape`` (recursing into groups)."""
    # Group: descend, composing this group's transform onto the accumulated one.
    if _is_group(shape):
        child_transform = _compose(transform, _group_child_transform(shape))
        for child in shape.shapes:  # type: ignore[attr-defined]
            _collect_image_anchors(child, child_transform, out)
        return

    # Table: any key in any cell is an invalid image location. Geometry is irrelevant
    # (a cell cannot hold a picture), so the table shape's own box is reported only as a
    # placeholder and the flag is what the caller consumes.
    if getattr(shape, "has_table", False):
        abs_box = _apply_transform(
            transform,
            shape.left or 0,
            shape.top or 0,
            shape.width or 0,
            shape.height or 0,
        )
        for row in shape.table.rows:  # type: ignore[attr-defined]
            for cell in row.cells:
                for paragraph in cell.text_frame.paragraphs:
                    merged = "".join(run.text for run in paragraph.runs)
                    for match in KEY_PATTERN.finditer(merged):
                        out.append(
                            ImageAnchor(
                                key=match.group(1).strip(),
                                left=abs_box[0],
                                top=abs_box[1],
                                width=abs_box[2],
                                height=abs_box[3],
                                shape=shape,
                                invalid_location=True,
                            )
                        )
        return

    # Plain text frame (text box, autoshape, title, ...): one anchor per key occurrence.
    if shape.has_text_frame:
        keys = _merged_keys_of_shape(shape)
        if not keys:
            return
        abs_box = _apply_transform(
            transform,
            shape.left or 0,
            shape.top or 0,
            shape.width or 0,
            shape.height or 0,
        )
        for key in keys:
            out.append(
                ImageAnchor(
                    key=key,
                    left=abs_box[0],
                    top=abs_box[1],
                    width=abs_box[2],
                    height=abs_box[3],
                    shape=shape,
                    invalid_location=False,
                )
            )
        return

    # Any other shape (picture, media, chart, SmartArt, line, ...) contributes nothing.


def list_image_anchors_on_slide(slide: "Slide") -> List[ImageAnchor]:
    """Return an :class:`ImageAnchor` per ``{{key}}`` occurrence on ``slide``.

    Walks every shape (descending into groups at any depth) and reports, for each key, the
    containing shape and its absolute EMU geometry. A key that appears in several shapes
    yields one anchor per shape, so each occurrence is filled later. Keys in table cells
    are flagged ``invalid_location=True``. SmartArt and chart text are out of scope (same
    as the text seam) and contribute no anchors.
    """
    anchors: List[ImageAnchor] = []
    for shape in slide.shapes:
        _collect_image_anchors(shape, _IDENTITY, anchors)
    return anchors
