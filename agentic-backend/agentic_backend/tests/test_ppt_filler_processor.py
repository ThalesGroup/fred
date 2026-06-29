# Copyright Thales 2025
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.

"""Unit tests for the ``ppt_filler`` ToolkitAssetProcessor (PPTFILL-04 / #1833).

Covers all three incoming-params states from the RFC "Save flow" table plus the
strip-before-persist invariant. Storage is FAKED/injected (no HTTP). Decks are built
in-test with python-pptx, mirroring the PPTFILL-01 parser tests (no checked-in binaries).
"""

import base64
import io
from typing import Dict, List, Optional, Tuple

import pytest
from pptx import Presentation
from pptx.util import Inches

from agentic_backend.core.tools.toolkit_asset_processor import (
    ToolkitAssetValidationError,
)
from agentic_backend.integrations.ppt_filler.folder_resolution import (
    CODE_FOLDER_NOT_FOUND,
)
from agentic_backend.integrations.ppt_filler.parser import KeyField, SlideSchema
from agentic_backend.integrations.ppt_filler.ppt_filler_params import (
    PPT_FILLER_TEMPLATE_KEY,
    PptFillerParams,
)
from agentic_backend.integrations.ppt_filler.ppt_filler_processor import (
    CODE_ASSET_REQUIRED,
    CODE_INVALID_UPLOAD,
    PptFillerAssetProcessor,
)

# (text-box body, notes text) per slide — same fixture shape as the PPTFILL-01 tests.
SlideSpec = Tuple[str, str]
_AGENT_ID = "agent-123"


def _build_deck(slides: List[SlideSpec]) -> bytes:
    """Build a ``.pptx`` with one text box and one notes block per slide spec."""
    presentation = Presentation()
    blank_layout = presentation.slide_layouts[6]  # blank
    for body, notes in slides:
        slide = presentation.slides.add_slide(blank_layout)
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
        textbox.text_frame.text = body
        if notes:
            notes_frame = slide.notes_slide.notes_text_frame
            assert notes_frame is not None  # python-pptx creates it on access
            notes_frame.text = notes
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


class _FakeStore:
    """In-memory fake of the ToolkitAssetStore port; records every upload call."""

    def __init__(self) -> None:
        self.calls: list[dict] = []

    async def upload_agent_config_blob(
        self,
        key: str,
        file_content,
        filename: str,
        agent_id: str,
        content_type: Optional[str] = None,
    ):
        self.calls.append(
            {
                "key": key,
                "file_content": file_content,
                "filename": filename,
                "agent_id": agent_id,
                "content_type": content_type,
            }
        )
        return {"key": key, "file_name": filename, "size": len(file_content)}


class _FakeResolver:
    """In-memory ``folder full-path -> tag id`` resolver; records each resolve call."""

    def __init__(self, known: Dict[str, str]):
        self.known = known
        self.calls: List[str] = []

    async def resolve(self, folder: str) -> Optional[str]:
        self.calls.append(folder)
        return self.known.get(folder)


def _b64_deck(slides: List[SlideSpec]) -> str:
    return base64.b64encode(_build_deck(slides)).decode("ascii")


# An image-key notes block: type:image + a folder, plus prose so it is well-formed.
_IMAGE_NOTES = "{{flag}}:\n- type: image\n- folder: images/flags\nPick a flag."


@pytest.mark.asyncio
async def test_upload_bytes_valid_uploads_blob_recomputes_schema_and_strips_bytes():
    # State 1: a VALID template (every key described). The processor must upload the blob
    # under the FIXED key, recompute schema server-side from the bytes, and strip the
    # transient upload field.
    deck_b64 = _b64_deck(
        [
            ("Hello {{name}} and {{role}}", "{{name}}:\nThe name\n{{role}}:\nThe role"),
            ("City: {{city}}", "{{city}}:\nThe city"),
        ]
    )
    # Seed an unrelated/stale schema to prove the SERVER recomputes from the bytes.
    params = PptFillerParams(
        schema=[SlideSchema(slide=99, keys=[KeyField(key="stale", description="x")])],
        template_upload_b64=deck_b64,
    )
    store = _FakeStore()

    result = await PptFillerAssetProcessor().process(
        params, agent_id=_AGENT_ID, store=store
    )

    # Blob uploaded once, under the fixed per-agent key, scoped to the agent.
    assert len(store.calls) == 1
    call = store.calls[0]
    assert call["key"] == PPT_FILLER_TEMPLATE_KEY
    assert call["agent_id"] == _AGENT_ID
    assert isinstance(call["file_content"], (bytes, bytearray))

    # Schema recomputed from the ACTUAL bytes (not the stale seed).
    assert [s.slide for s in result.schema_slides] == [1, 2]
    assert result.schema_slides[0].keys[0].key == "name"
    assert result.schema_slides[1].keys[0].key == "city"

    # Strip-before-persist invariant: the returned params never carry the upload bytes,
    # neither on the attribute nor in any serialized form.
    assert result.template_upload_b64 is None
    dumped = result.model_dump(by_alias=True)
    assert dumped.get("template_upload_b64") is None
    assert "template_upload_b64" not in {k for k, v in dumped.items() if v is not None}

    # The input model was not mutated in place.
    assert params.template_upload_b64 == deck_b64


@pytest.mark.asyncio
async def test_upload_bytes_invalid_template_raises_typed_error_with_structured_errors():
    # State 1 (invalid): a key without a description must reject with the typed error
    # carrying the structured per-slide errors (slide, key, code, message).
    deck_b64 = _b64_deck([("Hi {{name}} aged {{age}}", "{{name}}:\nThe name")])
    params = PptFillerParams(template_upload_b64=deck_b64)
    store = _FakeStore()

    with pytest.raises(ToolkitAssetValidationError) as excinfo:
        await PptFillerAssetProcessor().process(params, agent_id=_AGENT_ID, store=store)

    errors = excinfo.value.errors
    assert errors, "expected at least one structured error"
    for err in errors:
        assert {"slide", "key", "code", "message"} <= set(err.model_dump().keys())
    codes = {(e.code, e.key, e.slide) for e in errors}
    assert ("key_without_description", "age", 1) in codes

    # The 422 JSON body shape matches the analyze endpoint contract.
    payload = excinfo.value.errors_payload()
    assert payload and all(
        set(d.keys()) == {"slide", "key", "code", "message"} for d in payload
    )

    # Invalid template → NOTHING uploaded (no half-written blob).
    assert store.calls == []


@pytest.mark.asyncio
async def test_bytes_absent_schema_present_is_noop_passthrough():
    # State 2: ordinary edit — no upload bytes, schema already present. The processor must
    # leave the template/storage untouched and pass the params through unchanged.
    params = PptFillerParams(
        schema=[SlideSchema(slide=1, keys=[KeyField(key="name", description="n")])],
    )
    store = _FakeStore()

    result = await PptFillerAssetProcessor().process(
        params, agent_id=_AGENT_ID, store=store
    )

    # Storage NOT called: the template is untouched.
    assert store.calls == []
    # Params unchanged (same schema, still no upload field).
    assert result is params
    assert result.template_upload_b64 is None
    assert [s.slide for s in result.schema_slides] == [1]


@pytest.mark.asyncio
async def test_bytes_absent_schema_absent_required_rejects():
    # State 3: no upload bytes, no schema, asset required → reject create/update.
    params = PptFillerParams()  # empty schema, no upload
    store = _FakeStore()

    with pytest.raises(ToolkitAssetValidationError) as excinfo:
        await PptFillerAssetProcessor().process(params, agent_id=_AGENT_ID, store=store)

    codes = {e.code for e in excinfo.value.errors}
    assert CODE_ASSET_REQUIRED in codes
    assert store.calls == []


@pytest.mark.asyncio
async def test_unreadable_upload_rejects_without_uploading():
    # A non-.pptx (or corrupt) upload must be rejected as invalid, never uploaded.
    bad_b64 = base64.b64encode(b"this is not a pptx").decode("ascii")
    params = PptFillerParams(template_upload_b64=bad_b64)
    store = _FakeStore()

    with pytest.raises(ToolkitAssetValidationError) as excinfo:
        await PptFillerAssetProcessor().process(params, agent_id=_AGENT_ID, store=store)

    assert CODE_INVALID_UPLOAD in {e.code for e in excinfo.value.errors}
    assert store.calls == []


# --- space-aware folder resolution at save time (Story 05) -----------------------------


@pytest.mark.asyncio
async def test_image_template_with_known_folder_persists_tag_id_and_strips_bytes():
    """A valid image template whose folder resolves -> blob uploaded, schema persisted WITH
    the resolved folder_tag_id, transient upload bytes stripped."""
    deck_b64 = _b64_deck([("{{flag}}", _IMAGE_NOTES)])
    params = PptFillerParams(template_upload_b64=deck_b64)
    store = _FakeStore()
    resolver = _FakeResolver({"images/flags": "tag-123"})

    result = await PptFillerAssetProcessor().process(
        params,
        agent_id=_AGENT_ID,
        store=store,
        team_id="team-42",
        folder_resolver=resolver,
    )

    # Folder resolved exactly once, in the agent's space.
    assert resolver.calls == ["images/flags"]

    # Blob uploaded once under the fixed key.
    assert len(store.calls) == 1
    assert store.calls[0]["key"] == PPT_FILLER_TEMPLATE_KEY

    # Schema persisted with the resolved tag id; image fields preserved.
    flag = result.schema_slides[0].keys[0]
    assert flag.key == "flag"
    assert flag.type == "image"
    assert flag.folder == "images/flags"
    assert flag.folder_tag_id == "tag-123"

    # Strip-before-persist invariant still holds.
    assert result.template_upload_b64 is None


@pytest.mark.asyncio
async def test_image_template_unknown_folder_rejects_and_uploads_nothing():
    """An image key whose folder does not resolve -> ToolkitAssetValidationError carrying
    folder_not_found, and NOTHING uploaded (hard 422 at save time)."""
    deck_b64 = _b64_deck([("{{flag}}", _IMAGE_NOTES)])
    params = PptFillerParams(template_upload_b64=deck_b64)
    store = _FakeStore()
    resolver = _FakeResolver({"images/other": "tag-999"})  # different folder

    with pytest.raises(ToolkitAssetValidationError) as excinfo:
        await PptFillerAssetProcessor().process(
            params,
            agent_id=_AGENT_ID,
            store=store,
            team_id="team-42",
            folder_resolver=resolver,
        )

    codes = {(e.code, e.key, e.slide) for e in excinfo.value.errors}
    assert (CODE_FOLDER_NOT_FOUND, "flag", 1) in codes
    # Invalid → nothing written to storage.
    assert store.calls == []


@pytest.mark.asyncio
async def test_personal_scope_resolver_is_honored():
    """Without a team_id the processor still resolves via the (personal-scoped) resolver
    handed in; the resolver receives the folder string and the upload proceeds."""
    deck_b64 = _b64_deck([("{{flag}}", _IMAGE_NOTES)])
    params = PptFillerParams(template_upload_b64=deck_b64)
    store = _FakeStore()
    resolver = _FakeResolver({"images/flags": "tag-personal"})

    result = await PptFillerAssetProcessor().process(
        params,
        agent_id=_AGENT_ID,
        store=store,
        team_id=None,  # personal space
        folder_resolver=resolver,
    )

    assert resolver.calls == ["images/flags"]
    assert result.schema_slides[0].keys[0].folder_tag_id == "tag-personal"
    assert len(store.calls) == 1


@pytest.mark.asyncio
async def test_image_template_without_resolver_does_not_crash():
    """Backward-compat: with NO resolver threaded in (e.g. no folders configured) the
    processor must not crash on an image template; it skips resolution and persists the
    schema as-is (folder_tag_id stays unset)."""
    deck_b64 = _b64_deck([("{{flag}}", _IMAGE_NOTES)])
    params = PptFillerParams(template_upload_b64=deck_b64)
    store = _FakeStore()

    result = await PptFillerAssetProcessor().process(
        params, agent_id=_AGENT_ID, store=store
    )

    # No resolver → no resolution; the blob is still uploaded and the schema persisted.
    assert len(store.calls) == 1
    flag = result.schema_slides[0].keys[0]
    assert flag.type == "image"
    assert flag.folder == "images/flags"
    assert flag.folder_tag_id is None
    assert result.template_upload_b64 is None
