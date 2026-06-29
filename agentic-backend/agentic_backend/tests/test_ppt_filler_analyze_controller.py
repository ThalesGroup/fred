# Copyright Thales 2025
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.

"""Thin HTTP smoke test for the stateless PPT Filler analyze endpoint.

Core parsing is already covered by ``tests/test_ppt_filler_parser.py`` (PPTFILL-01); this
only asserts the ``200 { schema, errors }`` response shape — for a valid upload and for a
template that produces parse errors carrying ``slide``, ``key``, ``code``, ``message``.

Decks are built in-test with python-pptx (no checked-in binaries), mirroring the
PPTFILL-01 parser tests.
"""

import io
from typing import Dict, List, Optional, Tuple

from fastapi import status
from fastapi.testclient import TestClient
from fred_core.common import OwnerFilter
from pptx import Presentation
from pptx.util import Inches

from agentic_backend.core.agents import agent_controller

_ANALYZE_URL = "/agentic/v1/agents/ppt-filler/analyze"
# Test bearer token (not a real credential).
_TOKEN = "dummy-token"  # nosec B105
_HEADERS = {"Authorization": f"Bearer {_TOKEN}"}

# (text-box body, notes text) per slide — same fixture shape as the PPTFILL-01 tests.
SlideSpec = Tuple[str, str]


class _FakeResolver:
    """In-memory ``folder full-path -> tag id`` resolver; records each resolve call."""

    def __init__(self, known: Dict[str, str]):
        self.known = known
        self.calls: List[str] = []

    async def resolve(self, folder: str) -> Optional[str]:
        self.calls.append(folder)
        return self.known.get(folder)


class _RecordingFolderResolverFactory:
    """Stand-in for ``_build_folder_resolver``: records its scope args, returns a fake.

    Patched over the module-level factory so the analyze endpoint resolves folders against
    an in-memory fake (no live Knowledge Flow) and the test can assert the scope
    (``team_id`` -> TEAM, else PERSONAL) that the endpoint forwarded.
    """

    def __init__(self, resolver: _FakeResolver):
        self._resolver = resolver
        self.captured: List[dict] = []

    def __call__(self, access_token, team_id, user) -> _FakeResolver:
        owner_filter = OwnerFilter.TEAM if team_id else OwnerFilter.PERSONAL
        self.captured.append(
            {
                "access_token": access_token,
                "team_id": team_id,
                "owner_filter": owner_filter,
            }
        )
        return self._resolver


def _patch_resolver(
    monkeypatch, known: Dict[str, str]
) -> _RecordingFolderResolverFactory:
    """Patch the endpoint's resolver factory with a fake over ``known`` folders."""
    factory = _RecordingFolderResolverFactory(_FakeResolver(known))
    monkeypatch.setattr(agent_controller, "_build_folder_resolver", factory)
    return factory


def _build_deck(slides: List[SlideSpec]) -> bytes:
    """Build a ``.pptx`` with one text box and one notes block per slide spec."""
    presentation = Presentation()
    blank_layout = presentation.slide_layouts[6]  # blank
    for body, notes in slides:
        slide = presentation.slides.add_slide(blank_layout)
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
        textbox.text_frame.text = body
        if notes:
            notes_frame = slide.notes_slide.notes_text_frame
            assert notes_frame is not None  # python-pptx creates it on access
            notes_frame.text = notes
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _post_deck(client: TestClient, deck: bytes, team_id: Optional[str] = None):
    data = {"team_id": team_id} if team_id is not None else None
    return client.post(
        _ANALYZE_URL,
        headers=_HEADERS,
        files={
            "file": (
                "template.pptx",
                deck,
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
        },
        data=data,
    )


def test_analyze_valid_template_returns_200_schema_and_empty_errors(
    client: TestClient,
):
    deck = _build_deck(
        [
            (
                "Hello {{name}} and {{role}}",
                "{{name}}:\nThe person's name\n{{role}}:\nTheir role",
            ),
            ("City: {{city}}", "{{city}}:\nThe city"),
        ]
    )

    response = _post_deck(client, deck)

    assert response.status_code == status.HTTP_200_OK
    body = response.json()
    assert set(body.keys()) == {"schema", "errors"}
    assert body["errors"] == []
    assert body["schema"] == [
        {
            "slide": 1,
            "keys": [
                {"key": "name", "description": "The person's name"},
                {"key": "role", "description": "Their role"},
            ],
        },
        {"slide": 2, "keys": [{"key": "city", "description": "The city"}]},
    ]


def test_analyze_template_with_errors_returns_200_with_structured_errors(
    client: TestClient,
):
    # {{age}} appears on slide 1 but is not described; {{ghost}} is described but absent.
    deck = _build_deck(
        [
            (
                "Hi {{name}} aged {{age}}",
                "{{name}}:\nThe name\n{{ghost}}:\nA stale description",
            )
        ]
    )

    response = _post_deck(client, deck)

    # Analyze returns 200 even with errors (schema + errors shown together).
    assert response.status_code == status.HTTP_200_OK
    body = response.json()
    assert set(body.keys()) == {"schema", "errors"}

    # Schema is still returned (the text-box keys, deduped per slide).
    assert body["schema"] == [
        {
            "slide": 1,
            "keys": [
                {"key": "name", "description": "The name"},
                {"key": "age", "description": ""},
            ],
        }
    ]

    # Every error carries slide, key, code, message.
    for error in body["errors"]:
        assert set(error.keys()) == {"slide", "key", "code", "message"}

    codes = {(e["code"], e["key"], e["slide"]) for e in body["errors"]}
    assert ("key_without_description", "age", 1) in codes
    assert ("described_but_not_in_slide", "ghost", 1) in codes


def test_analyze_reports_each_error_against_its_own_slide(client: TestClient):
    """Errors are attributed per-slide: a missing-description error on slide 1 and a
    ghost-key error on slide 2 each carry their own slide number (not conflated)."""
    deck = _build_deck(
        [
            # Slide 1: {{role}} has no description -> key_without_description on slide 1.
            ("Hello {{name}} the {{role}}", "{{name}}:\nThe name"),
            # Slide 2: {{country}} is described but absent -> described_but_not_in_slide on slide 2.
            ("City: {{city}}", "{{city}}:\nThe city\n{{country}}:\nA ghost country"),
        ]
    )

    response = _post_deck(client, deck)

    assert response.status_code == status.HTTP_200_OK
    body = response.json()
    assert set(body.keys()) == {"schema", "errors"}

    for error in body["errors"]:
        assert set(error.keys()) == {"slide", "key", "code", "message"}

    codes = {(e["code"], e["key"], e["slide"]) for e in body["errors"]}
    # Each error is pinned to the slide it actually occurs on.
    assert ("key_without_description", "role", 1) in codes
    assert ("described_but_not_in_slide", "country", 2) in codes
    # The missing-description error must NOT leak onto slide 2, nor the ghost onto slide 1.
    assert ("key_without_description", "role", 2) not in codes
    assert ("described_but_not_in_slide", "country", 1) not in codes


def test_analyze_rejects_non_pptx_upload_with_400(client: TestClient):
    response = client.post(
        _ANALYZE_URL,
        headers=_HEADERS,
        files={"file": ("notes.txt", b"this is not a pptx", "text/plain")},
    )

    assert response.status_code == status.HTTP_400_BAD_REQUEST


# --- space-aware image-folder resolution (Story 05) ------------------------------------

_IMAGE_NOTES = "{{flag}}:\n- type: image\n- folder: images/flags\nPick a flag."


def test_analyze_known_folder_returns_image_type_and_tag_id(client, monkeypatch):
    """An image key whose folder is known -> 200, schema key carries type:'image' and a
    resolved folder_tag_id, and no folder_not_found error."""
    _patch_resolver(monkeypatch, {"images/flags": "tag-123"})
    deck = _build_deck([("{{flag}}", _IMAGE_NOTES)])

    response = _post_deck(client, deck)

    assert response.status_code == status.HTTP_200_OK
    body = response.json()
    assert set(body.keys()) == {"schema", "errors"}
    assert body["schema"] == [
        {
            "slide": 1,
            "keys": [
                {
                    "key": "flag",
                    "description": "Pick a flag.",
                    "type": "image",
                    "folder": "images/flags",
                    "folder_tag_id": "tag-123",
                }
            ],
        }
    ]
    assert [e["code"] for e in body["errors"]] == []


def test_analyze_unknown_folder_returns_folder_not_found(client, monkeypatch):
    """An image key whose folder is unknown -> 200 with folder_not_found(slide, key)."""
    _patch_resolver(monkeypatch, {"images/other": "tag-999"})  # different folder
    deck = _build_deck([("{{flag}}", _IMAGE_NOTES)])

    response = _post_deck(client, deck)

    assert response.status_code == status.HTTP_200_OK
    body = response.json()
    codes = {(e["code"], e["key"], e["slide"]) for e in body["errors"]}
    assert ("folder_not_found", "flag", 1) in codes
    # No tag id was resolved, so the schema key keeps folder_tag_id unset.
    flag = body["schema"][0]["keys"][0]
    assert flag["type"] == "image"
    assert flag.get("folder_tag_id") is None


def test_analyze_forwards_team_id_with_team_scope(client, monkeypatch):
    """A provided team_id is forwarded to the resolver factory with TEAM scope."""
    factory = _patch_resolver(monkeypatch, {"images/flags": "tag-123"})
    deck = _build_deck([("{{flag}}", _IMAGE_NOTES)])

    response = _post_deck(client, deck, team_id="team-42")

    assert response.status_code == status.HTTP_200_OK
    assert factory.captured == [
        {
            "access_token": _TOKEN,
            "team_id": "team-42",
            "owner_filter": OwnerFilter.TEAM,
        }
    ]


def test_analyze_without_team_id_uses_personal_scope(client, monkeypatch):
    """No team_id -> the resolver is built with PERSONAL scope (team_id None)."""
    factory = _patch_resolver(monkeypatch, {"images/flags": "tag-123"})
    deck = _build_deck([("{{flag}}", _IMAGE_NOTES)])

    response = _post_deck(client, deck)

    assert response.status_code == status.HTTP_200_OK
    assert factory.captured == [
        {
            "access_token": _TOKEN,
            "team_id": None,
            "owner_filter": OwnerFilter.PERSONAL,
        }
    ]


def test_analyze_text_only_template_never_calls_resolver(client, monkeypatch):
    """A purely text template resolves no folders (resolver untouched) and stays clean."""
    factory = _patch_resolver(monkeypatch, {"images/flags": "tag-123"})
    deck = _build_deck([("Hello {{name}}", "{{name}}:\nThe name")])

    response = _post_deck(client, deck)

    assert response.status_code == status.HTTP_200_OK
    body = response.json()
    assert body["errors"] == []
    # The factory was invoked (scope is always prepared) but no folder was resolved.
    assert factory._resolver.calls == []
