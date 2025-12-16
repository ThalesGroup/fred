import logging
import re

from pptx import Presentation

logger = logging.getLogger(__name__)


def fill_slide_from_structured_response(ppt_path, structured_response, output_path):
    prs = Presentation(ppt_path)
    slide_numbers = [3, 4, 5]
    responses = structured_response.values()
    for slide_number, response in zip(slide_numbers, responses):
        # Support 1-based slide numbering
        slide = prs.slides[slide_number - 1]

        # Pattern: {name}, {email}, {phone}, etc.
        pattern = re.compile(r"\{([^}]+)\}")

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            for paragraph in shape.text_frame.paragraphs:  # type: ignore
                i = 0
                while i < len(paragraph.runs):
                    # merge text fields (called 'runs') that contain split placeholder keys
                    run = paragraph.runs[i]
                    text = run.text
                    if "{" in run and "}" not in run:
                        for run2 in paragraph.runs[i + 1 :]:
                            text += run2.text
                            i += 1
                            if "}" in text:
                                break
                    i += 1
                    # find placeholder keys and replace them with text
                    matches = pattern.findall(text)
                    for key in matches:
                        if key in response.keys():
                            text = text.replace(f"{{{key}}}", str(response[key]))

                    paragraph.runs[i].text = text

    prs.save(output_path)
    return output_path
