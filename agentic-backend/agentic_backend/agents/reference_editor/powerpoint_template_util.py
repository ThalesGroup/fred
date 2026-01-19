import logging
import re
from io import BytesIO

from docx import Document
from docx.shared import Inches
from pptx import Presentation
from pptx.util import Inches as PptxInches

from agentic_backend.agents.reference_editor.image_search_util import get_image_for_technology

logger = logging.getLogger(__name__)


def fill_slide_from_structured_response(ppt_path, structured_responses, output_path, vector_search_client=None, kf_base_client=None):
    prs = Presentation(ppt_path)
    pattern = re.compile(r"\{([^}]+)\}")

    # Flatten the nested structure into a single dictionary
    flattened_data = {}
    for section_name, section_data in structured_responses.items():
        if isinstance(section_data, dict):
            # Add all nested key-value pairs to the flattened dictionary
            for key, value in section_data.items():
                flattened_data[key] = value
        else:
            # If it's not a dict, keep it as is
            flattened_data[section_name] = section_data

    logger.info(f"Flattened data keys: {list(flattened_data.keys())}")

    # Pre-fetch images for listeTechnologies if present
    tech_images = []
    if "listeTechnologies" in flattened_data and vector_search_client and kf_base_client:
        technologies_text = flattened_data["listeTechnologies"]
        logger.info(f"Detected listeTechnologies field with value: {technologies_text}")
        tech_images = parse_technologies_and_fetch_images(technologies_text, vector_search_client, kf_base_client)

    # Process only the first slide (slide 0) for single-slide templates
    slide = prs.slides[0]

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        for paragraph in shape.text_frame.paragraphs:  # type: ignore
            # Merge all runs to find complete placeholders (handles split placeholders)
            full_text = "".join(run.text for run in paragraph.runs)

            # Find all placeholders and their positions in the merged text
            placeholder_replacements = {}
            has_liste_technologies = False
            for match in pattern.finditer(full_text):
                key = match.group(1)
                if key in flattened_data:
                    # Check if this is the listeTechnologies placeholder
                    if key == "listeTechnologies" and tech_images:
                        has_liste_technologies = True
                        # We'll handle this separately with images
                        placeholder_replacements[match.start()] = {
                            "end": match.end(),
                            "placeholder": match.group(0),
                            "value": "",  # Clear the text, we'll add images
                        }
                    else:
                        placeholder_replacements[match.start()] = {
                            "end": match.end(),
                            "placeholder": match.group(0),
                            "value": str(flattened_data[key]),
                        }
                else:
                    logger.warning(f"Placeholder '{key}' not found in flattened data")

            if not placeholder_replacements:
                continue

            # Sort placeholders by position to handle multiple placeholders correctly
            sorted_placeholders = sorted(
                placeholder_replacements.items(), key=lambda x: x[0]
            )

            # Process each run and replace placeholder text
            char_position = 0
            for run in paragraph.runs:
                run_start = char_position
                run_end = char_position + len(run.text)
                new_run_text = run.text
                offset = 0  # Track text length changes within this run

                for placeholder_start, replacement in sorted_placeholders:
                    placeholder_end = replacement["end"]

                    # Case 1: Placeholder starts in this run
                    if run_start <= placeholder_start < run_end:
                        local_start = placeholder_start - run_start + offset

                        if placeholder_end <= run_end:
                            # Placeholder is completely within this run
                            placeholder_len = len(replacement["placeholder"])
                            new_run_text = (
                                new_run_text[:local_start]
                                + replacement["value"]
                                + new_run_text[local_start + placeholder_len :]
                            )
                            offset += len(replacement["value"]) - placeholder_len
                        else:
                            # Placeholder spans multiple runs - replace the start
                            new_run_text = (
                                new_run_text[:local_start] + replacement["value"]
                            )

                    # Case 2: This run is in the middle/end of a placeholder
                    elif placeholder_start < run_start < placeholder_end:
                        if placeholder_end >= run_end:
                            # Entire run is part of the placeholder - empty it
                            new_run_text = ""
                        else:
                            # Placeholder ends in this run - keep text after placeholder
                            local_end = placeholder_end - run_start
                            new_run_text = new_run_text[local_end:]

                run.text = new_run_text
                char_position = run_end

            # After text replacement, if this shape contained listeTechnologies, add images
            if has_liste_technologies and tech_images:
                logger.info(f"Adding {len(tech_images)} technology images to PowerPoint shape")
                add_images_to_pptx_shape(shape, tech_images, slide)

    prs.save(output_path)
    return output_path


def fill_word_from_structured_response(docx_path, structured_responses, output_path, vector_search_client=None, kf_base_client=None):
    doc = Document(docx_path)
    pattern = re.compile(r"\{([^}]+)\}")

    logger.info("=== STARTING WORD TEMPLATE FILL ===")
    logger.info(f"Template path: {docx_path}")
    logger.info(f"Output path: {output_path}")

    # Log template content before processing
    logger.info("=== TEMPLATE CONTENT BEFORE PROCESSING ===")
    total_placeholders_in_template = 0
    for i, para in enumerate(doc.paragraphs):
        if para.text.strip():
            logger.info(f"Template Para {i}: {para.text}")
            placeholders = pattern.findall(para.text)
            if placeholders:
                total_placeholders_in_template += len(placeholders)
                logger.info(f"  Placeholders found: {placeholders}")

    for i, table in enumerate(doc.tables):
        for row_idx, row in enumerate(table.rows):
            for col_idx, cell in enumerate(row.cells):
                for para in cell.paragraphs:
                    if para.text.strip():
                        logger.info(f"Template Table {i}[{row_idx},{col_idx}]: {para.text}")
                        placeholders = pattern.findall(para.text)
                        if placeholders:
                            total_placeholders_in_template += len(placeholders)
                            logger.info(f"  Placeholders found: {placeholders}")

    logger.info(f"Total placeholders in template: {total_placeholders_in_template}")

    # Flatten the nested structure into a single dictionary
    flattened_data = {}
    for section_name, section_data in structured_responses.items():
        if isinstance(section_data, dict):
            # Add all nested key-value pairs to the flattened dictionary
            for key, value in section_data.items():
                flattened_data[key] = value
        else:
            # If it's not a dict, keep it as is
            flattened_data[section_name] = section_data

    logger.info("=== DATA FOR REPLACEMENT ===")
    logger.info(f"Flattened data keys: {list(flattened_data.keys())}")
    logger.info(f"Flattened data: {flattened_data}")

    # Pre-fetch images for listeTechnologies if present
    tech_images = []
    if "listeTechnologies" in flattened_data and vector_search_client and kf_base_client:
        technologies_text = flattened_data["listeTechnologies"]
        logger.info(f"Detected listeTechnologies field with value: {technologies_text}")
        tech_images = parse_technologies_and_fetch_images(technologies_text, vector_search_client, kf_base_client)

    # Helper function to process paragraphs with formatting preservation
    def process_paragraph(paragraph):
        # Merge all runs to find complete placeholders (handles split placeholders)
        full_text = "".join(run.text for run in paragraph.runs)

        if not full_text:
            return

        logger.debug(f"Processing paragraph text: '{full_text}'")

        # Find all placeholders and their positions in the merged text
        matches = list(pattern.finditer(full_text))
        if not matches:
            return

        logger.info(f"Found {len(matches)} placeholders in paragraph")

        # Build a list of placeholder replacements
        placeholder_replacements = {}
        has_liste_technologies = False
        for match in matches:
            key = match.group(1)
            logger.info(f"Processing placeholder: {{{key}}}")
            if key in flattened_data:
                # Check if this is the listeTechnologies placeholder
                if key == "listeTechnologies" and tech_images:
                    has_liste_technologies = True
                    # We'll handle this separately with images
                    placeholder_replacements[match.start()] = {
                        "end": match.end(),
                        "placeholder": match.group(0),
                        "value": "",  # Clear the text, we'll add images
                    }
                else:
                    value = str(flattened_data[key])
                    logger.info(f"Replacing {{{key}}} with: {value}")
                    placeholder_replacements[match.start()] = {
                        "end": match.end(),
                        "placeholder": match.group(0),
                        "value": value,
                    }
            else:
                logger.warning(f"Placeholder '{key}' not found in flattened data")

        if not placeholder_replacements:
            return

        # Sort placeholders by position
        sorted_placeholders = sorted(placeholder_replacements.items(), key=lambda x: x[0])

        # Process each run and replace placeholder text while preserving formatting
        char_position = 0
        for run in paragraph.runs:
            run_start = char_position
            run_end = char_position + len(run.text)
            new_run_text = run.text
            offset = 0  # Track text length changes within this run

            for placeholder_start, replacement in sorted_placeholders:
                placeholder_end = replacement["end"]

                # Case 1: Placeholder starts in this run
                if run_start <= placeholder_start < run_end:
                    local_start = placeholder_start - run_start + offset

                    if placeholder_end <= run_end:
                        # Placeholder is completely within this run
                        placeholder_len = len(replacement["placeholder"])
                        new_run_text = (
                            new_run_text[:local_start]
                            + replacement["value"]
                            + new_run_text[local_start + placeholder_len :]
                        )
                        offset += len(replacement["value"]) - placeholder_len
                    else:
                        # Placeholder spans multiple runs - replace the start
                        new_run_text = (
                            new_run_text[:local_start] + replacement["value"]
                        )
                        offset = len(new_run_text) - local_start

                # Case 2: This run is in the middle/end of a placeholder
                elif placeholder_start < run_start < placeholder_end:
                    if placeholder_end >= run_end:
                        # Entire run is part of the placeholder - empty it
                        new_run_text = ""
                        offset = -len(run.text)
                    else:
                        # Placeholder ends in this run - keep text after placeholder
                        local_end = placeholder_end - run_start
                        new_run_text = new_run_text[local_end:]
                        offset = -local_end

            run.text = new_run_text
            char_position = run_end

        # After text replacement, if this paragraph contained listeTechnologies, add images
        if has_liste_technologies and tech_images:
            logger.info(f"Adding {len(tech_images)} technology images to Word paragraph")
            add_images_to_word_paragraph(paragraph, tech_images)

        logger.info(f"Final paragraph text after replacement: '{paragraph.text}'")

    # Process all paragraphs in the document body
    for paragraph in doc.paragraphs:
        process_paragraph(paragraph)

    # Process tables in the document
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    process_paragraph(paragraph)

    # Process headers
    for section in doc.sections:
        # Process header (first, even, default)
        header = section.header
        for paragraph in header.paragraphs:
            process_paragraph(paragraph)
        # Process header tables
        for table in header.tables:
            for row in table.rows:
                for cell in row.cells:
                    for paragraph in cell.paragraphs:
                        process_paragraph(paragraph)

    # Process footers
    for section in doc.sections:
        # Process footer (first, even, default)
        footer = section.footer
        for paragraph in footer.paragraphs:
            process_paragraph(paragraph)
        # Process footer tables
        for table in footer.tables:
            for row in table.rows:
                for cell in row.cells:
                    for paragraph in cell.paragraphs:
                        process_paragraph(paragraph)

    # Process textboxes (VML and DrawingML) - CRITICAL for templates with text boxes
    # Many Word templates use textboxes for layout, and python-docx doesn't handle them well
    logger.info("=== PROCESSING TEXTBOXES ===")

    textbox_count = 0
    replacements_made = 0

    # Process all textbox content elements (both VML and DrawingML formats)
    for txbxContent in doc.element.body.iter():
        # Check if this is a textbox content element
        tag_name = txbxContent.tag.split('}')[-1] if '}' in txbxContent.tag else txbxContent.tag
        if tag_name != 'txbxContent':
            continue

        textbox_count += 1
        logger.info(f"Processing textbox {textbox_count}")

        # Find all text elements in this textbox
        for t_elem in txbxContent.iter():
            t_tag = t_elem.tag.split('}')[-1] if '}' in t_elem.tag else t_elem.tag
            if t_tag != 't':
                continue

            if t_elem.text and ('{' in t_elem.text or '}' in t_elem.text):
                original_text = t_elem.text
                logger.info(f"  Found text with potential placeholders: '{original_text}'")

                # Find and replace all placeholders in this text element
                new_text = original_text
                for match in pattern.finditer(original_text):
                    key = match.group(1)
                    placeholder = match.group(0)

                    if key in flattened_data:
                        value = str(flattened_data[key])
                        logger.info(f"    Replacing {placeholder} with: {value}")
                        new_text = new_text.replace(placeholder, value)
                        replacements_made += 1
                    else:
                        logger.warning(f"    Placeholder '{key}' not found in flattened data")

                if new_text != original_text:
                    t_elem.text = new_text
                    logger.info(f"    Updated text: '{new_text}'")

    logger.info(f"Processed {textbox_count} textboxes, made {replacements_made} replacements")

    doc.save(output_path)
    return output_path


def parse_technologies_and_fetch_images(technologies_text: str, vector_search_client, kf_base_client) -> list[tuple[str, BytesIO | None]]:
    """
    Parse a comma-separated list of technologies and fetch their corresponding images.

    Args:
        technologies_text: Comma-separated list of technology names (e.g., "Nvidia, Apple, AWS, SharePoint")
        kf_client: Authenticated KfBaseClient instance

    Returns:
        List of tuples (technology_name, image_data), where image_data is BytesIO or None if not found
    """
    if not technologies_text or not isinstance(technologies_text, str):
        return []

    # Split by commas and clean each technology name
    technologies = [tech.strip() for tech in technologies_text.split(",") if tech.strip()]

    logger.info(f"Parsing {len(technologies)} technologies: {technologies}")

    results = []
    for tech_name in technologies:
        logger.info(f"Fetching image for technology: {tech_name}")
        image_data = get_image_for_technology(tech_name, vector_search_client, kf_base_client)

        if image_data:
            logger.info(f"Successfully fetched image for: {tech_name}")
            results.append((tech_name, image_data))
        else:
            logger.warning(f"Could not fetch image for: {tech_name}, will use text instead")
            results.append((tech_name, None))

    logger.info(f"Fetched {sum(1 for _, img in results if img is not None)} images out of {len(technologies)} technologies")
    return results


def add_images_to_pptx_shape(shape, images: list[tuple[str, BytesIO | None]], slide):
    """
    Add images to a PowerPoint shape placeholder, replacing the text content.

    Args:
        shape: The PowerPoint shape containing the placeholder
        images: List of tuples (technology_name, image_data)
        slide: The slide object to add images to
    """
    if not images:
        return

    # Clear the text content from the shape
    if shape.has_text_frame:
        shape.text_frame.clear()

    # Calculate positioning for images
    # Get the shape's position and size
    left = shape.left
    top = shape.top
    width = shape.width
    height = shape.height

    # Calculate image dimensions (we'll arrange them horizontally)
    num_images = len([img for _, img in images if img is not None])
    if num_images == 0:
        # No images, just set text
        shape.text = ", ".join([name for name, _ in images])
        return

    # Image size and spacing
    image_width = width // num_images
    image_height = min(height, PptxInches(0.8))  # Max height of 0.8 inches

    # Add each image
    current_left = left
    for tech_name, image_data in images:
        if image_data is not None:
            try:
                # Reset the BytesIO pointer to the beginning
                image_data.seek(0)

                # Add the image to the slide
                slide.shapes.add_picture(
                    image_data,
                    current_left,
                    top + (height - image_height) // 2,  # Center vertically
                    width=image_width,
                    height=image_height
                )
                current_left += image_width
                logger.info(f"Added image for technology: {tech_name}")
            except Exception as e:
                logger.error(f"Failed to add image for {tech_name}: {e}")


def add_images_to_word_paragraph(paragraph, images: list[tuple[str, BytesIO | None]]):
    """
    Add images to a Word paragraph, replacing the text content.

    Args:
        paragraph: The Word paragraph object
        images: List of tuples (technology_name, image_data)
    """
    if not images:
        return

    # Clear existing runs
    for run in paragraph.runs:
        run.text = ""

    # Add each image or text
    for i, (tech_name, image_data) in enumerate(images):
        if image_data is not None:
            try:
                # Reset the BytesIO pointer
                image_data.seek(0)

                # Add image to the paragraph
                run = paragraph.add_run()
                run.add_picture(image_data, width=Inches(0.5))

                # Add spacing between images
                if i < len(images) - 1:
                    paragraph.add_run("  ")

                logger.info(f"Added image for technology: {tech_name} to Word document")
            except Exception as e:
                logger.error(f"Failed to add image for {tech_name} to Word: {e}")
                # Fallback to text
                paragraph.add_run(tech_name)
                if i < len(images) - 1:
                    paragraph.add_run(", ")
        else:
            # No image found, use text
            paragraph.add_run(tech_name)
            if i < len(images) - 1:
                paragraph.add_run(", ")


referenceSchema = {
    "type": "object",
    "properties": {
        "informationsProjet": {
            "type": "object",
            "description": "Informations sur le projet",
            "properties": {
                "nomSociete": {
                    "type": "string",
                    "description": "Le nom de la société",
                    "maxLength": 50,
                },
                "nomProjet": {
                    "type": "string",
                    "description": "Le nom du projet",
                    "maxLength": 50,
                },
                "dateProjet": {
                    "type": "string",
                    "description": "La date de début et de fin du projet",
                },
                "nombrePersonnes": {
                    "type": "string",
                    "description": "Le nombre de personnes dans la direction (uniquement dans la direction)",
                },
                "enjeuFinancier": {
                    "type": "string",
                    "description": "Les coûts financiers ou la rentabilité du projet exprimé en euros, juste un seul chiffre clé comme le CA (jamais en nombre de personnes sinon ne rien mettre)",
                    "maxLength": 100,
                },
            },
        },
        "contexte": {
            "type": "object",
            "description": "Informations sur le contexte et le client",
            "properties": {
                "presentationClient": {
                    "type": "string",
                    "description": "Courte présentation du client. Longueur maximale: 300 caractères (une à deux phrases).",
                    "maxLength": 300,
                },
                "presentationContexte": {
                    "type": "string",
                    "description": "Courte présentation du contexte du projet. Longueur maximale: 300 caractères (une à deux phrases).",
                    "maxLength": 300,
                },
                "listeTechnologies": {
                    "type": "string",
                    "description": "Listes des technologies utilisés lors de ce projet",
                },
            },
        },
        "syntheseProjet": {
            "type": "object",
            "description": "Synthèse struturée du projet",
            "properties": {
                "enjeux": {
                    "type": "string",
                    "description": "Court résumé des enjeux du projet. Longueur maximale: 300 caractères (une à deux phrases).",
                    "maxLength": 300,
                },
                "activiteSolutions": {
                    "type": "string",
                    "description": "Court résumé des activités et solutions du projet. Longueur maximale: 300 caractères (une à deux phrases).",
                    "maxLength": 300,
                },
                "beneficeClients": {
                    "type": "string",
                    "description": "Court résumé des bénéfices pour le client. Longueur maximale: 300 caractères (une à deux phrases).",
                },
                "pointsForts": {
                    "type": "string",
                    "description": "Court résumé des points forts du projet. Longueur maximale: 300 caractères (une à deux phrases).",
                },
            },
        },
    },
}
