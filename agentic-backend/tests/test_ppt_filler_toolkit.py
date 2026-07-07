"""Offline, fixture-driven tests for the PPT Filler fill tool (PPTFILL-05 / #1834).

Decks are built in-test with python-pptx (no checked-in binaries), reusing the helper
style of ``tests/test_ppt_filler_parser.py``. The workspace client is faked the way
``tests/test_kf_vector_search_tools.py`` fakes its HTTP client, so everything is offline:
no network, no real ApplicationContext.

Coverage (the acceptance criteria of #1834):

- The dynamic ``args_schema`` is generated per-slide from a given ``schema_slides``, with
  each leaf carrying its note description.
- Filling uses the shared traversal and fills repeated keys on a slide with the same value.
- The filled deck is uploaded to SESSION-SCOPED user storage and a download ``LinkPart`` is
  returned.
- No remaining ``{{keys}}`` after filling all provided fields.
"""

import io
from typing import List, Optional, Tuple

import pytest
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from agentic_backend.common import kf_workspace_client as kf_ws
from agentic_backend.common.kf_document_client import RawContentBlob
from agentic_backend.common.kf_workspace_client import (
    UserStorageBlob,
    UserStorageUploadResult,
)
from agentic_backend.common.structures import AgentSettings, AgentTuning
from agentic_backend.core.agents.agent_spec import MCPServerRef
from agentic_backend.core.agents.runtime_context import RuntimeContext
from agentic_backend.core.chatbot.chat_schema import LinkKind, LinkPart
from agentic_backend.integrations.ppt_filler.parser import parse
from agentic_backend.integrations.ppt_filler.ppt_filler_params import (
    PPT_FILLER_TEMPLATE_KEY,
    PptFillerParams,
)
from agentic_backend.integrations.ppt_filler.toolkit import build_ppt_filler_tools
from agentic_backend.integrations.ppt_filler.traversal import (
    KEY_PATTERN,
    list_keys_on_slide,
)

_PPTX_CONTENT_TYPE = (
    "application/vnd.openxmlformats-officedocument.presentationml.presentation"
)


# --- Deck-builder fixtures (mirrors tests/test_ppt_filler_parser.py) --------------

# (text-box body, notes text) per slide.
SlideSpec = Tuple[str, str]


def _build_deck(slides: List[SlideSpec]) -> bytes:
    """Build a ``.pptx`` with one text box and one notes block per slide spec."""
    presentation = Presentation()
    blank_layout = presentation.slide_layouts[6]  # blank
    for body, notes in slides:
        slide = presentation.slides.add_slide(blank_layout)
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
        textbox.text_frame.text = body
        if notes:
            slide.notes_slide.notes_text_frame.text = notes
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _schema_slides(deck: bytes):
    """Parse a deck and return its persisted-style ``schema_slides`` list."""
    return parse(deck).slides


# --- Fakes ------------------------------------------------------------------------


class _FakeAgent:
    """Minimal KnowledgeFlowAgentContext stand-in carrying ppt_filler params."""

    def __init__(self, *, params: PptFillerParams, session_id: Optional[str] = None):
        self.agent_settings = AgentSettings(
            id="agent-1",
            name="PPT Filler",
            tuning=AgentTuning(
                role="r",
                description="d",
                mcp_servers=[MCPServerRef(id="mcp-ppt-filler", params=params)],
            ),
        )
        self.runtime_context = RuntimeContext(session_id=session_id, access_token="tok")

    def refresh_user_access_token(self) -> str:
        raise NotImplementedError("not exercised by these tests")


class _FakeWorkspaceClient:
    """Records fetch/upload calls and serves a configured template blob.

    Substituted for the real ``KfWorkspaceClient`` so the tool runs fully offline.
    """

    template_bytes: bytes = b""
    fetch_calls: list[dict] = []
    upload_calls: list[dict] = []

    def __init__(self, *args, **kwargs):
        # The tool constructs ``KfWorkspaceClient(agent=agent)``; accept and ignore.
        pass

    async def fetch_agent_config_blob(self, key, access_token=None, agent_id=None):
        type(self).fetch_calls.append(
            {"key": key, "access_token": access_token, "agent_id": agent_id}
        )
        return UserStorageBlob(
            bytes=type(self).template_bytes,
            content_type=_PPTX_CONTENT_TYPE,
            filename=key,
            size=len(type(self).template_bytes),
        )

    async def upload_user_blob(self, key, file_content, filename, content_type=None):
        content = (
            file_content
            if isinstance(file_content, (bytes, bytearray))
            else file_content.read()
        )
        type(self).upload_calls.append(
            {
                "key": key,
                "content": bytes(content),
                "filename": filename,
                "content_type": content_type,
            }
        )
        return UserStorageUploadResult(
            key=key,
            file_name=filename,
            size=len(content),
            document_uid="doc-uid-123",
            download_url=f"https://example.test/storage/user/{key}",
        )


@pytest.fixture
def fake_ws(monkeypatch):
    """Install a fresh fake KfWorkspaceClient for one test.

    By default the PDF preview is disabled (``convert_pptx_bytes_to_pdf`` stubbed to return
    ``None``) so the fill returns the plain download chip — this keeps the offline suite
    independent of a real ``soffice`` binary. Tests that exercise the preview path override
    the stub (see ``_enable_preview``).
    """

    class _Client(_FakeWorkspaceClient):
        template_bytes = b""
        fetch_calls = []
        upload_calls = []

    monkeypatch.setattr(kf_ws, "KfWorkspaceClient", _Client)
    # The toolkit imported the symbol by name, so patch it there too.
    import agentic_backend.integrations.ppt_filler.toolkit as toolkit_mod

    monkeypatch.setattr(toolkit_mod, "KfWorkspaceClient", _Client)

    # Preview off by default: no PDF is produced, so no second upload happens.
    async def _no_pdf(_bytes, *args, **kwargs):
        return None

    monkeypatch.setattr(toolkit_mod, "convert_pptx_bytes_to_pdf", _no_pdf)
    return _Client


def _enable_preview(monkeypatch, fake_ws, *, pdf_bytes: bytes = b"%PDF-1.5 fake"):
    """Turn the best-effort preview on for a test by stubbing conversion to yield PDF bytes.

    No presigned URL is stubbed: the tool no longer presigns at fill time (a presigned URL
    would expire while the part is persisted). It instead derives a DURABLE presign href from
    the .pptx download URL, which the fake upload already returns.
    """
    import agentic_backend.integrations.ppt_filler.toolkit as toolkit_mod

    async def _pdf(_bytes, *args, **kwargs):
        return pdf_bytes

    monkeypatch.setattr(toolkit_mod, "convert_pptx_bytes_to_pdf", _pdf)


def _the_tool(agent):
    tools = build_ppt_filler_tools(agent)
    assert len(tools) == 1
    return tools[0]


# --- Image fakes ------------------------------------------------------------------


def _png_bytes(width: int = 10, height: int = 10, color: str = "red") -> bytes:
    """A tiny, valid PNG with the given pixel dimensions (no checked-in binaries)."""
    buffer = io.BytesIO()
    Image.new("RGB", (width, height), color).save(buffer, format="PNG")
    return buffer.getvalue()


def _webp_bytes(width: int = 10, height: int = 10, color: str = "red") -> bytes:
    """A tiny, valid WEBP. Pillow decodes it but python-pptx's ``add_picture`` cannot
    embed WEBP, so the toolkit must transcode it to PNG before insertion."""
    buffer = io.BytesIO()
    Image.new("RGB", (width, height), color).save(buffer, format="WEBP")
    return buffer.getvalue()


class _FakeDocumentClient:
    """Serves configured raw bytes for ``fetch_raw_content``; records the uids asked for.

    Substituted for the real ``KfDocumentClient`` so image fetches run fully offline (the
    real one needs an initialized ApplicationContext, which the tests do not have).
    """

    # {document_uid: (bytes, content_type)} served by fetch_raw_content.
    docs: dict = {}
    fetch_uids: list = []

    def __init__(self, *args, **kwargs):
        pass

    async def fetch_raw_content(self, *, document_uid: str) -> RawContentBlob:
        type(self).fetch_uids.append(document_uid)
        if document_uid not in type(self).docs:
            raise RuntimeError(f"no such document: {document_uid}")
        content, content_type = type(self).docs[document_uid]
        return RawContentBlob(
            bytes=content,
            content_type=content_type,
            filename=f"{document_uid}.bin",
            size=len(content),
        )


@pytest.fixture
def fake_doc(monkeypatch):
    """Install a fresh fake KfDocumentClient for one test (patched on the toolkit)."""

    class _Client(_FakeDocumentClient):
        docs = {}
        fetch_uids = []

    import agentic_backend.integrations.ppt_filler.toolkit as toolkit_mod

    monkeypatch.setattr(toolkit_mod, "KfDocumentClient", _Client)
    return _Client


def _build_image_deck(slides):
    """Build a deck whose slides each carry one or more textbox specs.

    ``slides`` is a list of ``(notes, [textbox_text, ...])``. Each textbox is added at a
    distinct, fixed box so image anchors have a real geometry to fit inside.
    """
    presentation = Presentation()
    blank_layout = presentation.slide_layouts[6]  # blank
    for notes, textboxes in slides:
        slide = presentation.slides.add_slide(blank_layout)
        for offset, text in enumerate(textboxes):
            box = slide.shapes.add_textbox(
                Inches(1),
                Inches(1 + 2 * offset),
                Inches(4),  # box wider than tall -> a square image fits the HEIGHT
                Inches(2),
            )
            box.text_frame.text = text
        if notes:
            slide.notes_slide.notes_text_frame.text = notes
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _image_schema(deck: bytes, *, slide: int, key: str, tag_id: str = "tag-xyz"):
    """Parse a deck and stamp a resolved ``folder_tag_id`` on one image key.

    Story 05 (the save processor) is what resolves the folder to a tag id; here we set it
    directly on the parsed schema so the fill tool sees the same persisted shape.
    """
    schema_slides = parse(deck).slides
    for slide_schema in schema_slides:
        if slide_schema.slide != slide:
            continue
        for key_field in slide_schema.keys:
            if key_field.key == key:
                key_field.folder_tag_id = tag_id
    return schema_slides


def _picture_shapes(slide) -> list:
    """Every PICTURE shape on a slide (python-pptx exposes ``shape_type`` PICTURE=13)."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    return [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]


# --- A. Dynamic per-slide args_schema ---------------------------------------------


def test_args_schema_is_nested_by_slide_with_leaf_descriptions():
    deck = _build_deck(
        [
            (
                "Hello {{name}} and {{role}}",
                "{{name}}:\nThe person's name\n{{role}}:\nTheir role",
            ),
            ("City: {{city}}", "{{city}}:\nThe city"),
        ]
    )
    params = PptFillerParams(schema=_schema_slides(deck))
    tool = _the_tool(_FakeAgent(params=params))

    schema = tool.args_schema.model_json_schema()

    # Top-level groups are keyed by slide_<n> (1-based slide numbers), plus an optional
    # output_file_name the model can set to name the produced deck.
    assert set(schema["properties"]) == {"slide_1", "slide_2", "output_file_name"}
    # Slide groups are required; output_file_name is optional.
    assert set(schema["required"]) == {"slide_1", "slide_2"}
    assert "output_file_name" not in schema.get("required", [])

    # Resolve the $ref of each slide group to its leaf object and assert leaf keys +
    # descriptions are present.
    defs = schema["$defs"]

    def _leaf(slide_field: str) -> dict:
        ref = schema["properties"][slide_field]["$ref"]
        return defs[ref.split("/")[-1]]

    # Each TEXT leaf keeps its note as the lead and appends the inline-formatting hint so
    # the agent knows it may bold/italic part of the value.
    from agentic_backend.integrations.ppt_filler.toolkit import _TEXT_FORMATTING_HINT

    slide1 = _leaf("slide_1")
    assert set(slide1["properties"]) == {"name", "role"}
    assert (
        slide1["properties"]["name"]["description"]
        == f"The person's name {_TEXT_FORMATTING_HINT}"
    )
    assert (
        slide1["properties"]["role"]["description"]
        == f"Their role {_TEXT_FORMATTING_HINT}"
    )

    slide2 = _leaf("slide_2")
    assert set(slide2["properties"]) == {"city"}
    assert (
        slide2["properties"]["city"]["description"]
        == f"The city {_TEXT_FORMATTING_HINT}"
    )


def test_no_tool_when_schema_is_empty():
    """No persisted schema (template never configured) -> nothing to fill -> no tool."""
    agent = _FakeAgent(params=PptFillerParams())
    assert build_ppt_filler_tools(agent) == []


# --- B/C/D. Fill via shared traversal, upload, return LinkPart --------------------


@pytest.mark.asyncio
async def test_fill_repeated_keys_on_slide_uses_one_value(fake_ws):
    """A key appearing twice on one slide is filled with the same value (the shared
    run-merging traversal does this), and the deck is uploaded session-scoped with a
    download LinkPart returned."""
    deck = _build_deck(
        [
            (
                "Title {{name}} ... footer {{name}}",
                "{{name}}:\nThe name",
            ),
        ]
    )
    fake_ws.template_bytes = deck
    params = PptFillerParams(schema=_schema_slides(deck), template_key="custom.pptx")
    tool = _the_tool(_FakeAgent(params=params, session_id="sess-9"))

    content, artifact = await tool.coroutine(slide_1={"name": "Jane Doe"})

    # Template fetched from agent-config storage under the configured key + agent_id.
    assert fake_ws.fetch_calls == [
        {"key": "custom.pptx", "access_token": "tok", "agent_id": "agent-1"}
    ]

    # Uploaded to SESSION-SCOPED user storage.
    assert len(fake_ws.upload_calls) == 1
    upload = fake_ws.upload_calls[0]
    assert upload["key"].startswith("sess-9/")
    assert upload["content_type"] == _PPTX_CONTENT_TYPE

    # The uploaded deck has both occurrences filled with the same value and no leftover
    # placeholders.
    filled = Presentation(io.BytesIO(upload["content"]))
    text = "".join(
        run.text
        for shape in filled.slides[0].shapes
        if shape.has_text_frame
        for paragraph in shape.text_frame.paragraphs
        for run in paragraph.runs
    )
    assert text.count("Jane Doe") == 2
    assert not KEY_PATTERN.search(text)

    # The returned artifact carries a download LinkPart with sensible name/mime.
    assert artifact.is_error is False
    assert len(artifact.ui_parts) == 1
    link = artifact.ui_parts[0]
    assert isinstance(link, LinkPart)
    assert link.kind == LinkKind.download
    assert link.mime == _PPTX_CONTENT_TYPE
    assert link.file_name == "filled_presentation.pptx"
    assert link.href and link.href.startswith("https://example.test/storage/user/")
    assert isinstance(content, str) and content.strip()


# --- PDF preview (best-effort) ----------------------------------------------------


@pytest.mark.asyncio
async def test_successful_fill_emits_ppt_preview_and_uploads_pdf(fake_ws, monkeypatch):
    """A successful fill converts to PDF, uploads it beside the .pptx, and emits a
    ppt_preview part carrying a DURABLE presign href (not a frozen presigned URL) + a version
    + the .pptx download, instead of a standalone download LinkPart."""
    from agentic_backend.core.chatbot.chat_schema import PptPreviewPart

    _enable_preview(monkeypatch, fake_ws)
    deck = _build_deck([("{{name}}", "{{name}}:\nThe name")])
    fake_ws.template_bytes = deck
    params = PptFillerParams(schema=_schema_slides(deck))
    tool = _the_tool(_FakeAgent(params=params, session_id="sess-9"))

    _content, artifact = await tool.coroutine(slide_1={"name": "Jane"})

    assert artifact.is_error is False
    # Two uploads: the .pptx then the preview .pdf, both session-scoped and sibling keys.
    assert len(fake_ws.upload_calls) == 2
    pptx_up, pdf_up = fake_ws.upload_calls
    assert pptx_up["key"].endswith(".pptx") and pptx_up["key"].startswith("sess-9/")
    assert pdf_up["key"] == pptx_up["key"][: -len(".pptx")] + ".pdf"
    assert pdf_up["content_type"] == "application/pdf"

    # Exactly one ppt_preview part, no standalone download LinkPart.
    assert len(artifact.ui_parts) == 1
    part = artifact.ui_parts[0]
    assert isinstance(part, PptPreviewPart)
    assert not any(isinstance(p, LinkPart) for p in artifact.ui_parts)
    # The preview points at the DURABLE presign endpoint for the PDF key (the browser mints a
    # fresh presigned URL from it at open time), NOT a baked, expiring presigned URL.
    assert (
        part.pdf_presign_url
        == f"https://example.test/storage/user/presigned/{pdf_up['key']}"
    )
    assert part.version  # a freshness token is stamped
    assert part.pptx_download_url and part.pptx_download_url.startswith(
        "https://example.test/"
    )
    assert part.file_name == "filled_presentation.pptx"


@pytest.mark.asyncio
async def test_failed_conversion_still_returns_pptx_with_note(fake_ws):
    """When conversion is unavailable (the default), the fill still returns the .pptx via a
    download LinkPart and the message states the preview could not be generated."""
    deck = _build_deck([("{{name}}", "{{name}}:\nThe name")])
    fake_ws.template_bytes = deck
    params = PptFillerParams(schema=_schema_slides(deck))
    tool = _the_tool(_FakeAgent(params=params, session_id="sess-9"))

    content, artifact = await tool.coroutine(slide_1={"name": "Jane"})

    assert artifact.is_error is False
    # No PDF upload happened; only the .pptx was stored.
    assert len(fake_ws.upload_calls) == 1
    assert len(artifact.ui_parts) == 1
    assert isinstance(artifact.ui_parts[0], LinkPart)
    assert artifact.ui_parts[0].kind == LinkKind.download
    assert "preview" in content.lower()


@pytest.mark.asyncio
async def test_preview_version_changes_on_refill(fake_ws, monkeypatch):
    """A re-fill that changes the deck yields a different preview version token (the
    cache-busting key that makes the open pane update live)."""
    from agentic_backend.core.chatbot.chat_schema import PptPreviewPart

    deck = _build_deck([("{{name}}", "{{name}}:\nThe name")])
    fake_ws.template_bytes = deck
    params = PptFillerParams(schema=_schema_slides(deck))
    tool = _the_tool(_FakeAgent(params=params, session_id="sess-9"))

    # First fill: the PDF conversion yields one set of bytes.
    _enable_preview(monkeypatch, fake_ws, pdf_bytes=b"%PDF-1.5 first")
    _c1, a1 = await tool.coroutine(slide_1={"name": "Jane"})
    v1 = a1.ui_parts[0]
    assert isinstance(v1, PptPreviewPart)

    # Second fill producing a different PDF (a real re-fill re-renders the deck).
    _enable_preview(monkeypatch, fake_ws, pdf_bytes=b"%PDF-1.5 second-different")
    _c2, a2 = await tool.coroutine(slide_1={"name": "Bob"})
    v2 = a2.ui_parts[0]
    assert isinstance(v2, PptPreviewPart)

    assert v1.version != v2.version


@pytest.mark.asyncio
async def test_fill_all_fields_leaves_no_placeholders(fake_ws):
    """Filling every provided field across multiple slides leaves NO {{keys}} in the
    produced bytes (re-open and assert)."""
    deck = _build_deck(
        [
            ("Hello {{name}}, the {{role}}", "{{name}}, {{role}}:\nName then role"),
            (
                "Repeated {{name}} and {{city}}",
                "{{name}}:\nName again\n{{city}}:\nThe city",
            ),
        ]
    )
    fake_ws.template_bytes = deck
    params = PptFillerParams(schema=_schema_slides(deck))
    tool = _the_tool(_FakeAgent(params=params, session_id="sess-1"))

    content, artifact = await tool.coroutine(
        slide_1={"name": "Alice", "role": "Engineer"},
        slide_2={"name": "Alice", "city": "Paris"},
    )

    assert artifact.is_error is False
    upload = fake_ws.upload_calls[0]
    refilled = Presentation(io.BytesIO(upload["content"]))
    for slide in refilled.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            assert not KEY_PATTERN.search(shape.text_frame.text)


@pytest.mark.asyncio
async def test_fill_without_session_uses_unscoped_key(fake_ws):
    """No session_id -> the upload key is the bare output file name (no session prefix)."""
    deck = _build_deck([("{{name}}", "{{name}}:\nThe name")])
    fake_ws.template_bytes = deck
    params = PptFillerParams(schema=_schema_slides(deck))
    tool = _the_tool(_FakeAgent(params=params, session_id=None))

    await tool.coroutine(slide_1={"name": "Bob"})

    assert fake_ws.upload_calls[0]["key"] == "filled_presentation.pptx"


@pytest.mark.asyncio
async def test_default_template_key_is_used(fake_ws):
    """With the default params, the template is fetched under PPT_FILLER_TEMPLATE_KEY."""
    deck = _build_deck([("{{name}}", "{{name}}:\nThe name")])
    fake_ws.template_bytes = deck
    params = PptFillerParams(schema=_schema_slides(deck))
    tool = _the_tool(_FakeAgent(params=params, session_id="s"))

    await tool.coroutine(slide_1={"name": "X"})

    assert fake_ws.fetch_calls[0]["key"] == PPT_FILLER_TEMPLATE_KEY


@pytest.mark.asyncio
async def test_template_fetch_failure_returns_error_result(fake_ws):
    """A template fetch failure becomes an is_error tool result, not a raised exception."""
    deck = _build_deck([("{{name}}", "{{name}}:\nThe name")])
    params = PptFillerParams(schema=_schema_slides(deck))
    tool = _the_tool(_FakeAgent(params=params, session_id="s"))

    async def boom(self, *a, **k):
        raise RuntimeError("storage down")

    fake_ws.fetch_agent_config_blob = boom  # type: ignore[assignment]

    content, artifact = await tool.coroutine(slide_1={"name": "X"})

    assert artifact.is_error is True
    assert artifact.ui_parts == ()
    assert "template" in content.lower()


# --- E. Runtime invocation path preserves the download artifact -------------------


@pytest.mark.asyncio
async def test_ainvoke_with_plain_args_preserves_download_artifact(fake_ws):
    """Regression guard for the dropped download link.

    The v2 ReAct resolver invokes inprocess provider tools via
    ``tool.ainvoke(<plain args dict>)`` and then normalizes the raw return value
    (see ``react_tool_resolution._resolve_runtime_provider_tool``). If the tool were
    declared with ``response_format="content_and_artifact"``, LangChain would return
    ONLY the content string on a plain-args invoke and silently drop the artifact —
    so the download ``LinkPart`` would never reach the UI. This test exercises that
    exact path end to end: ``ainvoke`` must yield the bare ``(content, artifact)``
    tuple, and the resolver's normalization must surface the download link.
    """
    from agentic_backend.core.agents.v2.react.react_tool_rendering import (
        normalize_runtime_provider_artifact,
        render_tool_result,
        stringify_tool_output,
    )

    deck = _build_deck([("{{name}}", "{{name}}:\nThe name")])
    fake_ws.template_bytes = deck
    params = PptFillerParams(schema=_schema_slides(deck))
    tool = _the_tool(_FakeAgent(params=params, session_id="sess-rt"))

    # The runtime calls ainvoke with PLAIN ARGS (not a ToolCall dict). With
    # content_and_artifact this would collapse to a bare content string.
    raw = await tool.ainvoke({"slide_1": {"name": "Ada"}})
    assert isinstance(raw, tuple) and len(raw) == 2, (
        "fill tool must return a (content, artifact) tuple on a plain-args ainvoke; "
        "do not set response_format='content_and_artifact' (it drops the artifact)"
    )

    # Mirror the resolver's normalization branch and assert the link survives.
    rendered = stringify_tool_output(raw[0]).strip()
    artifact = normalize_runtime_provider_artifact(raw[1])
    content = rendered if rendered else render_tool_result(artifact)

    assert content.strip()
    assert artifact is not None
    assert artifact.is_error is False
    assert len(artifact.ui_parts) == 1
    link = artifact.ui_parts[0]
    assert isinstance(link, LinkPart)
    assert link.kind == LinkKind.download
    assert link.file_name == "filled_presentation.pptx"


# --- F. Model-chosen output file name ---------------------------------------------


@pytest.mark.asyncio
async def test_output_file_name_is_used_for_key_and_download(fake_ws):
    """A model-provided output_file_name names the stored blob and the download link,
    and the ".pptx" extension is added when omitted."""
    deck = _build_deck([("{{name}}", "{{name}}:\nThe name")])
    fake_ws.template_bytes = deck
    params = PptFillerParams(schema=_schema_slides(deck))
    tool = _the_tool(_FakeAgent(params=params, session_id="sess-7"))

    _content, artifact = await tool.coroutine(
        slide_1={"name": "Ada"}, output_file_name="Proposition ACME 2026"
    )

    # Stored under the session-scoped key using the chosen name + added extension.
    upload = fake_ws.upload_calls[0]
    assert upload["key"] == "sess-7/Proposition ACME 2026.pptx"
    assert upload["filename"] == "Proposition ACME 2026.pptx"
    # The download link carries the same name.
    assert artifact.ui_parts[0].file_name == "Proposition ACME 2026.pptx"


@pytest.mark.asyncio
async def test_output_file_name_existing_extension_not_doubled(fake_ws):
    """If the model already includes a (case-insensitive) .pptx, it is not doubled."""
    deck = _build_deck([("{{name}}", "{{name}}:\nThe name")])
    fake_ws.template_bytes = deck
    params = PptFillerParams(schema=_schema_slides(deck))
    tool = _the_tool(_FakeAgent(params=params, session_id="s"))

    await tool.coroutine(slide_1={"name": "X"}, output_file_name="Deck.PPTX")

    assert fake_ws.upload_calls[0]["filename"] == "Deck.pptx"


@pytest.mark.asyncio
async def test_output_file_name_strips_path_components(fake_ws):
    """Path separators are stripped (no traversal): only the final component is kept."""
    deck = _build_deck([("{{name}}", "{{name}}:\nThe name")])
    fake_ws.template_bytes = deck
    params = PptFillerParams(schema=_schema_slides(deck))
    tool = _the_tool(_FakeAgent(params=params, session_id="s"))

    await tool.coroutine(slide_1={"name": "X"}, output_file_name="../../etc/evil.pptx")

    upload = fake_ws.upload_calls[0]
    assert upload["filename"] == "evil.pptx"
    assert upload["key"] == "s/evil.pptx"


@pytest.mark.asyncio
async def test_blank_output_file_name_falls_back_to_default(fake_ws):
    """A blank/whitespace name (or just an extension) falls back to the default."""
    deck = _build_deck([("{{name}}", "{{name}}:\nThe name")])
    fake_ws.template_bytes = deck
    params = PptFillerParams(schema=_schema_slides(deck))
    tool = _the_tool(_FakeAgent(params=params, session_id="s"))

    await tool.coroutine(slide_1={"name": "X"}, output_file_name="   ")
    assert fake_ws.upload_calls[0]["filename"] == "filled_presentation.pptx"


# --- G. Notes: authoring guidance stripped; kept content preserved ----------------


def _slide_notes(presentation, index: int) -> str:
    slide = list(presentation.slides)[index]
    if not slide.has_notes_slide:
        return ""
    return slide.notes_slide.notes_text_frame.text or ""


@pytest.mark.asyncio
async def test_filled_deck_strips_authoring_notes(fake_ws):
    """The ``{{key}}:`` descriptions are internal guidance and must NOT appear in the
    filled deck. A slide whose notes are authoring-only ends up with empty notes."""
    deck = _build_deck([("Hello {{name}}", "{{name}}:\nThe person's name")])
    fake_ws.template_bytes = deck
    params = PptFillerParams(schema=_schema_slides(deck))
    tool = _the_tool(_FakeAgent(params=params, session_id="s"))

    await tool.coroutine(slide_1={"name": "Ada"})

    filled = Presentation(io.BytesIO(fake_ws.upload_calls[0]["content"]))
    notes = _slide_notes(filled, 0)
    assert "{{name}}" not in notes
    assert "The person's name" not in notes
    assert notes == ""


@pytest.mark.asyncio
async def test_filled_deck_keeps_notes_after_separator(fake_ws):
    """Content after a ``---`` keep-separator is preserved as the slide's real notes,
    while the authoring description above it is removed."""
    deck = _build_deck(
        [
            (
                "Hello {{name}}",
                "{{name}}:\nThe person's name\n---\nSpeaker note: pause here.",
            )
        ]
    )
    fake_ws.template_bytes = deck
    params = PptFillerParams(schema=_schema_slides(deck))
    tool = _the_tool(_FakeAgent(params=params, session_id="s"))

    await tool.coroutine(slide_1={"name": "Ada"})

    filled = Presentation(io.BytesIO(fake_ws.upload_calls[0]["content"]))
    notes = _slide_notes(filled, 0)
    assert notes == "Speaker note: pause here."
    assert "{{name}}" not in notes
    # And the body was still filled.
    body = "".join(
        run.text
        for shape in list(filled.slides)[0].shapes
        if shape.has_text_frame
        for paragraph in shape.text_frame.paragraphs
        for run in paragraph.runs
    )
    assert "Ada" in body


# --- H. Image keys: place chosen images / remove unused image slots ---------------

# Notes marking {{logo}} as an image key bound to a folder (story 01 metadata syntax).
_IMAGE_NOTES = "{{logo}}:\n- type: image\n- folder: Brand/Logos\nThe company logo"


def test_args_schema_image_leaf_is_optional_and_names_its_directory():
    """An image key's leaf is OPTIONAL and its description names the field's note plus the
    source directory by its working_directory PATH (the author's folder). The how-to-browse
    procedure is not repeated per field — it lives once in the main tool description, which
    insists the tree tool MUST be called for each image directory (asserted below)."""
    deck = _build_image_deck([(_IMAGE_NOTES, ["{{logo}}"])])
    params = PptFillerParams(
        schema=_image_schema(deck, slide=1, key="logo", tag_id="tag-logos")
    )
    tool = _the_tool(_FakeAgent(params=params))

    schema = tool.args_schema.model_json_schema()
    defs = schema["$defs"]
    ref = schema["properties"]["slide_1"]["$ref"]
    leaf = defs[ref.split("/")[-1]]

    # Every leaf is optional now (the slide_1 leaf object requires nothing).
    assert leaf.get("required", []) == []
    desc = leaf["properties"]["logo"]["description"]
    # The per-field description carries the note and points at the folder as the
    # working_directory PATH. It does NOT surface the opaque tag id (that is a search
    # filter, not a browsable path) nor repeat the browsing procedure (that moved to the
    # main tool description).
    assert "The company logo" in desc
    assert "working_directory" in desc
    assert "Brand/Logos" in desc
    assert "tag-logos" not in desc
    assert "list_document_tree" not in desc

    # The generic image-handling procedure lives once in the main tool description and
    # insists the tree tool MUST be called first.
    assert "list_document_tree" in tool.description
    assert "working_directory" in tool.description
    assert "MUST" in tool.description


@pytest.mark.asyncio
async def test_image_key_with_doc_id_places_picture_and_removes_placeholder(
    fake_ws, fake_doc
):
    """Providing a doc id for an image key -> the placeholder shape is removed and a
    picture is added; the saved deck shows a picture and no {{logo}} text."""
    deck = _build_image_deck([(_IMAGE_NOTES, ["{{logo}}"])])
    fake_ws.template_bytes = deck
    fake_doc.docs = {"doc-logo": (_png_bytes(10, 10, "red"), "image/png")}
    params = PptFillerParams(schema=_image_schema(deck, slide=1, key="logo"))
    tool = _the_tool(_FakeAgent(params=params, session_id="s"))

    _content, artifact = await tool.coroutine(slide_1={"logo": "doc-logo"})

    assert artifact.is_error is False
    assert fake_doc.fetch_uids == ["doc-logo"]

    filled = Presentation(io.BytesIO(fake_ws.upload_calls[0]["content"]))
    slide = list(filled.slides)[0]
    # A picture was inserted and the {{logo}} placeholder text is gone.
    assert len(_picture_shapes(slide)) == 1
    assert list_keys_on_slide(slide) == []


@pytest.mark.asyncio
async def test_image_fit_inside_preserves_aspect_within_box(fake_ws, fake_doc):
    """The inserted picture fits inside its box and preserves the image aspect ratio.

    A square (10x10) image in a 4in x 2in box (wider than tall) fits the HEIGHT: the
    picture is 2in x 2in, horizontally centered."""
    deck = _build_image_deck([(_IMAGE_NOTES, ["{{logo}}"])])
    fake_ws.template_bytes = deck
    fake_doc.docs = {"doc-logo": (_png_bytes(10, 10, "blue"), "image/png")}
    params = PptFillerParams(schema=_image_schema(deck, slide=1, key="logo"))
    tool = _the_tool(_FakeAgent(params=params, session_id="s"))

    await tool.coroutine(slide_1={"logo": "doc-logo"})

    filled = Presentation(io.BytesIO(fake_ws.upload_calls[0]["content"]))
    picture = _picture_shapes(list(filled.slides)[0])[0]

    box_w = Inches(4)
    box_h = Inches(2)
    box_left = Inches(1)
    # Fits inside the box.
    assert picture.width <= box_w + 2
    assert picture.height <= box_h + 2
    # Square image -> square picture (aspect preserved), sized to the box HEIGHT.
    assert abs(picture.width - picture.height) <= 2
    assert abs(picture.height - box_h) <= 2
    # Horizontally centered within the box.
    expected_left = box_left + (box_w - picture.width) / 2
    assert abs(picture.left - expected_left) <= 2


@pytest.mark.asyncio
async def test_image_key_in_two_shapes_fills_both(fake_ws, fake_doc):
    """An image key in TWO shapes -> a picture inserted in both; both placeholders gone."""
    deck = _build_image_deck([(_IMAGE_NOTES, ["{{logo}}", "{{logo}}"])])
    fake_ws.template_bytes = deck
    fake_doc.docs = {"doc-logo": (_png_bytes(10, 10, "green"), "image/png")}
    params = PptFillerParams(schema=_image_schema(deck, slide=1, key="logo"))
    tool = _the_tool(_FakeAgent(params=params, session_id="s"))

    _content, artifact = await tool.coroutine(slide_1={"logo": "doc-logo"})

    assert artifact.is_error is False
    filled = Presentation(io.BytesIO(fake_ws.upload_calls[0]["content"]))
    slide = list(filled.slides)[0]
    assert len(_picture_shapes(slide)) == 2
    assert list_keys_on_slide(slide) == []


@pytest.mark.asyncio
async def test_omitted_image_key_removes_placeholder_without_picture(fake_ws, fake_doc):
    """An omitted image key -> its placeholder shape is removed, NO picture is added."""
    deck = _build_image_deck([(_IMAGE_NOTES, ["{{logo}}"])])
    fake_ws.template_bytes = deck
    # No docs configured: an omitted key must not trigger any fetch.
    params = PptFillerParams(schema=_image_schema(deck, slide=1, key="logo"))
    tool = _the_tool(_FakeAgent(params=params, session_id="s"))

    _content, artifact = await tool.coroutine(slide_1={"logo": None})

    assert artifact.is_error is False
    assert fake_doc.fetch_uids == []  # omitted -> never fetched
    filled = Presentation(io.BytesIO(fake_ws.upload_calls[0]["content"]))
    slide = list(filled.slides)[0]
    assert _picture_shapes(slide) == []
    # The empty {{logo}} placeholder box is gone (no leftover key text).
    assert list_keys_on_slide(slide) == []


@pytest.mark.asyncio
async def test_bad_image_bytes_hard_fail(fake_ws, fake_doc):
    """Bad image bytes (not an image) -> the fill HARD-fails (is_error), no silent drop
    and no upload."""
    deck = _build_image_deck([(_IMAGE_NOTES, ["{{logo}}"])])
    fake_ws.template_bytes = deck
    fake_doc.docs = {"doc-logo": (b"not an image", "image/png")}
    params = PptFillerParams(schema=_image_schema(deck, slide=1, key="logo"))
    tool = _the_tool(_FakeAgent(params=params, session_id="s"))

    content, artifact = await tool.coroutine(slide_1={"logo": "doc-logo"})

    assert artifact.is_error is True
    assert artifact.ui_parts == ()
    assert "image" in content.lower()
    # Hard fail -> nothing uploaded.
    assert fake_ws.upload_calls == []


@pytest.mark.asyncio
async def test_image_fetch_failure_hard_fail(fake_ws, fake_doc):
    """A doc id that cannot be fetched -> the fill HARD-fails (is_error)."""
    deck = _build_image_deck([(_IMAGE_NOTES, ["{{logo}}"])])
    fake_ws.template_bytes = deck
    fake_doc.docs = {}  # the requested uid is absent -> fetch raises
    params = PptFillerParams(schema=_image_schema(deck, slide=1, key="logo"))
    tool = _the_tool(_FakeAgent(params=params, session_id="s"))

    content, artifact = await tool.coroutine(slide_1={"logo": "missing-doc"})

    assert artifact.is_error is True
    assert "missing-doc" in content
    assert fake_ws.upload_calls == []


@pytest.mark.asyncio
async def test_image_key_with_path_instead_of_id_gives_clear_error(fake_ws, fake_doc):
    """A path-like value (e.g. 'Brand/Logos/logo.png') passed for an image key -> a clear
    hard-fail BEFORE any fetch, telling the agent an id is required. No fetch, no upload."""
    deck = _build_image_deck([(_IMAGE_NOTES, ["{{logo}}"])])
    fake_ws.template_bytes = deck
    params = PptFillerParams(schema=_image_schema(deck, slide=1, key="logo"))
    tool = _the_tool(_FakeAgent(params=params, session_id="s"))

    content, artifact = await tool.coroutine(slide_1={"logo": "Brand/Logos/logo.png"})

    assert artifact.is_error is True
    assert artifact.ui_parts == ()
    # The message explains it looks like a path, that an id is required, and points at the
    # tree tool; the offending value is echoed.
    lowered = content.lower()
    assert "path" in lowered
    assert "document id" in lowered
    assert "list_document_tree" in content
    assert "Brand/Logos/logo.png" in content
    # Detected up front: no fetch was attempted and nothing was uploaded.
    assert fake_doc.fetch_uids == []
    assert fake_ws.upload_calls == []


@pytest.mark.asyncio
async def test_webp_image_is_transcoded_and_embedded(fake_ws, fake_doc):
    """A WEBP doc id -> the toolkit transcodes it to PNG and embeds it as a picture
    instead of hard-failing. python-pptx cannot embed WEBP directly, so without the
    in-memory transcode this would fail at add_picture with 'unsupported image format'."""
    deck = _build_image_deck([(_IMAGE_NOTES, ["{{logo}}"])])
    fake_ws.template_bytes = deck
    fake_doc.docs = {"doc-logo": (_webp_bytes(10, 10, "red"), "image/webp")}
    params = PptFillerParams(schema=_image_schema(deck, slide=1, key="logo"))
    tool = _the_tool(_FakeAgent(params=params, session_id="s"))

    _content, artifact = await tool.coroutine(slide_1={"logo": "doc-logo"})

    assert artifact.is_error is False
    filled = Presentation(io.BytesIO(fake_ws.upload_calls[0]["content"]))
    slide = list(filled.slides)[0]
    # The WEBP was embedded (as PNG) and the {{logo}} placeholder text is gone.
    assert len(_picture_shapes(slide)) == 1
    assert list_keys_on_slide(slide) == []


@pytest.mark.asyncio
async def test_mixed_text_and_image_keys_fill_independently(fake_ws, fake_doc):
    """A slide with a text key AND an image key: the text key fills as text, the image
    key becomes a picture; no placeholders remain."""
    deck = _build_image_deck(
        [
            (
                "{{title}}:\nThe title\n" + _IMAGE_NOTES,
                ["Title: {{title}}", "{{logo}}"],
            )
        ]
    )
    fake_ws.template_bytes = deck
    fake_doc.docs = {"doc-logo": (_png_bytes(20, 10, "red"), "image/png")}
    params = PptFillerParams(schema=_image_schema(deck, slide=1, key="logo"))
    tool = _the_tool(_FakeAgent(params=params, session_id="s"))

    _content, artifact = await tool.coroutine(
        slide_1={"title": "Quarterly", "logo": "doc-logo"}
    )

    assert artifact.is_error is False
    filled = Presentation(io.BytesIO(fake_ws.upload_calls[0]["content"]))
    slide = list(filled.slides)[0]
    assert len(_picture_shapes(slide)) == 1
    assert list_keys_on_slide(slide) == []
    body = "".join(
        run.text
        for shape in slide.shapes
        if shape.has_text_frame
        for paragraph in shape.text_frame.paragraphs
        for run in paragraph.runs
    )
    assert "Quarterly" in body
