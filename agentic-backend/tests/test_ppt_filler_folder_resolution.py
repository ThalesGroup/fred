# Copyright Thales 2025
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.

"""Offline tests for the folder-resolution / image-location validation layer.

Decks are built in-test with python-pptx (no checked-in binaries) and the folder
resolver is an in-memory fake, so these stay fully offline. The deck-building helpers
mirror the parser test file's ``_build_deck`` / ``_build_table_deck`` (copied here rather
than imported, to keep each test file self-contained).
"""

import io
from typing import Dict, List, Optional, Tuple

import pytest
from pptx import Presentation
from pptx.util import Inches

from agentic_backend.integrations.ppt_filler.folder_resolution import (
    CODE_FOLDER_NOT_FOUND,
    CODE_IMAGE_KEY_INVALID_LOCATION,
    resolve_and_validate_images,
)
from agentic_backend.integrations.ppt_filler.parser import (
    CODE_IMAGE_WITHOUT_FOLDER,
    parse,
)

# (text-box body, notes text) per slide.
SlideSpec = Tuple[str, str]


def _build_deck(slides: List[SlideSpec]) -> bytes:
    """Build a ``.pptx`` with one text box and one notes block per slide spec."""
    presentation = Presentation()
    blank_layout = presentation.slide_layouts[6]  # blank
    for body, notes in slides:
        slide = presentation.slides.add_slide(blank_layout)
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
        textbox.text_frame.text = body
        if notes:
            slide.notes_slide.notes_text_frame.text = notes
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _build_table_deck(notes: str, cell_texts: List[str]) -> bytes:
    """Build a one-slide deck whose only key-bearing shape is a TABLE.

    ``cell_texts`` fills the cells of a 1xN table row-major. A ``{{key}}`` placeholder in
    a table cell is discoverable (so it lands in the schema) but is an invalid image
    location (a cell cannot hold a picture).
    """
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    rows, cols = 1, len(cell_texts)
    table = slide.shapes.add_table(
        rows, cols, Inches(1), Inches(1), Inches(8), Inches(1)
    ).table
    for col, text in enumerate(cell_texts):
        table.cell(0, col).text = text
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


class _FakeResolver:
    """In-memory ``folder full-path -> tag id`` resolver; counts calls for dedupe asserts."""

    def __init__(self, known: Dict[str, str]):
        self.known = known
        self.calls: List[str] = []

    async def resolve(self, folder: str) -> Optional[str]:
        self.calls.append(folder)
        return self.known.get(folder)


def _field(result, slide_number: int, key: str):
    slide = next(s for s in result.slides if s.slide == slide_number)
    return next(f for f in slide.keys if f.key == key)


def _codes(result) -> List[str]:
    return [e.code for e in result.errors]


# --- folder resolution -----------------------------------------------------------------


@pytest.mark.asyncio
async def test_existing_folder_writes_tag_id_and_no_error():
    notes = "{{flag}}:\n- type: image\n- folder: images/flags\nPick a flag."
    deck = _build_deck([("{{flag}}", notes)])
    resolver = _FakeResolver({"images/flags": "tag-123"})

    result = await resolve_and_validate_images(deck, parse(deck), resolver)

    assert _field(result, 1, "flag").folder_tag_id == "tag-123"
    assert CODE_FOLDER_NOT_FOUND not in _codes(result)
    assert resolver.calls == ["images/flags"]


@pytest.mark.asyncio
async def test_existing_folder_tag_id_is_reflected_in_serialized_schema():
    """The resolved tag id must show up in the returned result's wire schema (the image
    KeyField now serializes ``folder_tag_id`` because it is no longer None)."""
    notes = "{{flag}}:\n- type: image\n- folder: images/flags\nPick a flag."
    deck = _build_deck([("{{flag}}", notes)])
    resolver = _FakeResolver({"images/flags": "tag-123"})

    result = await resolve_and_validate_images(deck, parse(deck), resolver)

    dumped = result.model_dump(by_alias=True)
    assert dumped["schema"] == [
        {
            "slide": 1,
            "keys": [
                {
                    "key": "flag",
                    "description": "Pick a flag.",
                    "type": "image",
                    "folder": "images/flags",
                    "folder_tag_id": "tag-123",
                }
            ],
        }
    ]


@pytest.mark.asyncio
async def test_missing_folder_reports_folder_not_found_and_leaves_tag_none():
    notes = "{{flag}}:\n- type: image\n- folder: images/missing\nPick a flag."
    deck = _build_deck([("{{flag}}", notes)])
    resolver = _FakeResolver({"images/flags": "tag-123"})  # different folder

    result = await resolve_and_validate_images(deck, parse(deck), resolver)

    matching = [e for e in result.errors if e.code == CODE_FOLDER_NOT_FOUND]
    assert len(matching) == 1
    assert matching[0].slide == 1
    assert matching[0].key == "flag"
    assert "images/missing" in matching[0].message
    assert _field(result, 1, "flag").folder_tag_id is None


@pytest.mark.asyncio
async def test_shared_folder_is_resolved_once_and_both_keys_get_tag():
    """Two image keys sharing one folder -> resolver called once (dedupe), both filled."""
    notes = (
        "{{flag_a}}:\n- type: image\n- folder: images/flags\nFirst.\n"
        "{{flag_b}}:\n- type: image\n- folder: images/flags\nSecond."
    )
    deck = _build_deck([("{{flag_a}} {{flag_b}}", notes)])
    resolver = _FakeResolver({"images/flags": "tag-shared"})

    result = await resolve_and_validate_images(deck, parse(deck), resolver)

    assert resolver.calls == ["images/flags"]  # resolved once, not per key
    assert _field(result, 1, "flag_a").folder_tag_id == "tag-shared"
    assert _field(result, 1, "flag_b").folder_tag_id == "tag-shared"
    assert CODE_FOLDER_NOT_FOUND not in _codes(result)


@pytest.mark.asyncio
async def test_text_keys_are_skipped_entirely():
    """A plain text key has no folder and no geometry constraint: resolver untouched."""
    deck = _build_deck([("{{name}}", "{{name}}:\nThe name.")])
    resolver = _FakeResolver({"images/flags": "tag-123"})

    result = await resolve_and_validate_images(deck, parse(deck), resolver)

    assert resolver.calls == []
    assert result.errors == []
    assert _field(result, 1, "name").folder_tag_id is None


@pytest.mark.asyncio
async def test_empty_folder_is_not_re_reported_as_folder_not_found():
    """An image with a blank folder is already image_without_folder (Story 01); this
    layer must not resolve it nor add folder_not_found."""
    notes = "{{flag}}:\n- type: image\n- folder:\nPick a flag."
    deck = _build_deck([("{{flag}}", notes)])
    resolver = _FakeResolver({"images/flags": "tag-123"})

    base = parse(deck)
    assert CODE_IMAGE_WITHOUT_FOLDER in _codes(base)  # precondition from Story 01

    result = await resolve_and_validate_images(deck, base, resolver)

    assert resolver.calls == []  # empty folder is never resolved
    assert CODE_FOLDER_NOT_FOUND not in _codes(result)
    # The original image_without_folder error is preserved (not dropped, not duplicated).
    assert _codes(result).count(CODE_IMAGE_WITHOUT_FOLDER) == 1


@pytest.mark.asyncio
async def test_base_errors_are_preserved_alongside_new_errors():
    """Base parse errors come through; resolution errors are appended, not replacing."""
    notes = "{{flag}}:\n- type: image\n- folder: images/missing\nPick a flag."
    # {{flag}} is described+present (resolves to missing folder); {{lonely}} is on the
    # slide but undescribed -> key_without_description from the base parse.
    deck = _build_deck([("{{flag}} {{lonely}}", notes)])
    resolver = _FakeResolver({})

    base = parse(deck)
    base_codes = set(_codes(base))

    result = await resolve_and_validate_images(deck, base, resolver)

    assert base_codes.issubset(set(_codes(result)))  # nothing dropped
    assert CODE_FOLDER_NOT_FOUND in _codes(result)  # new error appended


# --- invalid location (table cell) ------------------------------------------------------


@pytest.mark.asyncio
async def test_image_key_in_table_cell_is_invalid_location():
    notes = "{{flag}}:\n- type: image\n- folder: images/flags\nPick a flag."
    deck = _build_table_deck(notes, ["{{flag}}"])
    resolver = _FakeResolver({"images/flags": "tag-123"})

    result = await resolve_and_validate_images(deck, parse(deck), resolver)

    matching = [e for e in result.errors if e.code == CODE_IMAGE_KEY_INVALID_LOCATION]
    assert len(matching) == 1
    assert matching[0].slide == 1
    assert matching[0].key == "flag"
    # The heading tells the author how to fix it.
    assert "text box" in matching[0].message


@pytest.mark.asyncio
async def test_image_key_in_table_cell_still_resolves_its_folder():
    """A mislocated key is still resolved (the two checks are independent), so the author
    sees BOTH the location error and, if the folder is bad, folder_not_found."""
    notes = "{{flag}}:\n- type: image\n- folder: images/flags\nPick a flag."
    deck = _build_table_deck(notes, ["{{flag}}"])
    resolver = _FakeResolver({"images/flags": "tag-123"})

    result = await resolve_and_validate_images(deck, parse(deck), resolver)

    assert _field(result, 1, "flag").folder_tag_id == "tag-123"
    assert CODE_IMAGE_KEY_INVALID_LOCATION in _codes(result)


@pytest.mark.asyncio
async def test_text_key_in_table_cell_is_not_invalid_location():
    """A TEXT key in a table cell is fine (text is fillable in a cell); only image keys
    are constrained."""
    deck = _build_table_deck("{{label}}:\nA label.", ["{{label}}"])
    resolver = _FakeResolver({})

    result = await resolve_and_validate_images(deck, parse(deck), resolver)

    assert CODE_IMAGE_KEY_INVALID_LOCATION not in _codes(result)


@pytest.mark.asyncio
async def test_image_key_in_textbox_is_valid_location():
    notes = "{{flag}}:\n- type: image\n- folder: images/flags\nPick a flag."
    deck = _build_deck([("{{flag}}", notes)])
    resolver = _FakeResolver({"images/flags": "tag-123"})

    result = await resolve_and_validate_images(deck, parse(deck), resolver)

    assert CODE_IMAGE_KEY_INVALID_LOCATION not in _codes(result)
