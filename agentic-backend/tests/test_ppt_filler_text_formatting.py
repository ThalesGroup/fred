"""Offline tests for inline bold/italic formatting in the PPT filler text traversal.

A filled ``{{key}}`` value may carry inline Markdown (``**bold**`` / ``*italic*``);
the filler parses it into pptx runs that OVERLAY bold/italic on top of the placeholder
run's own font (size/color/name inherited). Markup is honored only in the substituted
value — never in the surrounding static template text.

Decks are built in-test with python-pptx (no checked-in binaries), matching the style of
``test_ppt_filler_parser.py``.
"""

from __future__ import annotations

import io
from typing import List

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from agentic_backend.integrations.ppt_filler.traversal import (
    list_keys_on_slide,
    replace_keys_on_slide,
)

_FONT_NAME = "Arial"
_FONT_SIZE = Pt(20)
_FONT_COLOR = RGBColor(0x12, 0x34, 0x56)


def _one_run_deck(body: str) -> bytes:
    """A one-slide deck whose only text box holds ``body`` in a single styled run.

    The run carries a distinctive size/color/name so tests can assert that emphasized
    spans INHERIT the base font and only toggle bold/italic.
    """
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
    run = textbox.text_frame.paragraphs[0].add_run()
    run.text = body
    run.font.name = _FONT_NAME
    run.font.size = _FONT_SIZE
    run.font.color.rgb = _FONT_COLOR
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _fill(deck: bytes, values: dict) -> "Presentation":
    presentation = Presentation(io.BytesIO(deck))
    slide = presentation.slides[0]
    replace_keys_on_slide(slide, lambda key: values[key])
    return presentation


def _runs(presentation) -> List:
    slide = presentation.slides[0]
    return [
        run
        for shape in slide.shapes
        if shape.has_text_frame
        for paragraph in shape.text_frame.paragraphs
        for run in paragraph.runs
    ]


def _run_named(runs, text):
    return next((r for r in runs if r.text == text), None)


def _paragraph_with(presentation, needle: str):
    slide = presentation.slides[0]
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            if needle in "".join(r.text for r in paragraph.runs):
                return paragraph
    raise AssertionError(f"no paragraph containing {needle!r}")


def test_bold_value_becomes_a_bold_run_inheriting_base_font():
    deck = _one_run_deck("Result: {{val}}")
    runs = _runs(_fill(deck, {"val": "up **40%** today"}))

    bold = _run_named(runs, "40%")
    assert bold is not None
    assert bold.font.bold is True
    # Base font is inherited on the emphasized span: size/color/name unchanged.
    assert bold.font.size == _FONT_SIZE
    assert bold.font.name == _FONT_NAME
    assert bold.font.color.rgb == _FONT_COLOR
    # No placeholder remains and the full text reads with markers stripped.
    assert list_keys_on_slide(_fill(deck, {"val": "up **40%** today"}).slides[0]) == []
    assert "".join(r.text for r in runs) == "Result: up 40% today"


def test_italic_value_becomes_an_italic_run_inheriting_base_font():
    deck = _one_run_deck("Product: {{name}}")
    runs = _runs(_fill(deck, {"name": "the *Acme* widget"}))

    italic = _run_named(runs, "Acme")
    assert italic is not None
    assert italic.font.italic is True
    assert italic.font.size == _FONT_SIZE
    assert italic.font.name == _FONT_NAME
    assert italic.font.color.rgb == _FONT_COLOR


def test_value_with_no_markup_yields_a_single_unchanged_run():
    deck = _one_run_deck("Hi {{name}}!")
    presentation = _fill(deck, {"name": "Jane"})
    runs = _runs(presentation)

    # A markup-free fill collapses onto one run (today's behavior), base font intact.
    assert len(runs) == 1
    assert runs[0].text == "Hi Jane!"
    assert runs[0].font.size == _FONT_SIZE
    assert runs[0].font.name == _FONT_NAME
    assert runs[0].font.color.rgb == _FONT_COLOR
    assert runs[0].font.bold in (None, False)
    assert runs[0].font.italic in (None, False)


def test_literal_star_in_static_template_text_is_left_untouched():
    # The author's static '*' must NOT be reinterpreted as emphasis — only the value is
    # parsed. The value here carries no markup, so nothing should become italic/bold.
    deck = _one_run_deck("Rating {{score}} out of *five*")
    presentation = _fill(deck, {"score": "4"})
    runs = _runs(presentation)

    full = "".join(r.text for r in runs)
    assert full == "Rating 4 out of *five*"  # literal asterisks preserved
    assert all(r.font.italic in (None, False) for r in runs)
    assert all(r.font.bold in (None, False) for r in runs)


def test_literal_star_in_static_text_survives_alongside_value_markup():
    # Static '*five*' stays literal even when the VALUE contains real markup that is
    # honored — the two never bleed into each other.
    deck = _one_run_deck("{{label}} out of *five*")
    presentation = _fill(deck, {"label": "scored **4**"})
    runs = _runs(presentation)

    assert "".join(r.text for r in runs) == "scored 4 out of *five*"
    bold = _run_named(runs, "4")
    assert bold is not None and bold.font.bold is True
    # The literal-asterisk run is plain text, not emphasized.
    star_run = next(r for r in runs if "*five*" in r.text)
    assert star_run.font.bold in (None, False)
    assert star_run.font.italic in (None, False)


def test_bold_italic_combined_value_renders_bold():
    deck = _one_run_deck("{{k}}")
    runs = _runs(_fill(deck, {"k": "***wow***"}))
    bold = _run_named(runs, "wow")
    assert bold is not None and bold.font.bold is True


def test_markup_does_not_affect_key_detection_roundtrip():
    # Markup lives in values, not keys: a value with markup still fills cleanly and leaves
    # no placeholder behind (the parse->fill->reparse invariant).
    deck = _one_run_deck("A {{x}} and a {{y}}")
    presentation = _fill(deck, {"x": "**bold**", "y": "_em_"})
    assert list_keys_on_slide(presentation.slides[0]) == []


def test_newline_becomes_a_line_break_not_literal_in_run_text():
    # A '\n' in a value must become an <a:br/> element, NOT a literal newline inside a
    # run's text. A literal '\n' renders unreliably and bleeds the PRECEDING run's bold
    # across the wrapped line (a bold title leaking onto the body below it).
    from pptx.oxml.ns import qn

    deck = _one_run_deck("{{body}}")
    presentation = _fill(deck, {"body": "**Title**\nplain body line"})
    paragraph = _paragraph_with(presentation, "Title")

    # No run text contains a raw newline...
    assert all("\n" not in r.text for r in paragraph.runs)
    # ...and there is a real <a:br/> between the title and the body.
    assert paragraph._p.findall(qn("a:br"))


def test_bold_title_does_not_bleed_onto_following_body_line():
    # The regression that motivated <a:br/>: with markup-toggled runs, the body run after
    # a bold title must stay NON-bold (the title's bold must not carry over the line wrap).
    deck = _one_run_deck("{{body}}")
    value = "**Innovation Hub Expansion**\nA new hub was inaugurated in Toulouse."
    presentation = _fill(deck, {"body": value})
    runs = _runs(presentation)

    title = _run_named(runs, "Innovation Hub Expansion")
    body = _run_named(runs, "A new hub was inaugurated in Toulouse.")
    assert title is not None and title.font.bold is True
    assert body is not None
    assert body.font.bold in (None, False)


def test_multiple_titles_and_bodies_keep_independent_bold():
    # The full screenshot scenario: three "**Title**\nbody\n\n" blocks. Every title bold,
    # every body plain — bold never bleeds past its own title.
    deck = _one_run_deck("{{body}}")
    value = "**First**\nbody one.\n\n**Second**\nbody two.\n\n**Third**\nbody three."
    presentation = _fill(deck, {"body": value})
    runs = _runs(presentation)

    for title in ("First", "Second", "Third"):
        run = _run_named(runs, title)
        assert run is not None and run.font.bold is True, title
    for body in ("body one.", "body two.", "body three."):
        run = _run_named(runs, body)
        assert run is not None and run.font.bold in (None, False), body
